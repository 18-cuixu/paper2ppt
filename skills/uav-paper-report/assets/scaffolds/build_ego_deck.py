from __future__ import annotations

import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import build_v60_deck as base


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "out" / "uav-ego-planner-report.pptx"
CROPS = ROOT / "ego-paper-test" / "crops"

BLACK = base.BLACK
RED = base.RED
GRAY = base.GRAY
BLUE = base.BLUE
BODY_FONT = "Times New Roman"
TABLE_SIZE = 13.5

KEY_TERMS = [
    "ESDF-free", "必要障碍", "B-spline", "L-BFGS",
    "0.81 ms", "0.37 ms", "0.89", "24.38 s", "42.24 m",
    "实时", "成功率", "局部重规划",
]


def add_box(slide, left, top, width, height):
    return base.add_box(slide, left, top, width, height)


def reject_manual_newline(text: str, context: str, *, allow_newlines: bool = False) -> None:
    if "\n" in text and not allow_newlines:
        raise ValueError(f"{context}: manual newline is not allowed in body text")
    if "\n" in text and any(not part.strip() for part in text.split("\n")):
        raise ValueError(f"{context}: empty line inside manual line break")


def plain(
    box,
    text: str,
    size: float,
    *,
    bold=False,
    color=BLACK,
    align=PP_ALIGN.LEFT,
    allow_newlines: bool = False,
) -> None:
    reject_manual_newline(text, f"plain({getattr(box, 'name', 'textbox')})", allow_newlines=allow_newlines)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 1.0
        run = p.add_run()
        run.text = line
        base.set_run_font(run, BODY_FONT)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def set_runs(paragraph, text: str, size: float, *, bold=False, terms=None) -> None:
    for piece, marked in base.split_terms(text, KEY_TERMS + (terms or [])):
        run = paragraph.add_run()
        run.text = compact_text(piece)
        base.set_run_font(run, BODY_FONT)
        run.font.size = Pt(size)
        run.font.bold = bold or marked
        run.font.color.rgb = RED if marked else BLACK


def compact_text(text: str) -> str:
    join = "\u2060"
    replacements = {
        "EGO-Planner": "EGO‑Planner",
        "ESDF-free": "ESDF‑free",
        "gradient-based": "gradient‑based",
        "collision-free": "collision‑free",
        "B-spline": "B‑spline",
        "L-BFGS": "L‑BFGS",
        "Fast-Planner": "Fast‑Planner",
        "real-world": "real‑world",
        "time re-allocation": "time re‑allocation",
        "order of magnitude": "order\u00a0of\u00a0magnitude",
        "0.81 ms": "0.81\u00a0ms",
        "0.37 ms": "0.37\u00a0ms",
        "24.38 s": "24.38\u00a0s",
        "42.24 m": "42.24\u00a0m",
        "3.56 m/s": "3.56\u00a0m/s",
        "0.5 m/s": "0.5\u00a0m/s",
        "IEEE Robotics and Automation Letters": "IEEE\u00a0Robotics\u00a0and\u00a0Automation\u00a0Letters",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    text = re.sub(r"(?<=\d)\s+(?=(ms|s|m/s|m)\b)", "\u00a0", text)
    text = re.sub(r"(?<=[A-Za-z0-9+\-/])\s+(?=[\u4e00-\u9fff])", "", text)
    text = re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[A-Za-z0-9])", "", text)
    for phrase in (
        "有限FOV", "局部重规划", "实时重规划", "必要障碍集合", "控制点约束",
        "碰撞惩罚", "动力学可行性", "时间重分配", "室内窄通道",
        "随机目标追踪", "森林场景", "轨迹优化", "全局路径",
    ):
        text = text.replace(phrase, join.join(phrase))
    return text


def set_textbox(box, lines: list[dict], *, default_size=20, space_after=1.0, line_spacing=0.96, terms=None) -> None:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for idx, item in enumerate(lines):
        text = item["text"]
        reject_manual_newline(text, f"bullet {idx + 1} in {getattr(box, 'name', 'textbox')}")
        level = item.get("level", 0)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.line_spacing = max(line_spacing, 1.05)
        p.space_before = Pt(0)
        p.space_after = Pt(item.get("space_after", space_after))
        ppr = p._p.get_or_add_pPr()
        if level == 0:
            ppr.set("marL", "210000")
            ppr.set("indent", "-140000")
            prefix = item.get("marker", "● ")
            size = float(item.get("size", default_size))
        elif level == 1:
            ppr.set("marL", "350000")
            ppr.set("indent", "-120000")
            prefix = item.get("marker", "• ")
            size = float(item.get("size", default_size))
        else:
            ppr.set("marL", "560000")
            ppr.set("indent", "-105000")
            prefix = item.get("marker", "– ")
            size = float(item.get("size", default_size))
        set_runs(p, prefix + text, size, bold=item.get("bold", level == 0), terms=terms)


def add_header(slide, part: str, title: str) -> None:
    base.add_header(slide, part, title)


def add_section_title(slide, num: str, title: str) -> None:
    base.add_section_title(slide, num, title)


def add_rule(slide, left: float, top: float, width: float, color=RGBColor(190, 190, 190)) -> None:
    base.add_rule(slide, left, top, width, color)


def add_pic(slide, path: Path, left, top, width, height):
    return base.add_pic(slide, path, left, top, width, height)


def label(slide, text: str, left, top, width):
    box = add_box(slide, left, top, width, 0.32)
    plain(box, text, 19, bold=True, color=GRAY)


def equation_block(slide, rows: list[tuple], left: float, top: float, width: float, height: float, *, size=22.0) -> None:
    count = max(len(rows), 1)
    row_h = height / count
    for idx, row in enumerate(rows):
        eq_label, text = row[0], row[1]
        row_size = float(row[2]) if len(row) > 2 else size
        box = add_box(slide, left, top + idx * row_h, width, row_h * 0.88)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = False
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.0
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        r1 = p.add_run()
        r1.text = f"{eq_label}    "
        base.set_run_font(r1, BODY_FONT)
        r1.font.size = Pt(row_size)
        r1.font.bold = True
        r1.font.color.rgb = GRAY
        r2 = p.add_run()
        r2.text = text
        base.set_run_font(r2, BODY_FONT)
        r2.font.size = Pt(row_size)
        r2.font.color.rgb = BLACK


def add_table(slide, left, top, width, height, rows, *, font_size=13.5, red_values=()):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(value)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.text_frame.word_wrap = False
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(242, 242, 242) if ri == 0 else RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                base.set_run_font(run, BODY_FONT)
                run.font.size = Pt(font_size)
                run.font.bold = ri == 0 or ci == 0
                run.font.color.rgb = RED if str(value) in set(red_values) else BLACK
    return shape


def metric_row(slide, items: list[tuple[str, str, str]], left: float, top: float, width: float, *, size=20.0) -> None:
    col_w = width / len(items)
    for idx, (value, label_text, note) in enumerate(items):
        x = left + idx * col_w
        value_box = add_box(slide, x, top, col_w - 0.12, 0.42)
        value_box.name = f"METRIC_VALUE_{idx}"
        plain(value_box, value, size + 6, bold=True, color=RED, align=PP_ALIGN.CENTER)
        label_box = add_box(slide, x, top + 0.44, col_w - 0.12, 0.28)
        label_box.name = f"METRIC_LABEL_{idx}"
        plain(label_box, label_text, size - 2, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        note_box = add_box(slide, x, top + 0.76, col_w - 0.12, 0.34)
        note_box.name = f"METRIC_NOTE_{idx}"
        plain(note_box, note, size - 4, color=BLACK, align=PP_ALIGN.CENTER)


def divider_slide(prs, part: str, title: str, bullets: list[dict]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, title)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(1.44), Inches(13.333), Inches(1.16))
    band.fill.solid()
    band.fill.fore_color.rgb = base.HEADER_GRAY
    band.line.fill.background()
    title_box = add_box(slide, 0.78, 1.72, 11.90, 0.52)
    plain(title_box, title, 28, bold=True, align=PP_ALIGN.CENTER)
    box = add_box(slide, 1.16, 3.10, 11.00, 2.52)
    normal_bullets = [{**item, "bold": False} for item in bullets]
    set_textbox(box, normal_bullets, default_size=22.0, space_after=0.08, line_spacing=1.04)
    return slide


def split_slide(prs, part, part_title, sec, sec_title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    label(slide, left_title, 0.76, 1.74, 5.20)
    left = add_box(slide, 0.72, 2.10, 11.92, 2.00)
    set_textbox(left, left_bullets, default_size=18.7, space_after=0.55, line_spacing=1.05)
    add_rule(slide, 0.76, 4.18, 11.82)
    label(slide, right_title, 0.76, 4.36, 5.20)
    right = add_box(slide, 0.72, 4.72, 11.92, 2.04)
    set_textbox(right, right_bullets, default_size=18.7, space_after=0.55, line_spacing=1.05)
    return slide


def bullet_with_metrics_slide(prs, part, part_title, sec, sec_title, bullets, metrics, *, side_title=None, side_notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    body_lines = list(bullets)
    if side_notes:
        body_lines.append({"level": 0, "text": side_title or "结构链路"})
        heading = None
        for note in side_notes:
            if note.get("level", 0) == 0:
                heading = note["text"]
                continue
            text = f"{heading}：{note['text']}" if heading else note["text"]
            body_lines.append({"level": 1, "text": text})
            heading = None
    box = add_box(slide, 0.62, 1.76, 12.00, 3.74)
    set_textbox(box, body_lines, default_size=19.2 if side_notes else 20.0, space_after=0.65, line_spacing=1.05)
    add_rule(slide, 0.78, 5.70, 11.78)
    metric_row(slide, metrics, 0.86, 5.98, 11.58, size=19.0)
    return slide


def table_with_notes_slide(prs, part, part_title, sec, sec_title, title, rows, notes, *, red_values=(), font_size=13.2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    top = add_box(slide, 0.62, 1.76, 12.05, 0.70)
    set_textbox(top, [
        {"level": 0, "text": title},
    ], default_size=20.2, space_after=0, line_spacing=0.96)
    add_table(slide, 0.72, 2.62, 11.96, 1.90, rows, font_size=font_size, red_values=red_values)
    add_rule(slide, 0.78, 4.92, 11.78)
    note_box = add_box(slide, 0.76, 5.12, 11.86, 1.55)
    set_textbox(note_box, notes, default_size=19.6, space_after=1.0, line_spacing=0.94)
    return slide


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if base.LOGO.exists():
        slide.shapes.add_picture(str(base.LOGO), Inches(4.10), Inches(0.64), width=Inches(5.08), height=Inches(1.41))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.40), Inches(13.333), Inches(3.10))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    title = add_box(slide, 0.55, 3.02, 12.25, 1.05)
    plain(title, "EGO-Planner: An ESDF-free Gradient-based\nLocal Planner for Quadrotors", 30, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, allow_newlines=True)
    subtitle = add_box(slide, 0.0, 4.24, 13.33, 0.74)
    plain(subtitle, "Xin Zhou, Zhepei Wang, Hongkai Ye, Chao Xu, Fei Gao\nIEEE Robotics and Automation Letters, 2020", 18, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, allow_newlines=True)
    date = add_box(slide, 5.50, 6.62, 2.35, 0.45)
    plain(date, "2026年5月8日", 18, color=BLUE, align=PP_ALIGN.CENTER)
    return slide


def thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if base.LOGO.exists():
        slide.shapes.add_picture(str(base.LOGO), Inches(4.10), Inches(0.64), width=Inches(5.08), height=Inches(1.41))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.40), Inches(13.333), Inches(3.10))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    title = add_box(slide, 0.0, 3.34, 13.33, 0.96)
    plain(title, "谢谢！", 46, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    date = add_box(slide, 5.50, 6.62, 2.35, 0.45)
    plain(date, "2026年5月8日", 18, color=BLUE, align=PP_ALIGN.CENTER)
    return slide


def body_slide(prs, part, part_title, sec, sec_title, bullets, *, size=20.5):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    box = add_box(slide, 0.56, 1.76, 12.15, 5.58)
    set_textbox(box, bullets, default_size=size, space_after=1.0, line_spacing=0.96)
    return slide


def visual_slide(prs, part, part_title, sec, sec_title, top_bullets, image_specs, bottom_bullets, *, top_height=0.88):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    top = add_box(slide, 0.62, 1.76, 12.05, top_height)
    set_textbox(top, top_bullets, default_size=20.0, space_after=0, line_spacing=0.93)
    for spec in image_specs:
        add_pic(slide, *spec)
    add_rule(slide, 0.74, 5.70, 11.80)
    bottom = add_box(slide, 0.70, 5.88, 11.95, 1.04)
    set_textbox(bottom, bottom_bullets, default_size=18.8, space_after=0, line_spacing=0.88)
    return slide


def bottleneck_slide(prs):
    return split_slide(
        prs,
        "Part. 01",
        "研究背景及动机",
        "1.2",
        "ESDF 瓶颈与局部极小",
        "计算瓶颈",
        [
            {"level": 0, "text": "ESDF 更新范围通常大于轨迹实际搜索空间。"},
            {"level": 1, "text": "局部地图每次融合深度图后都要维护距离场。"},
            {"level": 1, "text": "轨迹优化只需要当前控制点附近的梯度。"},
            {"level": 1, "text": "大量 ESDF 更新没有进入优化目标。"},
        ],
        "优化瓶颈",
        [
            {"level": 0, "text": "有限 FOV 让障碍背面不可见，距离场可能误导轨迹。"},
            {"level": 1, "text": "初始轨迹穿过障碍时，梯度方向不一定指向可行通道。"},
            {"level": 1, "text": "单独依赖 ESDF 时，局部优化容易停留在狭窄区域。"},
            {"level": 1, "text": "EGO 显式记录碰撞约束，并通过控制点优化将轨迹推离障碍。"},
        ],
    )


def technical_points_slide(prs):
    return bullet_with_metrics_slide(
        prs,
        "Part. 01",
        "研究背景及动机",
        "1.3",
        "技术要点",
        [
            {"level": 0, "text": "方法要点集中在 ESDF-free 局部规划框架、控制点优化和动力学可行性修正。"},
            {"level": 1, "text": "提出 ESDF-free 局部规划框架，直接从障碍生成排斥梯度。"},
            {"level": 1, "text": "把 B-spline 控制点、碰撞惩罚和动力学惩罚放入统一优化目标。"},
            {"level": 1, "text": "提出时间重分配和各向异性拟合，使优化后轨迹重新满足动力学约束。"},
            {"level": 1, "text": "集成到自主四旋翼系统，并在仿真、室内和森林场景中验证。"},
        ],
        [
            ("0.37 ms", "无 ESDF 总耗时", "EGO 与 EI 成功率同为 0.89"),
            ("0.81 ms", "对比规划时间", "低于 EWOK 与 Fast-Planner"),
            ("24.38 s", "平均飞行时间", "实验表中低于两个基线"),
        ],
        side_title="方法链路",
        side_notes=[
            {"level": 0, "text": "必要障碍表示"},
            {"level": 1, "text": "仅将当前轨迹相关的碰撞障碍写入优化。"},
            {"level": 0, "text": "控制点优化"},
            {"level": 1, "text": "用 B-spline 控制点统一承载平滑、避障和动力学代价。"},
            {"level": 0, "text": "系统验证"},
            {"level": 1, "text": "在仿真、室内窄通道和森林飞行中验证实时性。"},
        ],
    )


def pv_representation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.1", "{p,v} 障碍表示")
    top = add_box(slide, 0.62, 1.76, 12.05, 0.82)
    set_textbox(top, [
        {"level": 0, "text": "碰撞控制点不再查询 ESDF，而是绑定障碍表面锚点 p 和方向向量 v。"},
        {"level": 1, "text": "{p,v} 对定义了局部线性距离场，也给出可直接使用的梯度方向。"},
    ], default_size=20.2, line_spacing=0.96)
    add_pic(slide, CROPS / "fig3_pv_pairs.png", 0.82, 2.78, 7.58, 2.46)
    add_pic(slide, CROPS / "fig4_convex_hull.png", 8.62, 2.94, 3.54, 2.04)
    add_rule(slide, 0.74, 5.60, 11.82)
    bottom = add_box(slide, 0.76, 5.78, 11.92, 1.00)
    set_textbox(bottom, [
        {"level": 1, "text": "轨迹穿过障碍时，为相关控制点生成多个 {p,v} 对；优化阶段仅处理这些必要障碍。"},
        {"level": 1, "text": "B-spline 凸包性质让控制点约束可间接约束曲线段，比逐点采样更稳定。"},
    ], default_size=19.2, space_after=0, line_spacing=0.92)
    return slide


def bspline_basis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.2", "B-spline 控制点表示")
    top = add_box(slide, 0.68, 1.70, 12.00, 1.20)
    set_textbox(top, [
        {"level": 0, "text": "轨迹由控制点 Q 和均匀节点间隔 Δt 表示。"},
        {"level": 1, "text": "位置曲线由相邻控制点加权得到。"},
        {"level": 1, "text": "速度、加速度和 jerk 写成控制点差分，动力学约束不需要对连续曲线逐点检查。"},
    ], default_size=19.6, space_after=0.5, line_spacing=1.05)
    label(slide, "控制点差分关系", 0.86, 3.08, 4.00)
    equation_block(slide, [
        ("", "p(t) = Σᵢ Nᵢ,ₖ(t) Qᵢ", 24.0),
        ("", "Vᵢ = (Qᵢ₊₁ − Qᵢ)/Δt", 23.0),
        ("", "Aᵢ = (Vᵢ₊₁ − Vᵢ)/Δt", 23.0),
        ("", "Jᵢ = (Aᵢ₊₁ − Aᵢ)/Δt", 23.0),
    ], 1.52, 3.48, 10.68, 1.56, size=22.0)
    add_rule(slide, 0.76, 5.32, 11.82)
    detail = add_box(slide, 0.78, 5.52, 11.78, 1.14)
    set_textbox(detail, [
        {"level": 1, "text": "B-spline 的凸包性质把连续曲线安全性转成控制点附近的局部约束。"},
        {"level": 1, "text": "速度、加速度和 jerk 都由差分得到，Jd 直接惩罚超限控制点。"},
        {"level": 1, "text": "局部支撑意味着修改一个控制点只影响有限时间段，适合高频重规划。"},
    ], default_size=17.8, space_after=0, line_spacing=1.05)
    return slide


def formula_slide_distance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.3", "ESDF-free 障碍梯度")
    top = add_box(slide, 0.62, 1.76, 12.05, 0.86)
    set_textbox(top, [
        {"level": 0, "text": "碰撞控制点直接关联 {p,v} 障碍对，用局部几何关系构造距离和排斥梯度。"},
        {"level": 1, "text": "p 是障碍表面锚点，v 是从控制点指向 p 的排斥方向。"},
    ], default_size=20.4, line_spacing=0.95)
    label(slide, "距离定义与碰撞惩罚", 0.86, 2.64, 4.00)
    equation_block(slide, [
        ("(1)", "dᵢⱼ = (Qᵢ − pᵢⱼ) · vᵢⱼ", 24.0),
        ("(5a)", "cᵢⱼ = s₀ − dᵢⱼ", 23.0),
        ("(5b)", "jᶜ(i,j)=0,                              cᵢⱼ≤0", 21.2),
        ("(5c)", "jᶜ(i,j)=cᵢⱼ³,                         0<cᵢⱼ≤s₀", 21.2),
        ("(5d)", "jᶜ(i,j)=3s₀cᵢⱼ²−3s₀²cᵢⱼ+s₀³,   cᵢⱼ>s₀", 20.4),
    ], 1.02, 2.96, 11.68, 1.92, size=22.0)
    add_rule(slide, 0.74, 5.24, 11.80)
    notes = add_box(slide, 0.78, 5.44, 11.82, 1.22)
    set_textbox(notes, [
        {"level": 1, "text": "dᵢⱼ>0 表示控制点已离开该障碍半空间；dᵢⱼ 越小，碰撞惩罚越大。"},
        {"level": 1, "text": "分段三次函数保持连续可导，避免优化过程中梯度突变。"},
        {"level": 1, "text": "当 cᵢⱼ>s₀ 时继续增大惩罚，相当于强制把控制点推回安全侧。"},
        {"level": 1, "text": "每个碰撞控制点可关联多个 {p,v} 对，用来描述附近多个障碍面。"},
    ], default_size=17.8, space_after=0, line_spacing=1.05)
    return slide


def formula_slide_objective(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.4", "轨迹优化目标")
    top = add_box(slide, 0.62, 1.76, 12.05, 0.78)
    set_textbox(top, [
        {"level": 0, "text": "轨迹用均匀 B-spline 控制点 Q 表示，优化目标由平滑、避障和动力学可行性组成。"},
        {"level": 1, "text": "论文采用 L-BFGS 做数值优化，适合快速重启和近似二阶更新。"},
    ], default_size=20.3, line_spacing=0.95)
    label(slide, "核心目标函数", 0.86, 2.56, 3.60)
    equation_block(slide, [
        ("(3)", "min  J(Q) = λₛJₛ + λᶜJᶜ + λᵈJᵈ", 24.0),
        ("(4)", "Jₛ = Σ||Aᵢ||₂² + Σ||Jᵢ||₂²", 23.0),
        ("(6)", "Jᶜ = Σᵢ jᶜ(Qᵢ),    jᶜ(Qᵢ)=Σⱼ jᶜ(i,j)", 22.0),
        ("(8)", "Jᵈ = Σᵢ wᵥF(Vᵢ)+Σᵢ wₐF(Aᵢ)+Σᵢ wⱼF(Jᵢ)", 21.4),
    ], 1.04, 2.90, 11.70, 1.72, size=22.0)
    add_rule(slide, 0.74, 4.90, 11.80)
    notes = add_box(slide, 0.78, 5.10, 11.82, 1.58)
    set_textbox(notes, [
        {"level": 1, "text": "Jₛ 控制加速度和 jerk，保证轨迹平滑；Jc 将控制点推离障碍。"},
        {"level": 1, "text": "Jd 限制速度、加速度和 jerk，使轨迹满足四旋翼动力学约束。"},
        {"level": 1, "text": "三个权重 λₛ、λᶜ、λᵈ 决定平滑性、安全距离和可飞性的取舍。"},
        {"level": 1, "text": "优化变量集中在 B-spline 控制点 Q，避免对密集采样轨迹点逐一优化。"},
    ], default_size=18.4, space_after=0, line_spacing=1.05)
    return slide


def optimizer_slide(prs):
    return table_with_notes_slide(
        prs,
        "Part. 02",
        "研究方法",
        "2.5",
        "L-BFGS 数值优化",
        "优化器性能直接影响局部重规划频率和高频控制循环中的稳定性。",
        [
            ["Optimizer", "Success", "Avg time", "Max time", "Eval Avg"],
            ["BB", "0.86", "0.50 ms", "1.14 ms", "268.2"],
            ["T-NEWTON", "0.62", "0.79 ms", "3.59 ms", "344.29"],
            ["L-BFGS", "0.89", "0.37 ms", "0.80 ms", "79.04"],
        ],
        [
            {"level": 1, "text": "L-BFGS 不显式求 Hessian，而是通过历史梯度近似二阶方向，适合小规模反复重启的局部问题。"},
            {"level": 1, "text": "表中 L-BFGS 同时给出最高成功率和最低平均耗时，因此成为实验部分默认优化器。"},
            {"level": 1, "text": "可导代价函数为 L-BFGS 提供稳定梯度，使控制点能够快速收敛到无碰撞轨迹。"},
        ],
        red_values=("L-BFGS", "0.89", "0.37 ms", "0.80 ms", "79.04"),
        font_size=14.0,
    )


def refinement_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.6", "时间重分配与轨迹修正")
    intro = add_box(slide, 0.62, 1.76, 12.05, 0.78)
    set_textbox(intro, [
        {"level": 0, "text": "优化后轨迹若违反动力学限制，EGO-Planner 会拉长时间分配并重新拟合曲线。"},
        {"level": 1, "text": "各向异性拟合使轨迹沿切向调整更灵活、径向偏移更保守，从而保持避障裕度。"},
    ], default_size=19.8, line_spacing=1.05)
    add_pic(slide, CROPS / "fig5_refinement.png", 0.90, 2.66, 4.34, 1.20)
    label(slide, "重分配与拟合公式", 0.90, 4.02, 4.00)
    equation_block(slide, [
        ("(14)", "rₑ = max{|Vᵢ,r/vₘ|, √|Aⱼ,r/aₘ|, ³√|Jₖ,r/jₘ|, 1}", 19.0),
        ("(15)", "Δt′ = rₑ Δt", 21.0),
        ("(16)", "min  J′(Q)=λₛJₛ + λᵈJᵈ + λᶠJᶠ", 20.6),
        ("(18)", "Jᶠ = ∫₀¹ [dₐ(αT′)²/a² + dᵣ(αT′)²/b²] dα", 19.0),
    ], 1.04, 4.38, 11.62, 1.56, size=20.0)
    add_rule(slide, 0.74, 6.16, 11.80)
    notes = add_box(slide, 0.78, 6.34, 11.82, 0.50)
    set_textbox(notes, [
        {"level": 1, "text": "rₑ 根据速度、加速度和 jerk 超限比例确定；拟合目标 Jf 保留原始安全轨迹形状，同时继续压平高阶导数。"},
    ], default_size=17.2, space_after=0, line_spacing=1.05)
    return slide


def algorithm_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.7", "Rebound Planning 流程")
    top = add_box(slide, 0.62, 1.76, 12.05, 0.78)
    set_textbox(top, [
        {"level": 0, "text": "Algorithm 2 将障碍信息补充、梯度优化和动力学修正组成在线闭环。"},
        {"level": 1, "text": "仅在碰撞仍存在时继续添加障碍信息，因此维护的是“必要障碍”集合。"},
    ], default_size=20.0, line_spacing=0.96)
    add_pic(slide, CROPS / "algo2_rebound.png", 3.76, 2.68, 5.60, 2.72)
    add_rule(slide, 0.74, 5.58, 11.80)
    bottom = add_box(slide, 0.78, 5.76, 11.82, 1.02)
    set_textbox(bottom, [
        {"level": 1, "text": "流程分成三步：FindInit 生成初始控制点，CheckAndAddObstacleInfo 记录新增碰撞约束，OneStepOptimize 通过 L-BFGS 更新控制点直到轨迹无碰撞。"},
        {"level": 1, "text": "若 IsFeasible 不满足，再进入 ReAllocateTime 与 CurveFittingOptimize 做动力学修正。"},
    ], default_size=17.8, space_after=0, line_spacing=1.05)
    return slide


def esdf_compare_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "研究结果分析")
    add_section_title(slide, "3.1", "ESDF-free 消融对比")
    top = add_box(slide, 0.62, 1.76, 12.05, 0.78)
    set_textbox(top, [
        {"level": 0, "text": "消融实验比较 EGO、ENI 和 EI 三种方式，核心差别是是否维护 ESDF、是否改进初始路径。"},
        {"level": 1, "text": "EGO 的成功率接近 EI，但总时间显著更低。"},
    ], default_size=20.0, line_spacing=0.96)
    add_table(slide, 0.82, 2.68, 11.70, 1.68, [
        ["Method", "Success", "Energy", "Vel Avg", "Plan", "ESDF", "Total"],
        ["EGO", "0.89", "49.92", "2.12", "0.37 ms", "/", "0.37 ms"],
        ["ENI", "0.69", "35.55", "2.09", "0.43 ms", "5.03 ms", "5.46 ms"],
        ["EI", "0.89", "42.27", "2.11", "0.48 ms", "5.07 ms", "5.55 ms"],
    ], font_size=13.2, red_values=("EGO", "0.89", "0.37 ms"))
    metric_row(slide, [
        ("15.0×", "总耗时下降", "5.55 ms → 0.37 ms"),
        ("0.89", "成功率保持", "EGO 与 EI 相同"),
        ("5.07 ms", "ESDF 维护成本", "EI 的主要时间来源"),
    ], 0.90, 4.74, 11.46, size=18.8)
    return slide


def planner_compare_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "研究结果分析")
    add_section_title(slide, "3.2", "与 EWOK / Fast-Planner 对比")
    top = add_box(slide, 0.62, 1.76, 12.05, 0.72)
    set_textbox(top, [
        {"level": 0, "text": "在相同地图和默认参数下，EGO-Planner 的路径更短、飞行时间更低，规划时间最低。"},
        {"level": 1, "text": "EGO-Planner 的轨迹能量高于 Fast-Planner，但仍保持可执行性和较高成功率。"},
    ], default_size=19.8, line_spacing=0.96)
    add_pic(slide, CROPS / "fig7_time_bar.png", 0.84, 2.56, 4.70, 1.74)
    add_table(slide, 5.92, 2.56, 6.42, 1.62, [
        ["Planner", "t(s)", "Length", "Energy", "tESDF", "tplan"],
        ["EWOK", "31.00", "59.05", "246.12", "6.43", "1.39"],
        ["Fast-Planner", "30.76", "45.18", "135.21", "4.01", "3.29"],
        ["EGO-Planner", "24.38", "42.24", "196.64", "/", "0.81"],
    ], font_size=12.6, red_values=("EGO-Planner", "24.38", "42.24", "0.81"))
    add_rule(slide, 0.74, 4.58, 11.80)
    add_pic(slide, ROOT / "ego-paper-test" / "extracted" / "page7-img1.jpeg", 1.38, 4.76, 10.48, 1.18)
    bottom = add_box(slide, 0.78, 6.16, 11.82, 0.56)
    set_textbox(bottom, [
        {"level": 1, "text": "绿色轨迹对应 EGO-Planner，路径长度 42.24 m，飞行时间 24.38 s，规划时间 0.81 ms。"},
    ], default_size=17.4, space_after=0, line_spacing=1.05)
    return slide


def simulation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "研究结果分析")
    add_section_title(slide, "3.3", "仿真轨迹对比")
    top = add_box(slide, 0.62, 1.76, 12.05, 0.78)
    set_textbox(top, [
        {"level": 0, "text": "仿真轨迹用于补充表格指标，直观展示不同规划器的路径形态。"},
        {"level": 1, "text": "EGO 的局部轨迹能绕开障碍密集区域，同时保持较短路径。"},
    ], default_size=20.0, line_spacing=0.96)
    add_pic(slide, CROPS / "fig8_sim_traj.png", 0.78, 2.58, 6.56, 3.00)
    add_pic(slide, ROOT / "ego-paper-test" / "extracted" / "page7-img1.jpeg", 7.56, 2.70, 4.74, 2.58)
    add_rule(slide, 0.74, 5.96, 11.82)
    bottom = add_box(slide, 0.76, 6.14, 11.90, 0.68)
    set_textbox(bottom, [
        {"level": 1, "text": "左侧展示三维仿真轨迹，右侧展示同一组对比方法的路径形态；绿色 EGO 轨迹更短且没有明显绕行。"},
    ], default_size=19.0, space_after=0, line_spacing=0.90)
    return slide


def real_world_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "研究结果分析")
    add_section_title(slide, "3.4", "实机飞行验证")
    top = add_box(slide, 0.62, 1.76, 12.05, 0.82)
    set_textbox(top, [
        {"level": 0, "text": "实机实验覆盖室内窄通道、随机目标追踪和森林场景，验证有限 FOV 下的在线规划能力。"},
        {"level": 1, "text": "室内杂乱环境中最高速度达到 3.56 m/s，系统仍能保持在线重规划。"},
    ], default_size=19.8, line_spacing=0.96)
    add_pic(slide, CROPS / "fig6_local_traj.png", 0.82, 2.62, 6.40, 2.98)
    add_pic(slide, ROOT / "ego-paper-test" / "extracted" / "page8-img7.jpeg", 7.36, 2.62, 4.94, 1.54)
    add_pic(slide, ROOT / "ego-paper-test" / "extracted" / "page8-img9.png", 7.36, 4.34, 4.94, 1.34)
    add_rule(slide, 0.74, 6.08, 11.80)
    bottom = add_box(slide, 0.78, 6.24, 11.82, 0.58)
    set_textbox(bottom, [
        {"level": 1, "text": "实机结果表明，在深度相机 FOV 受限时，轨迹能够随新障碍和新目标持续调整。"},
    ], default_size=18.8, space_after=0, line_spacing=0.90)
    return slide


def limitation_slide(prs):
    return split_slide(
        prs,
        "Part. 04",
        "结论与思考",
        "4.1",
        "优势与局限",
        "优势",
        [
            {"level": 0, "text": "主要收益来自降低地图维护开销，同时保留梯度优化的轨迹质量。"},
            {"level": 1, "text": "必要障碍集合使优化聚焦于当前轨迹附近的碰撞约束。"},
            {"level": 1, "text": "B-spline 控制点让连续轨迹问题变成低维控制点问题。"},
            {"level": 1, "text": "L-BFGS 和时间重分配保证了在线运行速度。"},
        ],
        "局限",
        [
            {"level": 0, "text": "方法仍然是局部规划器，受初始轨迹和局部环境影响。"},
            {"level": 1, "text": "A* 初始路径若落入不良拓扑，局部优化可能无法跳出。"},
            {"level": 1, "text": "动态障碍只在慢速场景中讨论，复杂交互仍需要预测模块。"},
            {"level": 1, "text": "时间重分配可能带来保守速度，需要更细的段级时间优化。"},
        ],
    )


def conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 04", "结论与思考")
    add_section_title(slide, "4.2", "结论归纳")
    box = add_box(slide, 0.66, 1.76, 12.04, 5.12)
    set_textbox(box, [
        {"level": 0, "text": "EGO-Planner 的核心做法是把 ESDF 查询换成必要障碍约束，降低在线规划计算量。"},
        {"level": 1, "text": "方法层：用 {p,v} 对直接构造碰撞距离和梯度，避免维护大范围 ESDF。"},
        {"level": 1, "text": "优化层：用平滑、碰撞、动力学三类惩罚统一更新 B-spline 控制点。"},
        {"level": 1, "text": "后处理层：通过时间重分配和各向异性拟合恢复动力学可行性。"},
        {"level": 1, "text": "结果层：规划时间约 0.81 ms，对比方法中最低；无 ESDF 版本总耗时 0.37 ms。"},
        {"level": 1, "text": "计算层面：从维护完整 ESDF 转为维护必要障碍集合。"},
        {"level": 1, "text": "优化层面：从连续曲线约束转为控制点上的可导代价。"},
        {"level": 1, "text": "系统层面：感知、重规划和动力学修正形成在线闭环。"},
        {"level": 1, "text": "应用层面，EGO 适合作为未知环境无人机系统中的快速局部重规划模块，并与全局拓扑规划、动态障碍预测和多机避碰配合。"},
    ], default_size=19.0, space_after=0.5, line_spacing=1.05)
    return slide


def outlook_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 05", "可借鉴方向")
    add_section_title(slide, "5.1", "对无人机规划系统的启发")
    box = add_box(slide, 0.70, 1.80, 12.00, 5.10)
    set_textbox(box, [
        {"level": 0, "text": "该方法可作为无人机系统中的快速局部重规划模块。"},
        {"level": 1, "text": "全局层给出航点或走廊，EGO-Planner 负责在局部未知障碍中生成可飞轨迹。"},
        {"level": 1, "text": "任务语义可转化为目标点、禁飞区、优先级和安全距离等规划约束。"},
        {"level": 1, "text": "规划器把这些语义约束转成局部目标和代价项，在线感知负责补充障碍信息。"},
        {"level": 1, "text": "ESDF-free 思路适用于资源受限无人机，可进一步降低地图维护负担。"},
        {"level": 1, "text": "系统分工：任务层输出目标、约束和优先级；全局层生成航点和粗路径；局部层由 EGO 实时避障。"},
        {"level": 1, "text": "可扩展方向包括动态障碍预测、多机避碰和全局拓扑切换。"},
    ], default_size=19.4, space_after=0.7, line_spacing=1.05)
    return slide


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    bullet_with_metrics_slide(prs, "Part. 01", "研究背景及动机", "1.1", "论文与问题", [
        {"level": 0, "text": "论文关注四旋翼在未知或局部已知环境中的在线局部轨迹规划。"},
        {"level": 1, "text": "传统 gradient-based 规划依赖 ESDF 提供距离和梯度。"},
        {"level": 1, "text": "轨迹优化只访问 ESDF 的局部子空间，而 ESDF 维护需要更新更大范围地图。"},
        {"level": 1, "text": "EGO-Planner 通过去除 ESDF 维护，降低计算开销并保持安全、平滑和动力学可行性。"},
        {"level": 2, "text": "论文场景包括室内窄通道、有限视场相机和森林飞行等实时任务。"},
    ], [
        ("RA-L 2020", "论文来源", "IEEE Robotics and Automation Letters"),
        ("Quadrotor", "对象", "局部轨迹规划与重规划"),
        ("ROS", "工程形态", "论文同步开源系统"),
    ],
        side_title="问题链路",
        side_notes=[
            {"level": 0, "text": "传统路线"},
            {"level": 1, "text": "深度图融合后维护局部 ESDF，再把距离和梯度交给优化器。"},
            {"level": 0, "text": "主要矛盾"},
            {"level": 1, "text": "ESDF 更新范围大，但优化只访问控制点附近的少量空间。"},
            {"level": 0, "text": "论文切入点"},
            {"level": 1, "text": "保留影响当前轨迹的障碍信息，直接生成避障梯度。"},
        ],
    )
    visual_slide(prs, "Part. 01", "研究背景及动机", "1.2", "为什么 ESDF 会成为瓶颈", [
        {"level": 0, "text": "图 1 对比了 ESDF 更新范围和轨迹优化覆盖空间。"},
        {"level": 1, "text": "完整局部 ESDF 的维护范围较大，其中大量区域不会直接影响最终轨迹。"},
    ], [
        (CROPS / "fig1_esdf_range.png", 1.00, 2.72, 5.80, 2.58),
        (CROPS / "fig2_local_minimum.png", 7.14, 2.68, 4.70, 2.42),
    ], [
        {"level": 1, "text": "左图：红色框为 ESDF 更新范围，紫色区域为轨迹优化实际覆盖空间。"},
        {"level": 1, "text": "右图：有限 FOV 下障碍背面不可见，ESDF 梯度也可能把轨迹推入局部极小。"},
        {"level": 1, "text": "方法转向碰撞触发的必要障碍提取，减少无关地图维护。"},
    ])
    bottleneck_slide(prs)
    technical_points_slide(prs)
    pv_representation_slide(prs)
    bspline_basis_slide(prs)
    formula_slide_distance(prs)
    formula_slide_objective(prs)
    optimizer_slide(prs)
    refinement_slide(prs)
    algorithm_slide(prs)
    esdf_compare_slide(prs)
    planner_compare_slide(prs)
    simulation_slide(prs)
    real_world_slide(prs)
    limitation_slide(prs)
    conclusion_slide(prs)
    outlook_slide(prs)
    thanks(prs)
    return prs


def main() -> None:
    OUT.parent.mkdir(exist_ok=True)
    prs = build()
    base.assert_layout(prs)
    prs.save(OUT)


if __name__ == "__main__":
    main()
