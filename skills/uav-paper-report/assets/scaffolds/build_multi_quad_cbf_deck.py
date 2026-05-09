from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

import build_rl_privileged_deck as base


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "out" / "uav-multi-quad-cbf-report.pptx"
CROPS = ROOT / "multi-quad-cbf-test" / "clean-crops"

SLIDE_W = 13.333
SLIDE_H = 7.5
BLUE = base.BLUE
BLACK = base.BLACK
GRAY = base.GRAY
MID_GRAY = base.MID_GRAY
RED = base.RED
LIGHT_BLUE = base.LIGHT_BLUE
LIGHT_RED = base.LIGHT_RED
BODY_FONT = "Times New Roman"

TERMS = [
    "CBF", "MPC", "Duality", "convex polytopes", "100%", "0%", "60%", "70%",
    "三架四旋翼", "六自由度", "对偶约束", "全局参考", "实机",
]


def add_box(slide, left, top, width, height):
    return base.add_box(slide, left, top, width, height)


def plain(box, text, size, *, bold=False, color=BLACK, align=PP_ALIGN.LEFT, allow_newlines=False):
    base.plain(box, text, size, bold=bold, color=color, align=align, allow_newlines=allow_newlines)


def set_textbox(box, lines, *, default_size=18.8, space_after=0.06, line_spacing=0.96):
    for item in lines:
        item.setdefault("terms", TERMS)
    base.set_textbox(box, lines, default_size=default_size, space_after=space_after, line_spacing=line_spacing)


def add_header(slide, part, title):
    base.add_header(slide, part, title)


def add_section_title(slide, num, title):
    base.add_section_title(slide, num, title)


def add_rule(slide, left, top, width):
    base.add_rule(slide, left, top, width)


def label(slide, text, left, top, width, *, color=BLUE):
    base.label(slide, text, left, top, width, color=color)


def add_pic(slide, name, left, top, width, height):
    path = CROPS / name
    with Image.open(path) as im:
        ratio = im.width / im.height
    box_ratio = width / height
    if ratio > box_ratio:
        w = width
        h = width / ratio
        x = left
        y = top + (height - h) / 2
    else:
        h = height
        w = height * ratio
        x = left + (width - w) / 2
        y = top
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def math_rows(slide, rows, left, top, width, row_h=0.48):
    base.equation_math_rows(slide, rows, left, top, width, row_h=row_h)


def table(slide, left, top, width, height, rows, *, font_size=13.0, red_values=(), green_values=(), col_widths=None):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = shape.table
    if col_widths:
        for idx, w in enumerate(col_widths):
            tbl.columns[idx].width = Inches(w)
    reds = set(red_values)
    greens = set(green_values)
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(value)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.015)
            cell.margin_bottom = Inches(0.015)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(242, 242, 242) if ri == 0 else RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(0)
            p.line_spacing = 0.95
            for run in p.runs:
                color = RED if value in reds else RGBColor(0, 140, 65) if value in greens else BLACK
                base.set_run_font(run, font_size, bold=ri == 0 or ci == 0 or value in reds or value in greens, color=color)
    return shape


def metric_row(slide, items, left, top, width, *, size=18.0):
    base.metric_row(slide, items, left, top, width, size=size)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if base.LOGO.exists():
        slide.shapes.add_picture(str(base.LOGO), Inches(4.10), Inches(0.64), width=Inches(5.08), height=Inches(1.41))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.40), Inches(SLIDE_W), Inches(3.10))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    title = add_box(slide, 0.42, 2.82, 12.50, 1.20)
    plain(title, "Optimal Trajectory Planning for Cooperative Manipulation\nwith Multiple Quadrotors Using CBFs", 28.6, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, allow_newlines=True)
    subtitle = add_box(slide, 0.0, 4.20, 13.33, 0.72)
    plain(subtitle, "Pallar, Li, Sarvaiya and Loianno / arXiv:2503.10695v2 / 2025", 18.2, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    date = add_box(slide, 5.50, 6.62, 2.35, 0.45)
    plain(date, "2026年5月8日", 18, color=BLUE, align=PP_ALIGN.CENTER)


def thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if base.LOGO.exists():
        slide.shapes.add_picture(str(base.LOGO), Inches(4.10), Inches(0.64), width=Inches(5.08), height=Inches(1.41))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.40), Inches(SLIDE_W), Inches(3.10))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    title = add_box(slide, 0.0, 3.34, 13.33, 0.96)
    plain(title, "谢谢！", 46, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    date = add_box(slide, 5.50, 6.62, 2.35, 0.45)
    plain(date, "2026年5月8日", 18, color=BLUE, align=PP_ALIGN.CENTER)


def body_slide(prs, part, part_title, sec, title, bullets, *, size=18.8):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, title)
    box = add_box(slide, 0.70, 1.78, 11.92, 5.42)
    set_textbox(box, bullets, default_size=size, space_after=0.14, line_spacing=0.99)


def background_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 01", "研究背景与问题")
    add_section_title(slide, "1.1", "多机缆绳搬运的规划难点")
    left = add_box(slide, 0.70, 1.78, 6.08, 2.52)
    set_textbox(left, [
        {"level": 0, "text": "论文研究多架四旋翼通过缆绳协同搬运刚体载荷，并在狭窄障碍环境中生成安全轨迹。"},
        {"level": 1, "text": "系统包含四旋翼、缆绳和六自由度载荷，碰撞对象不再只是单个机体。"},
        {"level": 1, "text": "载荷由缆绳张力间接驱动，规划载荷姿态会改变每根缆绳方向和四旋翼位置。"},
        {"level": 1, "text": "狭窄门洞要求系统整体旋转穿越，传统点质量避障难以覆盖载荷边角和缆绳。"},
    ], default_size=18.1, space_after=0.02)
    add_pic(slide, "fig1_real_gap.png", 7.20, 1.76, 4.74, 2.94)
    add_rule(slide, 0.76, 4.94, 11.80)
    metric_row(slide, [
        ("3 UAVs", "实机规模", "三机缆绳载荷"),
        ("6 DoF", "载荷状态", "位置与姿态同时规划"),
        ("100 Hz", "反馈频率", "Vicon 实验闭环"),
    ], 0.98, 5.30, 11.36, size=17.2)
    bottom = add_box(slide, 0.74, 6.62, 11.90, 0.48)
    set_textbox(bottom, [
        {"level": 1, "text": "本文将复杂外形、窄通道和动力学耦合统一到安全轨迹规划问题中，在狭窄环境中生成可执行轨迹。"},
    ], default_size=17.8, space_after=0.0, line_spacing=0.92)


def related_gap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 01", "研究背景与问题")
    add_section_title(slide, "1.2", "已有方法难以同时处理整体几何与动态可行性")
    table(slide, 0.76, 1.86, 11.82, 2.24, [
        ["方法类别", "常见假设", "主要不足", "本文处理方式"],
        ["RRT*/DMP", "简化障碍或固定队形", "仿真简单，难覆盖载荷-缆绳整体", "连续优化 + CBF 约束"],
        ["势场/避障控制", "偏重局部反应", "容易保守或陷入局部构型", "全局参考路径缓解局部极小"],
        ["CLF/CBF 控制", "只约束邻近机器人或点模型", "未覆盖载荷、缆绳和四旋翼整体", "多组件 convex polytope 建模"],
    ], font_size=12.8, red_values=("多组件 convex polytope 建模",), col_widths=[2.25, 2.70, 3.18, 3.69])
    add_rule(slide, 0.76, 4.42, 11.80)
    box = add_box(slide, 0.74, 4.68, 11.90, 1.46)
    set_textbox(box, [
        {"level": 0, "text": "本文将“整个搬运系统的最小距离”定义为 CBF 安全函数，而不是仅约束载荷中心或单架无人机。"},
        {"level": 1, "text": "CBF 给出前向不变性的安全约束，对偶形式把点到凸多面体距离转成可嵌入优化的问题。"},
        {"level": 1, "text": "A* 全局参考只提供粗路径，局部优化仍负责载荷姿态、缆绳方向和四旋翼位置。"},
    ], default_size=18.2, space_after=0.04)
    metric_row(slide, [
        ("全系统", "避障对象", "payload/cables/quadrotors"),
        ("CBF", "安全表达", "距离安全集"),
        ("A*", "参考路径", "避免局部极小"),
    ], 1.08, 6.20, 11.18, size=16.8)


def technical_points_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 01", "研究背景与问题")
    add_section_title(slide, "1.3", "技术要点")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.78)
    set_textbox(top, [
        {"level": 0, "text": "核心做法是把协同搬运的复杂几何、安全约束和非线性动力学统一到可在线求解的轨迹规划框架中。"},
    ], default_size=18.8, space_after=0.03)
    add_rule(slide, 0.76, 2.74, 11.80)
    labels = ["系统建模", "CBF 约束", "实验验证"]
    texts = [
        [
            {"level": 1, "text": "载荷、缆绳和四旋翼均用几何体表达，规划变量包含载荷位姿、速度、力和力矩。"},
            {"level": 1, "text": "缆绳张力通过 wrench mapping 联系到载荷动力学。"},
        ],
        [
            {"level": 1, "text": "使用 convex polytopes 表示载荷、障碍物和四旋翼-缆绳组合。"},
            {"level": 1, "text": "通过 Duality theorem 将距离 CBF 变为等价的优化不等式。"},
        ],
        [
            {"level": 1, "text": "在 6 个仿真环境中测试随机起终点，加入全局参考后成功率均为 100%。"},
            {"level": 1, "text": "三机实机通过狭窄门洞，轨迹跟踪显示载荷姿态调整有效。"},
        ],
    ]
    for i, title in enumerate(labels):
        x = 0.78 + 4.02 * i
        label(slide, title, x, 3.00, 2.80)
        box = add_box(slide, x, 3.40, 3.58, 1.66)
        set_textbox(box, texts[i], default_size=17.8, space_after=0.02, line_spacing=0.93)
    add_rule(slide, 0.76, 5.42, 11.80)
    bottom = add_box(slide, 0.74, 5.70, 11.92, 1.10)
    set_textbox(bottom, [
        {"level": 1, "text": "本文把规划目标从“到达目标点”扩展为“在优化过程中保持系统不进入障碍集合”。"},
        {"level": 1, "text": "载荷外形与缆绳几何直接进入安全约束，避免点模型或球模型漏掉边角碰撞。"},
    ], default_size=17.8, space_after=0.03)


def system_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.1", "系统约定与缆绳载荷模型")
    add_pic(slide, "fig2_system_convention.png", 0.82, 1.80, 5.00, 2.76)
    right = add_box(slide, 6.32, 1.78, 5.94, 2.52)
    set_textbox(right, [
        {"level": 0, "text": "系统由 n 架四旋翼通过缆绳牵引刚体载荷。"},
        {"level": 1, "text": "I、L、Bk 分别表示世界坐标系、载荷坐标系和第 k 架四旋翼机体系。"},
        {"level": 1, "text": "ρk 是载荷上的连接点，ξk 表示缆绳方向，μk 表示第 k 根缆绳张力。"},
        {"level": 1, "text": "载荷姿态改变会同时改变所有连接点和四旋翼的位置。"},
    ], default_size=17.8, space_after=0.02)
    add_rule(slide, 0.76, 4.72, 11.80)
    math_rows(slide, [
        {"label": "(1)", "size": 21.0, "pieces": ["X", ("L", "sub"), " = [ x", ("L", "sub"), ", q", ("L", "sub"), " ]", ("T", "sup"), " ,   V", ("L", "sub"), " = [ x", ("L", "sub"), "dot, Ω", ("L", "sub"), " ]", ("T", "sup")]},
        {"label": "(2)", "size": 21.0, "pieces": ["X", ("L", "sub"), "dot = [ x", ("L", "sub"), "dot, 1/2 q", ("L", "sub"), " ⊗ Ω", ("L", "sub"), " ]", ("T", "sup")]},
        {"label": "(3)", "size": 21.0, "pieces": ["m", ("L", "sub"), " x", ("L", "sub"), "ddot = F - m", ("L", "sub"), "g,    J", ("L", "sub"), " Ω", ("L", "sub"), "dot = M - Ω", ("L", "sub"), " × J", ("L", "sub"), "Ω", ("L", "sub")]},
    ], 1.00, 5.00, 7.76, row_h=0.48)
    explain = add_box(slide, 9.02, 4.92, 3.26, 1.54)
    set_textbox(explain, [
        {"level": 1, "text": "规划器实际控制的是载荷所受合力 F 与力矩 M。"},
        {"level": 1, "text": "缆绳张力再由力/力矩反推得到。"},
    ], default_size=17.8, space_after=0.01, line_spacing=0.92)


def wrench_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.2", "载荷 wrench 与缆绳张力映射")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.78)
    set_textbox(top, [
        {"level": 0, "text": "每根缆绳张力共同决定载荷上的合力与力矩，因此规划载荷运动等价于规划一组满足几何约束的张力。"},
    ], default_size=18.6, space_after=0.03)
    add_rule(slide, 0.76, 2.72, 11.80)
    math_rows(slide, [
        {"label": "(4)", "size": 22.0, "pieces": ["F = Σ μ", ("k", "sub"), " ,    M = Σ ρ", ("k", "sub"), " × R", ("L", "sub"), ("T", "sup"), " μ", ("k", "sub")]},
        {"label": "(5)", "size": 22.0, "pieces": ["F", ("L", "sup"), " = R", ("L", "sub"), ("T", "sup"), "F = Σ μ", ("k", "sub"), ("L", "sup")]},
        {"label": "(6)", "size": 22.0, "pieces": ["W", ("L", "sup"), " = [ F", ("L", "sup"), " ; M ] = P μ", ("L", "sup"), " ,    μ", ("L", "sup"), " = [ μ", ("1", "sub"), ("L", "sup"), " ... μ", ("n", "sub"), ("L", "sup"), " ]", ("T", "sup")]},
        {"label": "(8)", "size": 22.0, "pieces": ["μ", ("L", "sup"), " = P", ("T", "sup"), "(P P", ("T", "sup"), ")", ("-1", "sup"), " W", ("L", "sup")]},
    ], 0.98, 2.98, 8.82, row_h=0.50)
    right = add_box(slide, 9.96, 2.98, 2.46, 2.12)
    set_textbox(right, [
        {"level": 1, "text": "P 将所有缆绳张力映射到载荷六维 wrench。"},
        {"level": 1, "text": "张力方向决定四旋翼相对载荷的位置。"},
        {"level": 1, "text": "该耦合是避障约束必须覆盖缆绳与四旋翼的根本原因。"},
    ], default_size=17.8, space_after=0.01, line_spacing=0.91)
    add_rule(slide, 0.76, 5.32, 11.80)
    table(slide, 0.88, 5.56, 11.48, 1.22, [
        ["变量", "技术含义", "与安全规划的关系"],
        ["Wᒪ", "载荷坐标系下的合力与力矩", "决定载荷六自由度运动"],
        ["P", "缆绳张力到 wrench 的映射矩阵", "把每根缆绳耦合到载荷姿态"],
        ["μᒪ", "所有缆绳张力的堆叠向量", "反推出四旋翼与缆绳几何位置"],
    ], font_size=12.4, red_values=("P",), col_widths=[1.36, 4.56, 5.56])


def polytope_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.3", "几何包络：载荷、缆绳与无人机统一成凸多面体")
    add_pic(slide, "fig3_polytope.png", 0.82, 1.82, 5.26, 3.04)
    right = add_box(slide, 6.46, 1.78, 5.92, 2.94)
    set_textbox(right, [
        {"level": 0, "text": "论文用 convex polytopes 近似系统组件。"},
        {"level": 1, "text": "载荷在自身坐标系 L 中定义固定多面体。"},
        {"level": 1, "text": "每个四旋翼与对应缆绳被包络为一个长方体多面体。"},
        {"level": 1, "text": "障碍物也表达为线性不等式 Ay ≤ B，便于统一求距离。"},
    ], default_size=17.8, space_after=0.02)
    add_rule(slide, 0.76, 5.08, 11.80)
    math_rows(slide, [
        {"label": "(20)", "size": 20.5, "pieces": ["R", ("QC,k", "sub"), " = [ (e", ("1", "sub"), " × ξ", ("k", "sub"), ") × ξ", ("k", "sub"), " ,  -e", ("1", "sub"), " × ξ", ("k", "sub"), " ,  -ξ", ("k", "sub"), " ]"]},
        {"label": "", "size": 20.5, "pieces": ["x", ("QC,k", "sub"), " = x", ("L", "sub"), " + R", ("L", "sub"), "ρ", ("k", "sub"), " - l", ("k", "sub"), "ξ", ("k", "sub"), "/2"]},
    ], 0.98, 5.38, 7.94, row_h=0.50)
    note = add_box(slide, 9.22, 5.32, 3.10, 0.98)
    set_textbox(note, [
        {"level": 1, "text": "局部坐标系 QCk 沿缆绳方向建立，长方体包络覆盖缆绳和机体旋翼范围。"},
    ], default_size=17.8, space_after=0.01)


def optimization_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.4", "轨迹优化问题")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.78)
    set_textbox(top, [
        {"level": 0, "text": "规划器在固定预测时域内同时优化系统状态序列和输入序列。"},
    ], default_size=18.6, space_after=0.03)
    add_rule(slide, 0.76, 2.72, 11.80)
    math_rows(slide, [
        {"label": "(12a)", "size": 21.4, "pieces": ["min", ("X,U", "sub"), "  Σ", ("i=0", "sub"), ("N-1", "sup"), " h(X", ("i", "sub"), ", U", ("i", "sub"), ") + h", ("N", "sub"), "(X", ("N", "sub"), ")"]},
        {"label": "(12b)", "size": 21.4, "pieces": ["s.t.    X", ("i+1", "sub"), " = f(X", ("i", "sub"), ", U", ("i", "sub"), ")"]},
        {"label": "(12c)", "size": 21.4, "pieces": ["X", ("0", "sub"), " = X(t", ("0", "sub"), "),    g(X", ("i", "sub"), ", U", ("i", "sub"), ") ≤ 0"]},
        {"label": "(13)", "size": 21.4, "pieces": ["e", ("XN", "sub"), ("T", "sup"), "Q", ("XN", "sub"), "e", ("XN", "sub"), " + Σ(e", ("Xi", "sub"), ("T", "sup"), "Q", ("X", "sub"), "e", ("Xi", "sub"), " + e", ("Ui", "sub"), ("T", "sup"), "Q", ("U", "sub"), "e", ("Ui", "sub"), ")"]},
    ], 0.98, 3.04, 8.80, row_h=0.52)
    right = add_box(slide, 10.02, 3.04, 2.42, 2.16)
    set_textbox(right, [
        {"level": 1, "text": "状态 X 包含载荷位姿、速度、合力和力矩。"},
        {"level": 1, "text": "输入 U 是合力和力矩的时间导数。"},
        {"level": 1, "text": "目标项约束终端姿态，同时惩罚状态和输入偏差。"},
    ], default_size=17.8, space_after=0.01)
    add_rule(slide, 0.76, 5.64, 11.80)
    table(slide, 0.88, 5.88, 11.48, 1.04, [
        ["模块", "作用", "对搬运系统的意义"],
        ["状态 X", "载荷位姿、速度、合力和力矩", "让轨迹同时满足姿态与动力学"],
        ["输入 U", "合力和力矩的时间导数", "约束控制变化，避免张力突变"],
        ["A* 参考", "只给载荷位置粗路径", "局部优化补齐载荷姿态和缆绳张力"],
    ], font_size=12.2, red_values=("A* 参考",), col_widths=[1.58, 4.28, 5.62])


def cbf_distance_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.5", "距离 CBF：把障碍物安全距离写成前向不变条件")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.86)
    set_textbox(top, [
        {"level": 0, "text": "CBF 的安全函数 h(X) 用系统多面体与障碍物多面体之间的最小距离定义。"},
        {"level": 1, "text": "h(X) ≥ 0 表示当前系统组件距离障碍物至少 dsafe。"},
    ], default_size=18.3, space_after=0.02)
    add_rule(slide, 0.76, 2.84, 11.80)
    math_rows(slide, [
        {"label": "(14)", "size": 21.0, "pieces": ["h(X", ("i+1", "sub"), ") ≥ γ", ("i", "sub"), " h(X", ("i", "sub"), "),     0 ≤ γ", ("i", "sub"), " < 1"]},
        {"label": "(15a)", "size": 20.8, "pieces": ["h(X) = min", ("yL,yobs", "sub"), " || y", ("L", "sub"), " - y", ("obs", "sub"), " ||", ("2", "sub"), ("2", "sup"), " - d", ("safe", "sub"), ("2", "sup")]},
        {"label": "(15b)", "size": 20.8, "pieces": ["A", ("L", "sub"), "(X)y", ("L", "sub"), " ≤ B", ("L", "sub"), "(X),    A", ("obs", "sub"), "y", ("obs", "sub"), " ≤ B", ("obs", "sub")]},
    ], 0.98, 3.16, 8.62, row_h=0.62)
    right = add_box(slide, 9.74, 3.18, 2.82, 1.72)
    set_textbox(right, [
        {"level": 1, "text": "γi：裕度收缩率。"},
        {"level": 1, "text": "h(X)≥0：处于安全集。"},
        {"level": 1, "text": "递推：保持前向不变。"},
    ], default_size=17.0, space_after=0.0, line_spacing=0.90)
    add_rule(slide, 0.76, 5.18, 11.80)
    table(slide, 0.88, 5.46, 11.48, 1.26, [
        ["对象", "约束形式", "解释"],
        ["载荷", "hL(X) ≥ 0", "载荷多面体与障碍物保持 dsafe"],
        ["四旋翼-缆绳", "hQC,k(X) ≥ 0", "每根缆绳和机体外形也进入距离约束"],
        ["预测步", "h(Xi+1) ≥ γi h(Xi)", "安全裕度按 γi 收缩，轨迹保持在安全集合内"],
    ], font_size=12.3, red_values=("hQC,k(X) ≥ 0",), col_widths=[1.72, 3.10, 6.66])


def duality_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.6", "对偶约束：把最小距离计算嵌入轨迹优化")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.76)
    set_textbox(top, [
        {"level": 0, "text": "为了避免在优化中嵌套求最小距离，论文用 Duality theorem 推导等价 CBF 约束。"},
    ], default_size=18.5, space_after=0.03)
    add_rule(slide, 0.76, 2.70, 11.80)
    math_rows(slide, [
        {"label": "(16a)", "size": 20.8, "pieces": ["g(X) = max", ("λL,λobs", "sub"), " -λ", ("obs", "sup"), "B", ("obs", "sub"), " - λ", ("L", "sup"), "B", ("L", "sub"), "(X)"]},
        {"label": "(16b)", "size": 20.8, "pieces": ["λ", ("obs", "sup"), "A", ("obs", "sub"), " + λ", ("L", "sup"), "A", ("L", "sub"), "(X) = 0"]},
        {"label": "(16c)", "size": 20.8, "pieces": ["λ", ("obs", "sub"), " ≥ 0,   λ", ("L", "sub"), " ≥ 0,   ||λ", ("obs", "sup"), " * A", ("obs", "sub"), "||", ("2", "sub"), " ≤ 1"]},
        {"label": "(18a)", "size": 20.8, "pieces": ["-λ", ("i", "sub"), ("obs", "sup"), "B", ("obs", "sub"), " - λ", ("i", "sub"), ("L", "sup"), "B", ("L", "sub"), "(X", ("i+1", "sub"), ") ≥ γ", ("i", "sub"), "h(X", ("i", "sub"), ")"]},
    ], 0.98, 3.02, 8.82, row_h=0.48)
    right = add_box(slide, 10.02, 3.02, 2.42, 2.10)
    set_textbox(right, [
        {"level": 1, "text": "λ 为对偶变量，刻画多面体间距离。"},
        {"level": 1, "text": "凸问题满足强对偶，g(X)=h(X)。"},
        {"level": 1, "text": "转换后 CBF 进入 MPC 约束。"},
    ], default_size=17.8, space_after=0.01)
    add_rule(slide, 0.76, 5.36, 11.80)
    table(slide, 0.88, 5.60, 11.48, 1.08, [
        ["转换对象", "处理结果", "进入优化的形式"],
        ["多面体距离", "强对偶改写", "有限维 CBF 不等式"],
        ["系统几何", "载荷 / 缆绳 / 障碍物统一表达", "polytope 约束随状态更新"],
    ], font_size=12.8, red_values=("有限维 CBF 不等式",), col_widths=[2.28, 4.66, 4.54])


def exp_cbf_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.7", "指数 CBF 与松弛变量")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.84)
    set_textbox(top, [
        {"level": 0, "text": "指数 CBF 用初始安全函数 h(X0) 作为参考，并通过 αi 缓解严格约束导致的不可行。"},
        {"level": 1, "text": "当环境狭窄或初始状态接近障碍物时，松弛变量能让求解器找到可恢复的轨迹。"},
    ], default_size=18.1, space_after=0.02)
    add_rule(slide, 0.76, 2.84, 11.80)
    math_rows(slide, [
        {"label": "(19)", "size": 22.0, "pieces": ["-λ", ("i", "sub"), ("obs", "sup"), "B", ("obs", "sub"), " - λ", ("i", "sub"), ("L", "sup"), "B", ("L", "sub"), "(X", ("i+1", "sub"), ") ≥ α", ("i", "sub"), " Π", ("j=0", "sub"), ("i", "sup"), "γ", ("j", "sub"), "h(X", ("0", "sub"), ")"]},
        {"label": "", "size": 21.0, "pieces": ["α", ("i", "sub"), " 与状态、输入一起优化，h(X", ("0", "sub"), ") 每次 MPC 重规划时重新计算"]},
    ], 0.98, 3.34, 10.60, row_h=0.58)
    add_rule(slide, 0.76, 4.52, 11.80)
    table(slide, 0.88, 4.78, 11.48, 1.42, [
        ["项", "含义", "工程影响"],
        ["Πγj h(X0)", "以初始安全裕度为基准，随预测步逐渐收缩", "避免每一步都只看局部瞬时距离"],
        ["αi", "与状态和输入一起优化的松弛变量", "狭窄门洞附近减少不可行，安全裕度退化由 αi 反映"],
        ["X0 重算", "每次 MPC 重规划重新计算初始安全函数", "反馈状态变化后安全边界同步更新"],
    ], font_size=12.2, red_values=("αi",), col_widths=[1.40, 5.30, 4.78])
    bottom = add_box(slide, 0.74, 6.36, 11.90, 0.48)
    set_textbox(bottom, [
        {"level": 1, "text": "在载荷先旋转再穿越狭窄区域时，指数 CBF 给出从当前安全裕度逐步恢复的优化目标。"},
    ], default_size=17.8, space_after=0.0, line_spacing=0.92)


def algorithm_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.8", "在线规划流程")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.76)
    set_textbox(top, [
        {"level": 0, "text": "在线阶段以当前状态为初值，反复求解带 CBF 约束的有限时域优化问题。"},
    ], default_size=18.5, space_after=0.03)
    add_rule(slide, 0.76, 2.70, 11.80)
    nodes = [
        (["当前状态", "X(t0)"], 0.90, 3.18, 1.75),
        (["A* 参考", "载荷路径"], 3.00, 3.18, 1.80),
        (["MPC 优化", "X0...XN, U0...UN-1"], 5.20, 3.05, 2.22),
        (["CBF 约束", "多面体避障"], 7.94, 3.18, 1.92),
        (["执行首个输入", "重规划"], 10.28, 3.18, 1.74),
    ]
    for text, x, y, w in nodes:
        base.diagram_box(slide, text, x, y, w, 0.78, fill=LIGHT_BLUE, border=BLUE)
    for x1, x2 in [(2.65, 3.00), (4.80, 5.20), (7.42, 7.94), (9.86, 10.28)]:
        base.connector(slide, x1, 3.57, x2, 3.57)
    add_rule(slide, 0.76, 4.54, 11.80)
    table(slide, 0.88, 4.78, 11.48, 1.30, [
        ["步骤", "输入/输出", "作用"],
        ["状态反馈", "X(t0)", "固定优化初值，闭环更新预测问题"],
        ["障碍筛选", "距离小于 0.6 m 的活跃障碍", "减少无关 CBF 约束数量"],
        ["MPC 执行", "只执行第一个控制输入", "滚动重规划，吸收跟踪误差和状态变化"],
    ], font_size=12.3, red_values=("距离小于 0.6 m 的活跃障碍",), col_widths=[1.74, 4.76, 4.98])
    bottom = add_box(slide, 0.74, 6.26, 11.90, 0.58)
    set_textbox(bottom, [
        {"level": 1, "text": "论文报告单次求解约 9-15 s，当前实现对应轨迹生成和低频重规划；高速在线避障受求解速度限制。"},
    ], default_size=17.8, space_after=0.0, line_spacing=0.92)


def setup_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验结果")
    add_section_title(slide, "3.1", "仿真与实机设置")
    add_pic(slide, "fig5_block_diagram.png", 0.82, 1.84, 5.86, 2.06)
    right = add_box(slide, 7.18, 1.78, 5.08, 2.34)
    set_textbox(right, [
        {"level": 0, "text": "实验采用 RotorTM 仿真和三机实机验证。"},
        {"level": 1, "text": "实机空间为 10×6×4 m，Vicon 以 100 Hz 提供载荷、连接点和无人机反馈。"},
        {"level": 1, "text": "载荷为 196 g 三角形平台，由 3 架四旋翼通过缆绳连接。"},
    ], default_size=17.8, space_after=0.02)
    add_rule(slide, 0.76, 4.40, 11.80)
    metric_row(slide, [
        ("25 steps", "预测时域", "CasADi 优化"),
        ("0.6 m", "障碍筛选", "只保留近距离障碍"),
        ("7.5 cm", "实机膨胀", "补偿环境不确定性"),
    ], 1.02, 4.74, 11.24, size=17.2)
    bottom = add_box(slide, 0.74, 6.08, 11.90, 0.72)
    set_textbox(bottom, [
        {"level": 1, "text": "系统框图显示规划器输出期望载荷状态，控制器再分配缆绳张力并生成每架无人机的控制命令。"},
    ], default_size=17.8, space_after=0.01)


def sim_visual_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验结果")
    add_section_title(slide, "3.2", "仿真环境中的轨迹生成")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.72)
    set_textbox(top, [
        {"level": 0, "text": "仿真环境包含多种窄通道、扭曲通道和随机障碍，用于测试规划器的局部几何适应能力。"},
    ], default_size=18.3, space_after=0.02)
    add_pic(slide, "fig4_sim_envs.png", 0.78, 2.48, 6.86, 3.50)
    right = add_box(slide, 8.10, 2.64, 4.18, 2.96)
    set_textbox(right, [
        {"level": 0, "text": "轨迹显示载荷会主动旋转以通过狭窄区域。"},
        {"level": 1, "text": "蓝色障碍物附近，系统通过改变载荷姿态和缆绳方向保持整体安全距离。"},
        {"level": 1, "text": "该行为来自载荷外形、缆绳方向和四旋翼位置的联合规划，单点避障模型难以生成相同姿态变化。"},
    ], default_size=17.8, space_after=0.02)
    add_rule(slide, 0.76, 6.30, 11.80)
    bottom = add_box(slide, 0.74, 6.48, 11.90, 0.42)
    set_textbox(bottom, [
        {"level": 1, "text": "仿真结果表明，CBF 约束能在复杂几何中维持全系统碰撞安全，而不只是约束载荷中心。"},
    ], default_size=17.8, space_after=0.0, line_spacing=0.92)


def success_table_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验结果")
    add_section_title(slide, "3.3", "全局参考路径显著提升规划成功率")
    table(slide, 0.84, 1.82, 5.92, 2.70, [
        ["Environment", "No Global", "With Global"],
        ["ENV 1", "60%", "100%"],
        ["ENV 2", "60%", "100%"],
        ["ENV 3", "50%", "100%"],
        ["ENV 4", "70%", "100%"],
        ["ENV 5", "0%", "100%"],
        ["ENV 6", "20%", "100%"],
    ], font_size=13.8, red_values=("0%", "20%"), green_values=("100%",), col_widths=[1.90, 1.96, 2.06])
    right = add_box(slide, 7.26, 1.92, 5.04, 2.36)
    set_textbox(right, [
        {"level": 0, "text": "没有全局参考时，优化器容易陷入局部最优或不可恢复状态。"},
        {"level": 1, "text": "ENV 5 中无全局参考的成功率为 0%，局部优化未能找到绕行拓扑。"},
        {"level": 1, "text": "加入粗路径后，6 个环境全部达到 100% 成功率。"},
    ], default_size=17.8, space_after=0.02)
    add_rule(slide, 0.76, 4.72, 11.80)
    metric_row(slide, [
        ("6/6", "环境成功", "加入全局参考"),
        ("10", "随机任务", "每个环境起终点"),
        ("100%", "最高成功率", "所有环境均达到"),
    ], 1.04, 5.04, 11.20, size=17.8)
    bottom = add_box(slide, 0.74, 6.34, 11.90, 0.54)
    set_textbox(bottom, [
        {"level": 1, "text": "结果表明，CBF-MPC 主要负责局部安全与动态可行性；全局路径仍是避免拓扑死锁的必要条件。"},
    ], default_size=17.8, space_after=0.0)


def real_tracking_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验结果")
    add_section_title(slide, "3.4", "三机实机：载荷姿态调整穿越狭窄门洞")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.72)
    set_textbox(top, [
        {"level": 0, "text": "实机轨迹跟踪图记录载荷位置和姿态在穿越窄缝时的变化。"},
    ], default_size=18.4, space_after=0.02)
    add_pic(slide, "fig6_tracking.png", 0.82, 2.66, 8.46, 3.02)
    right = add_box(slide, 9.70, 2.70, 2.62, 2.54)
    set_textbox(right, [
        {"level": 1, "text": "约 0.05 s 处虚线对应载荷开始进入窄缝。"},
        {"level": 1, "text": "roll、pitch、yaw 均出现明显调整，用于让载荷姿态对齐通道。"},
        {"level": 1, "text": "位置跟踪误差整体可控，轨迹能够由下层控制器执行。"},
    ], default_size=17.8, space_after=0.01, line_spacing=0.91)
    add_rule(slide, 0.76, 5.94, 11.80)
    bottom = add_box(slide, 0.74, 6.20, 11.90, 0.72)
    set_textbox(bottom, [
        {"level": 1, "text": "实机结果表明，规划器生成了包含载荷姿态变化的避障轨迹，而不是仅让载荷中心通过门洞。"},
    ], default_size=17.8, space_after=0.0)


def result_interpret_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验结果")
    add_section_title(slide, "3.5", "实验结果的含义与边界")
    left = add_box(slide, 0.74, 1.82, 5.72, 2.00)
    set_textbox(left, [
        {"level": 0, "text": "仿真验证全局参考 + CBF-MPC 组合有效。"},
        {"level": 1, "text": "全局参考负责拓扑方向，避免狭窄环境中的局部极小。"},
        {"level": 1, "text": "CBF-MPC 负责在参考附近生成满足系统几何和动力学的安全轨迹。"},
    ], default_size=17.8, space_after=0.02)
    right = add_box(slide, 6.78, 1.82, 5.58, 2.00)
    set_textbox(right, [
        {"level": 0, "text": "实机验证方法具备物理可执行性。"},
        {"level": 1, "text": "系统能通过旋转载荷穿越窄缝，三架四旋翼协调调整位置。"},
        {"level": 1, "text": "实验仍依赖 Vicon 和较慢的重规划速度，离完全机载实时部署还有差距。"},
    ], default_size=17.8, space_after=0.02)
    add_rule(slide, 0.76, 3.72, 11.80)
    table(slide, 0.90, 4.02, 11.34, 1.46, [
        ["维度", "已验证", "仍需改进"],
        ["安全几何", "payload/cable/quadrotor 整体避障", "动态障碍与非凸障碍组合"],
        ["动力学", "载荷六自由度与缆绳张力映射", "更快求解与机载算力约束"],
        ["实验", "三机实机穿越窄缝", "去 Vicon、强扰动和长航时任务"],
    ], font_size=12.6, red_values=("更快求解与机载算力约束",), col_widths=[1.56, 4.54, 5.24])
    bottom = add_box(slide, 0.74, 5.82, 11.90, 0.66)
    set_textbox(bottom, [
        {"level": 1, "text": "本文完成了可认证轨迹规划框架的验证，机载实时部署仍受感知和求解速度限制。"},
    ], default_size=17.8, space_after=0.0)


def limitation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 04", "结论与思考")
    add_section_title(slide, "4.1", "方法局限")
    top = add_box(slide, 0.70, 1.78, 11.92, 0.74)
    set_textbox(top, [
        {"level": 0, "text": "主要局限来自求解速度、环境假设和感知闭环。"},
    ], default_size=18.8, space_after=0.02)
    add_rule(slide, 0.76, 2.72, 11.80)
    left = add_box(slide, 0.74, 2.98, 5.72, 2.08)
    set_textbox(left, [
        {"level": 1, "text": "论文报告 CasADi 单次求解约 9-15 s，当前速度对应轨迹生成与低频重规划，不满足高速在线避障。"},
        {"level": 1, "text": "障碍物用 convex polytopes 表达，复杂非凸环境需要分解或近似，可能带来约束数量增长。"},
        {"level": 1, "text": "CBF 约束能保持局部安全，但不能单独解决全局可达性和任务级死锁。"},
    ], default_size=17.8, space_after=0.02)
    right = add_box(slide, 6.78, 2.98, 5.62, 2.08)
    set_textbox(right, [
        {"level": 1, "text": "实机依赖 Vicon 反馈，尚未展示仅凭机载视觉估计载荷、缆绳和障碍物几何。"},
        {"level": 1, "text": "松弛变量提升可行性，但安全边界退化会反映为 αi 的变化。"},
        {"level": 1, "text": "动态障碍和强扰动尚未在该框架中完整闭环处理。"},
    ], default_size=17.8, space_after=0.02)
    add_rule(slide, 0.76, 5.34, 11.80)
    metric_row(slide, [
        ("9-15 s", "求解时间", "偏离高速实时控制"),
        ("Vicon", "实机反馈", "未完全机载感知"),
        ("CBF", "局部安全", "仍需全局参考路径"),
    ], 1.04, 5.66, 11.20, size=17.2)


def inspiration_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 04", "结论与思考")
    add_section_title(slide, "4.2", "方法可迁移性")
    top = add_box(slide, 0.70, 1.76, 11.92, 0.86)
    set_textbox(top, [
        {"level": 0, "text": "本文的建模方式是把系统外形直接写入安全约束。"},
        {"level": 1, "text": "对多无人机协同任务，危险往往发生在连接结构、载荷边角或队形外缘，而不是机体中心。"},
    ], default_size=18.3, space_after=0.02)
    add_rule(slide, 0.76, 2.86, 11.80)
    label(slide, "可迁移环节", 0.82, 3.12, 2.80)
    left = add_box(slide, 0.74, 3.48, 5.68, 1.72)
    set_textbox(left, [
        {"level": 1, "text": "无人机队形、载荷或吊挂物可用 convex hull / polytope 表达外轮廓。"},
        {"level": 1, "text": "CBF 可作为学习式规划器或采样规划器之后的安全约束层。"},
        {"level": 1, "text": "全局粗路径给出拓扑方向，局部优化负责安全与动力学可行性。"},
    ], default_size=17.8, space_after=0.02)
    label(slide, "工程限制", 6.92, 3.12, 2.80)
    right = add_box(slide, 6.78, 3.48, 5.68, 1.72)
    set_textbox(right, [
        {"level": 1, "text": "实时性受求解器速度影响，可通过 warm start、稀疏约束和并行求解改善。"},
        {"level": 1, "text": "机载感知误差会直接影响 h(X) 的安全判断。"},
        {"level": 1, "text": "松弛变量对应安全边界退化，应与故障检测或保守停机策略配合。"},
    ], default_size=17.8, space_after=0.02)
    add_rule(slide, 0.76, 5.50, 11.80)
    bottom = add_box(slide, 0.74, 5.76, 11.90, 0.96)
    set_textbox(bottom, [
        {"level": 1, "text": "对协同搬运、近距离编队和带载避障任务，该框架提供了系统级几何安全约束的建模路径。"},
        {"level": 1, "text": "工程验证重点包括约束构造、求解时间以及感知误差对安全函数 h(X) 的影响。"},
    ], default_size=17.8, space_after=0.01)


def summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 04", "结论与思考")
    add_section_title(slide, "4.3", "汇报总结")
    table(slide, 0.86, 1.82, 11.62, 2.92, [
        ["问题", "核心方法", "实验结论"],
        ["多四旋翼通过缆绳搬运载荷", "载荷动力学 + 缆绳张力映射", "三机实机能穿越窄缝"],
        ["系统整体避障而非点避障", "convex polytopes + CBF 距离约束", "仿真中全局参考后成功率 100%"],
        ["复杂几何导致优化困难", "Duality theorem 转换距离约束", "仍受求解速度和感知闭环限制"],
    ], font_size=13.8, red_values=("成功率 100%",), col_widths=[3.35, 4.16, 4.11])
    add_rule(slide, 0.76, 5.08, 11.80)
    box = add_box(slide, 0.74, 5.36, 11.90, 1.42)
    set_textbox(box, [
        {"level": 0, "text": "本文把协同搬运系统的几何外形和动力学耦合写入 CBF-MPC，使载荷、缆绳和四旋翼作为整体穿越障碍。"},
        {"level": 1, "text": "方法重点集中在安全约束建模，工程部署受求解速度、感知可靠性和实时闭环实现限制。"},
    ], default_size=18.0, space_after=0.03)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    cover(prs)
    background_slide(prs)
    related_gap_slide(prs)
    technical_points_slide(prs)
    system_model_slide(prs)
    wrench_slide(prs)
    polytope_slide(prs)
    optimization_slide(prs)
    cbf_distance_slide(prs)
    duality_slide(prs)
    exp_cbf_slide(prs)
    algorithm_slide(prs)
    setup_slide(prs)
    sim_visual_slide(prs)
    success_table_slide(prs)
    real_tracking_slide(prs)
    result_interpret_slide(prs)
    limitation_slide(prs)
    inspiration_slide(prs)
    summary_slide(prs)
    thanks(prs)
    return prs


def assert_layout(prs: Presentation) -> None:
    base.assert_layout(prs)


def main() -> None:
    OUT.parent.mkdir(exist_ok=True)
    prs = build()
    assert_layout(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
