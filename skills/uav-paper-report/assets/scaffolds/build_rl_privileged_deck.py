from __future__ import annotations

import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "out" / "uav-rl-privileged-report.pptx"
CROPS = ROOT / "rl-privileged-test" / "clean-crops"
LOGO = ROOT / "header-logo.png"

SLIDE_W = 13.333
SLIDE_H = 7.5
BODY_FONT = "Times New Roman"
MATH_FONT = "Times New Roman"

BLUE = RGBColor(0, 121, 192)
HEADER_GRAY = RGBColor(242, 242, 242)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(78, 78, 78)
MID_GRAY = RGBColor(155, 155, 155)
LIGHT_BLUE = RGBColor(226, 241, 251)
LIGHT_RED = RGBColor(255, 238, 238)
RED = RGBColor(220, 0, 0)

KEY_TERMS = [
    "ToA", "Yaw", "GRU", "APG", "BPTT", "FMM",
    "86%", "34%", "36%", "589 m", "20 次", "4 m/s",
    "特权信息", "域随机化", "无碰撞", "可微动力学",
]


def add_box(slide, left: float, top: float, width: float, height: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.margin_left = Inches(0)
    box.text_frame.margin_right = Inches(0)
    box.text_frame.margin_top = Inches(0)
    box.text_frame.margin_bottom = Inches(0)
    return box


def set_shape_name(shape, name: str) -> None:
    shape.name = name


def set_run_font(run, size: float, *, bold: bool = False, color: RGBColor = BLACK, font: str = BODY_FONT) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        element = r_pr.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            r_pr.append(element)
        element.set("typeface", font)


def compact_text(text: str) -> str:
    replacements = {
        "learning-based": "learning\u2011based",
        "end-to-end": "end\u2011to\u2011end",
        "time-of-arrival": "time\u2011of\u2011arrival",
        "point-mass": "point\u2011mass",
        "sim-to-real": "sim\u2011to\u2011real",
        "body rate": "body\u00a0rate",
        "86%": "86%",
        "34%": "34%",
        "36%": "36%",
        "589 m": "589\u00a0m",
        "4 m/s": "4\u00a0m/s",
        "50 Hz": "50\u00a0Hz",
        "60 Hz": "60\u00a0Hz",
        "1.7 kg": "1.7\u00a0kg",
        "1.15g": "1.15g",
        "1.3g": "1.3g",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    text = re.sub(r"(?<=\d)\s+(?=(ms|s|m|m/s|Hz|kg|%)\b)", "\u00a0", text)
    text = re.sub(r"(?<=[A-Za-z0-9+\-/])\s+(?=[\u4e00-\u9fff])", "", text)
    text = re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[A-Za-z0-9])", "", text)
    return text


def split_terms(text: str, terms: list[str]) -> list[tuple[str, bool]]:
    terms = sorted({t for t in terms if t}, key=len, reverse=True)
    out: list[tuple[str, bool]] = []
    i = 0
    while i < len(text):
        hit = next((term for term in terms if text.startswith(term, i)), None)
        if hit:
            out.append((hit, True))
            i += len(hit)
            continue
        j = i + 1
        while j < len(text) and not any(text.startswith(term, j) for term in terms):
            j += 1
        out.append((text[i:j], False))
        i = j
    return out


def add_runs(paragraph, text: str, size: float, *, bold: bool = False, terms: list[str] | None = None) -> None:
    for piece, marked in split_terms(compact_text(text), KEY_TERMS + (terms or [])):
        run = paragraph.add_run()
        run.text = piece
        set_run_font(run, size, bold=bold or marked, color=RED if marked else BLACK)


def plain(box, text: str, size: float, *, bold: bool = False, color: RGBColor = BLACK, align=PP_ALIGN.LEFT) -> None:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        r = p.add_run()
        r.text = line
        set_run_font(r, size, bold=bold, color=color)


def set_textbox(
    box,
    lines: list[dict],
    *,
    default_size: float = 18.8,
    space_after: float = 0.08,
    line_spacing: float = 0.98,
) -> None:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for idx, item in enumerate(lines):
        text = item["text"].replace("\n", " ")
        level = item.get("level", 0)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(item.get("space_after", space_after))
        ppr = p._p.get_or_add_pPr()
        if level == 0:
            ppr.set("marL", "210000")
            ppr.set("indent", "-140000")
            marker = item.get("marker", "● ")
            size = float(item.get("size", default_size))
        elif level == 1:
            ppr.set("marL", "350000")
            ppr.set("indent", "-120000")
            marker = item.get("marker", "• ")
            size = float(item.get("size", default_size - 1.0))
        else:
            ppr.set("marL", "560000")
            ppr.set("indent", "-105000")
            marker = item.get("marker", "– ")
            size = float(item.get("size", default_size - 2.0))
        add_runs(p, marker + text, size, bold=item.get("bold", level == 0), terms=item.get("terms", []))


def add_header(slide, part: str, title: str) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.87))
    band.fill.solid()
    band.fill.fore_color.rgb = HEADER_GRAY
    band.line.fill.background()
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.49), Inches(0.06), width=Inches(2.65), height=Inches(0.74))
    box = add_box(slide, 3.25, 0.11, 8.1, 0.62)
    box.text_frame.word_wrap = False
    plain(box, f"{part}  {title}", 28, bold=True)


def add_section_title(slide, num: str, title: str) -> None:
    box = add_box(slide, 0.70, 1.12, 9.2, 0.54)
    box.text_frame.word_wrap = False
    plain(box, f"{num} {title}", 20.5, bold=True)


def add_rule(slide, left: float, top: float, width: float) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left), Inches(top), Inches(left + width), Inches(top))
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(1.0)


def add_pic(slide, name: str, left: float, top: float, width: float, height: float):
    path = CROPS / name
    with Image.open(path) as im:
        ratio = im.width / im.height
    box_ratio = width / height
    if ratio > box_ratio:
        w = width
        h = width / ratio
        x = left
        y = top + (height - h) / 2
    else:
        h = height
        w = height * ratio
        x = left + (width - w) / 2
        y = top
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    return pic


def label(slide, text: str, left: float, top: float, width: float, *, color: RGBColor = BLUE) -> None:
    box = add_box(slide, left, top, width, 0.32)
    plain(box, text, 16.8, bold=True, color=color)


def equation_rows(slide, rows: list[dict], left: float, top: float, width: float, row_h: float = 0.42) -> None:
    for idx, row in enumerate(rows):
        y = top + idx * row_h
        lab = add_box(slide, left, y, 0.70, row_h * 0.86)
        set_shape_name(lab, "MATH_LABEL")
        plain(lab, row.get("label", ""), row.get("size", 20.0), bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
        body = add_box(slide, left + 0.86, y, width - 0.86, row_h * 0.86)
        set_shape_name(body, "MATH_BODY")
        body.text_frame.word_wrap = False
        plain(body, row["text"], row.get("size", 20.0), color=BLACK)


def set_baseline(run, value: int) -> None:
    run._r.get_or_add_rPr().set("baseline", str(value))


def math_run(paragraph, text: str, size: float, *, mode: str | None = None, bold: bool = False) -> None:
    run = paragraph.add_run()
    run.text = text
    run_size = size * 0.72 if mode in {"sub", "sup"} else size
    set_run_font(run, run_size, bold=bold, color=BLACK, font=MATH_FONT)
    if mode == "sub":
        set_baseline(run, -25000)
    elif mode == "sup":
        set_baseline(run, 30000)


def equation_math_rows(slide, rows: list[dict], left: float, top: float, width: float, row_h: float = 0.44) -> None:
    for idx, row in enumerate(rows):
        y = top + idx * row_h
        size = float(row.get("size", 20.0))
        lab = add_box(slide, left, y, 0.72, row_h * 0.86)
        set_shape_name(lab, "MATH_LABEL")
        plain(lab, row.get("label", ""), size, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
        body = add_box(slide, left + 0.86, y, width - 0.86, row_h * 0.86)
        set_shape_name(body, "MATH_BODY")
        tf = body.text_frame
        tf.clear()
        tf.word_wrap = False
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        for piece in row["pieces"]:
            if isinstance(piece, str):
                math_run(p, piece, size)
            else:
                text, mode = piece[0], piece[1]
                math_run(p, text, size, mode=mode)


def small_table(slide, left: float, top: float, width: float, height: float, rows: list[list[str]], *, font_size: float = 12.8):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = value
            cell.margin_left = Inches(0.025)
            cell.margin_right = Inches(0.025)
            cell.margin_top = Inches(0.012)
            cell.margin_bottom = Inches(0.012)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(242, 242, 242) if ri == 0 else RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(0)
            for run in p.runs:
                color = RED if value in {"Ours", "86%", "0", "589 m", "20"} else BLACK
                set_run_font(run, font_size, bold=ri == 0 or ci == 0 or value in {"Ours", "86%"}, color=color)
    return shape


def metric_row(slide, items: list[tuple[str, str, str]], left: float, top: float, width: float, *, size: float = 18.0) -> None:
    col_w = width / len(items)
    for idx, (value, name, note) in enumerate(items):
        x = left + idx * col_w
        value_box = add_box(slide, x, top, col_w - 0.12, 0.43)
        plain(value_box, value, size + 5, bold=True, color=RED, align=PP_ALIGN.CENTER)
        name_box = add_box(slide, x, top + 0.44, col_w - 0.12, 0.30)
        plain(name_box, name, size - 2.0, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        note_box = add_box(slide, x, top + 0.77, col_w - 0.12, 0.36)
        plain(note_box, note, size - 3.6, color=BLACK, align=PP_ALIGN.CENTER)


def diagram_box(slide, text: str, left: float, top: float, width: float, height: float, *, fill=LIGHT_BLUE, border=BLUE):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    set_shape_name(shp, "DIAG_BOX")
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1.0)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(0)
    p.line_spacing = 0.95
    for idx, part in enumerate(text.split("\n")):
        if idx:
            p.add_line_break()
        r = p.add_run()
        r.text = part
        set_run_font(r, 14.3 if idx else 16.0, bold=idx == 0)
    return shp


def connector(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    set_shape_name(line, "DIAG_CONN")
    line.line.color.rgb = GRAY
    line.line.width = Pt(1.15)
    line.line.end_arrowhead = True


def cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.10), Inches(0.64), width=Inches(5.08), height=Inches(1.41))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.40), Inches(SLIDE_W), Inches(3.10))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    title = add_box(slide, 0.44, 2.82, 12.46, 1.14)
    plain(title, "Quadrotor Navigation using Reinforcement Learning\nwith Privileged Information", 29.0, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    subtitle = add_box(slide, 0.0, 4.18, 13.33, 0.74)
    plain(subtitle, "Lee et al. / arXiv:2509.08177v2 / 2025", 18.2, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    date = add_box(slide, 5.50, 6.62, 2.35, 0.45)
    plain(date, "2026年5月8日", 18, color=BLUE, align=PP_ALIGN.CENTER)


def thanks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.10), Inches(0.64), width=Inches(5.08), height=Inches(1.41))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.40), Inches(SLIDE_W), Inches(3.10))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    title = add_box(slide, 0.0, 3.34, 13.33, 0.96)
    plain(title, "谢谢！", 46, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    date = add_box(slide, 5.50, 6.62, 2.35, 0.45)
    plain(date, "2026年5月8日", 18, color=BLUE, align=PP_ALIGN.CENTER)


def body_slide(prs, part, part_title, sec, title, bullets, *, size=18.8):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, title)
    box = add_box(slide, 0.72, 1.78, 11.92, 5.36)
    set_textbox(box, bullets, default_size=size, space_after=0.16, line_spacing=1.0)
    return slide


def problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 01", "研究背景与问题")
    add_section_title(slide, "1.1", "端到端无人机导航的瓶颈")
    top = add_box(slide, 0.72, 1.78, 11.86, 1.18)
    set_textbox(top, [
        {"level": 0, "text": "论文面向复杂户外环境中的四旋翼视觉导航，目标是在大障碍、急转弯和死胡同场景中实现低延迟自主飞行。"},
        {"level": 1, "text": "端到端策略能减少传统感知、规划、控制分模块串联带来的延迟，但局部反应式策略容易被大墙体或凹形障碍困住。"},
    ], default_size=18.9, space_after=0.05)
    add_rule(slide, 0.76, 3.18, 11.80)
    for title, x in [("局部感知不足", 0.86), ("航向受限", 4.58), ("实机迁移误差", 8.30)]:
        label(slide, title, x, 3.38, 3.05)
    cols = [
        (0.74, 3.76, [
            {"level": 1, "text": "深度图只提供当前视野内几何，ESDF 等局部代价难以给出绕行方向。"},
            {"level": 1, "text": "在凹形区域中，单纯避障损失会使策略停滞或选择过近路径。"},
        ]),
        (4.46, 3.76, [
            {"level": 1, "text": "固定朝向目标会让机体持续看向被遮挡方向，转弯时难以建立稳定视觉输入。"},
            {"level": 1, "text": "论文新增 Yaw 对齐损失，使机头朝向与实际运动方向一致。"},
        ]),
        (8.18, 3.76, [
            {"level": 1, "text": "训练采用 point-mass 动力学，实机需要跟踪姿态、推力和 body rate。"},
            {"level": 1, "text": "质量、推力系数和电压变化会产生稳态推力偏差，需要域随机化补偿。"},
        ]),
    ]
    for x, top_y, lines in cols:
        box = add_box(slide, x, top_y, 3.46, 1.58)
        set_textbox(box, lines, default_size=16.9, space_after=0.02, line_spacing=0.93)
    add_rule(slide, 0.76, 5.62, 11.80)
    metric_row(slide, [
        ("86%", "仿真成功率", "11 个 OOD 环境"),
        ("34%", "优于基线", "摘要报告提升幅度"),
        ("589 m", "实机飞行", "20 次无碰撞"),
    ], 0.96, 5.90, 11.36, size=18.0)


def related_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 01", "研究背景与问题")
    add_section_title(slide, "1.2", "与已有方法的技术差异")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.78)
    set_textbox(top, [{"level": 0, "text": "已有端到端飞行策略主要解决窄障碍快速通过，对需要绕行的大障碍和迷宫式区域处理不足。"}], default_size=18.8)
    add_rule(slide, 0.76, 2.74, 11.80)
    rows = [
        ["方法", "训练信号", "部署输入", "主要不足"],
        ["Zhang et al.", "可微物理 + 局部损失", "深度图 + 状态", "航向固定，绕大障碍困难"],
        ["YOPO", "ESDF 引导学习", "深度图 + primitive", "偏局部避障，缺少全局可达方向"],
        ["本文", "ToA 特权信息 + Yaw 损失", "深度图 + 目标 + 状态", "长时回退场景仍然困难"],
    ]
    small_table(slide, 0.78, 3.02, 11.78, 1.70, rows, font_size=13.3)
    add_rule(slide, 0.76, 5.08, 11.80)
    bottom = add_box(slide, 0.78, 5.34, 11.72, 1.34)
    set_textbox(bottom, [
        {"level": 1, "text": "本文的核心区别是训练时引入全局路径方向，推理时保持反应式低延迟控制。"},
        {"level": 1, "text": "ToA map 只提供训练监督；实机部署不显式维护地图，因此计算链路仍然接近端到端策略。"},
        {"level": 1, "text": "Yaw 预测补足视觉朝向，ToA 梯度补足绕行方向，两者共同解决大障碍导航。"},
    ], default_size=17.2, space_after=0.03, line_spacing=0.96)


def technical_points_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 01", "研究背景与问题")
    add_section_title(slide, "1.3", "技术要点")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.76)
    set_textbox(top, [{"level": 0, "text": "论文把可微仿真、航向预测、ToA 特权信息和实机控制补偿组合成端到端四旋翼导航训练框架。"}], default_size=18.8)
    add_rule(slide, 0.76, 2.76, 11.80)
    headers = ["航向与控制输出", "ToA 特权监督", "实机可迁移性"]
    lines = [
        [
            {"level": 1, "text": "策略输出质量归一化推力 t_k 和预测航向 ψ_k。"},
            {"level": 1, "text": "Yaw 对齐损失让机体在绕障时朝向运动方向，改善急转弯场景。"},
        ],
        [
            {"level": 1, "text": "ToA map 在训练期提供最短到达时间梯度，诱导策略学习绕开凹形障碍。"},
            {"level": 1, "text": "部署阶段不需要 ToA map，策略仍由深度图与状态直接输出动作。"},
        ],
        [
            {"level": 1, "text": "使用 body rate 姿态控制弥合 point-mass 训练与刚体实机之间的差异。"},
            {"level": 1, "text": "通过重力与状态噪声随机化，使策略适应推力建模误差。"},
        ],
    ]
    for i, head in enumerate(headers):
        x = 0.78 + i * 3.98
        label(slide, head, x, 3.02, 3.3)
        box = add_box(slide, x, 3.42, 3.58, 1.72)
        set_textbox(box, lines[i], default_size=16.8, space_after=0.02, line_spacing=0.92)
    add_rule(slide, 0.76, 5.48, 11.80)
    metric_row(slide, [
        ("50 Hz", "策略运行频率", "机载 Orin NX"),
        ("60 Hz", "深度输入", "RealSense D456"),
        ("1.7 kg", "实机平台", "自研四旋翼"),
    ], 1.04, 5.78, 11.20, size=17.5)


def overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.1", "方法闭环")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.70)
    set_textbox(top, [{"level": 0, "text": "训练闭环把可微动力学展开为时间序列，通过损失函数对策略参数反向传播；部署闭环只保留神经策略和低层控制器。"}], default_size=18.6)
    y = 3.02
    diagram_box(slide, "深度图\n64×64", 0.76, y, 1.58, 0.74)
    diagram_box(slide, "目标信息\nv_goal, d_goal", 0.76, y + 1.04, 1.58, 0.74)
    diagram_box(slide, "状态估计\nv, R", 0.76, y + 2.08, 1.58, 0.74)
    diagram_box(slide, "特征抽取\nLayerNorm", 3.04, y + 0.78, 1.68, 0.92)
    diagram_box(slide, "GRU 记忆\nh_k", 5.26, y + 0.78, 1.52, 0.92)
    diagram_box(slide, "动作输出\nt_k, ψ_k", 7.30, y + 0.78, 1.54, 0.92, fill=LIGHT_RED, border=RED)
    diagram_box(slide, "可微动力学\ns_{k+1}=f(s_k,u_k)", 9.34, y + 0.78, 2.10, 0.92)
    diagram_box(slide, "损失累积\nL = Σ λ_i L_i", 5.76, y + 2.15, 2.20, 0.82, fill=RGBColor(248, 248, 248), border=GRAY)
    for yy in [y + 0.37, y + 1.41, y + 2.45]:
        connector(slide, 2.34, yy, 3.04, y + 1.24)
    connector(slide, 4.72, y + 1.24, 5.26, y + 1.24)
    connector(slide, 6.78, y + 1.24, 7.30, y + 1.24)
    connector(slide, 8.84, y + 1.24, 9.34, y + 1.24)
    connector(slide, 10.38, y + 1.70, 7.92, y + 2.15)
    connector(slide, 6.86, y + 2.15, 6.02, y + 1.70)
    add_rule(slide, 0.76, 6.18, 11.80)
    bottom = add_box(slide, 0.78, 6.36, 11.72, 0.72)
    set_textbox(bottom, [
        {"level": 1, "text": "训练期引入 ToA 梯度、避障、平滑和航向损失；部署期不显式计算地图，策略直接输出推力和航向。"},
    ], default_size=17.2, space_after=0.0)


def dynamics_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.2", "可微 point-mass 动力学")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.74)
    set_textbox(top, [{"level": 0, "text": "策略训练采用 velocity Verlet 积分的 point-mass 模型，用较低计算代价保留对动作的可微梯度。"}], default_size=18.7)
    equation_math_rows(slide, [
        {"label": "(1)", "size": 21.0, "pieces": ["p", ("k+1", "sub"), " = p", ("k", "sub"), " + v", ("k", "sub"), " Δt + 1/2 a", ("k", "sub"), "(Δt)", ("2", "sup")]},
        {"label": "(2)", "size": 21.0, "pieces": ["v", ("k+1", "sub"), " = v", ("k", "sub"), " + (a", ("k", "sub"), " + a", ("k+1", "sub"), ") Δt / 2"]},
        {"label": "(3)", "size": 21.0, "pieces": ["a", ("k+1", "sub"), " = t", ("k+1", "sub"), " − [0, 0, g]", ("T", "sup")]},
    ], 0.96, 2.72, 11.40, row_h=0.50)
    add_rule(slide, 0.76, 4.36, 11.80)
    left = add_box(slide, 0.78, 4.60, 5.84, 1.66)
    set_textbox(left, [
        {"level": 1, "text": "p_k、v_k、a_k 分别表示位置、速度和加速度，动作中的 t_k 是质量归一化推力。"},
        {"level": 1, "text": "简化模型缩短训练时间，同时保留 APG 对策略参数的直接梯度。"},
    ], default_size=17.1, space_after=0.02, line_spacing=0.94)
    right = add_box(slide, 6.92, 4.60, 5.52, 1.66)
    set_textbox(right, [
        {"level": 1, "text": "模型不显式积分姿态，机体 z 轴由推力方向确定，x 轴由预测 Yaw 确定。"},
        {"level": 1, "text": "这种表示让策略同时学习推进方向和相机朝向。"},
    ], default_size=17.1, space_after=0.02, line_spacing=0.94)
    add_rule(slide, 0.76, 6.46, 11.80)
    bottom = add_box(slide, 0.78, 6.62, 11.72, 0.42)
    set_textbox(bottom, [{"level": 1, "text": "因此，训练得到的策略既有可微优化效率，又保留了对机体朝向的显式控制入口。"}], default_size=16.8, space_after=0.0)


def network_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.3", "策略网络结构")
    top = add_box(slide, 0.72, 1.78, 11.88, 0.76)
    set_textbox(top, [{"level": 0, "text": "网络把视觉、目标和状态编码到同一 192 维潜变量，再用 GRU 维持时序一致性并输出控制动作。"}], default_size=18.7)
    y = 3.08
    diagram_box(slide, "Depth\n64×64", 0.78, y, 1.28, 0.70)
    diagram_box(slide, "Conv + Pool\n视觉特征", 2.62, y, 1.64, 0.70)
    diagram_box(slide, "Target\n速度+距离", 0.78, y + 1.10, 1.28, 0.70)
    diagram_box(slide, "Linear\n目标特征", 2.62, y + 1.10, 1.64, 0.70)
    diagram_box(slide, "State\n速度+姿态", 0.78, y + 2.20, 1.28, 0.70)
    diagram_box(slide, "Linear\n状态特征", 2.62, y + 2.20, 1.64, 0.70)
    diagram_box(slide, "求和 + LayerNorm\n(B, 192)", 4.96, y + 1.04, 1.90, 0.82)
    diagram_box(slide, "GRUCell\nh_{k-1}→h_k", 7.42, y + 1.04, 1.72, 0.82)
    diagram_box(slide, "Linear 输出\n[t_k, ψ_k]", 9.78, y + 1.04, 1.72, 0.82, fill=LIGHT_RED, border=RED)
    for yy in [y + 0.35, y + 1.45, y + 2.55]:
        connector(slide, 2.06, yy, 2.62, yy)
        connector(slide, 4.26, yy, 4.96, y + 1.45)
    connector(slide, 6.86, y + 1.45, 7.42, y + 1.45)
    connector(slide, 9.14, y + 1.45, 9.78, y + 1.45)
    add_rule(slide, 0.76, 6.10, 11.80)
    bottom = add_box(slide, 0.78, 6.30, 11.70, 0.80)
    set_textbox(bottom, [
        {"level": 1, "text": "相较固定朝向策略，本文额外预测 ψ_k；相较纯前馈网络，GRU 能缓解单帧深度图带来的局部感知不完整。"},
    ], default_size=17.2, space_after=0.0)


def loss_safety_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.4", "避障与平滑损失")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.72)
    set_textbox(top, [{"level": 0, "text": "损失函数同时约束障碍距离、碰撞速度、加速度、jerk 和角速度，使单步动作序列形成可执行轨迹。"}], default_size=18.5)
    equation_math_rows(slide, [
        {"label": "(4)", "size": 19.0, "pieces": ["L", ("clearance", "sub"), " = (1/T) Σ β", ("1", "sub"), " ln(1 + exp(β", ("2", "sub"), "(d", ("k", "sub"), " − r)))"]},
        {"label": "(5)", "size": 19.0, "pieces": ["L", ("collision", "sub"), " = (1/T) Σ ||v", ("c", "sup"), ("k", "sub"), "|| max(1 − (d", ("k", "sub"), " − r), 0)", ("2", "sup")]},
        {"label": "(6)", "size": 19.0, "pieces": ["L", ("acc", "sub"), " = (1/T) Σ ||a", ("k", "sub"), "||", ("2", "sup")]},
        {"label": "(7)", "size": 19.0, "pieces": ["L", ("jerk", "sub"), " = (1/(T−1)) Σ ||(a", ("k+1", "sub"), " − a", ("k", "sub"), ")/Δt||", ("2", "sup")]},
        {"label": "(8)", "size": 19.0, "pieces": ["L", ("ω", "sub"), " = (1/T) Σ ||ω", ("k", "sub"), "||", ("2", "sup")]},
    ], 0.84, 2.58, 11.70, row_h=0.42)
    add_rule(slide, 0.76, 4.82, 11.80)
    box = add_box(slide, 0.78, 5.04, 11.72, 0.92)
    set_textbox(box, [
        {"level": 1, "text": "距离项提供软安全距离，碰撞项抑制朝障碍方向的速度分量。"},
        {"level": 1, "text": "加速度、jerk 和角速度项让控制输出在时间上连续，降低高速飞行时的姿态突变。"},
        {"level": 1, "text": "这组损失主要处理局部安全与可执行性，仍需要 ToA 和 Yaw 项补足全局绕行与视觉朝向。"},
    ], default_size=17.0, space_after=0.01, line_spacing=0.90)
    add_rule(slide, 0.76, 6.14, 11.80)
    metric_row(slide, [
        ("距离", "安全裕度", "远离障碍表面"),
        ("速度", "碰撞风险", "抑制朝障碍运动"),
        ("平滑", "控制连续", "保证轨迹可执行"),
    ], 1.02, 6.34, 11.28, size=16.0)


def loss_target_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.5", "目标速度与航向对齐")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.78)
    set_textbox(top, [{"level": 0, "text": "目标速度损失负责沿期望方向推进，Yaw 对齐损失负责让机头朝向与运动方向一致。"}], default_size=18.7)
    equation_math_rows(slide, [
        {"label": "(9)", "size": 20.0, "pieces": ["L", ("v", "sub"), " = (1/T) Σ SmoothL1( ||v", ("set", "sup"), ("k", "sub"), " − v̄", ("k", "sub"), "||, 0 )"]},
        {"label": "(10)", "size": 20.0, "pieces": ["L", ("vmax", "sub"), " = (1/T) Σ max( ||v", ("k", "sub"), "|| − v", ("max", "sub"), ", 0)", ("2", "sup")]},
        {"label": "(11)", "size": 20.6, "pieces": ["L", ("yaw", "sub"), " = − (1/T) Σ x", ("B", "sup"), ("k", "sub"), " · ṽ", ("k", "sub")]},
    ], 0.94, 2.84, 11.40, row_h=0.52)
    add_rule(slide, 0.76, 4.66, 11.80)
    left = add_box(slide, 0.78, 4.88, 5.66, 1.04)
    set_textbox(left, [
        {"level": 1, "text": "移动平均速度用于避免单步避障动作被错误惩罚。"},
        {"level": 1, "text": "速度上限项限制过大速度，防止推力输出发散。"},
    ], default_size=16.8, space_after=0.01, line_spacing=0.91)
    right = add_box(slide, 6.86, 4.88, 5.62, 1.04)
    set_textbox(right, [
        {"level": 1, "text": "机体系前向轴与指数加权速度方向越一致，航向损失越小。"},
        {"level": 1, "text": "该项使相机视野更贴合真实运动方向，改善绕大障碍时的视觉闭环。"},
    ], default_size=16.8, space_after=0.01, line_spacing=0.91)
    add_rule(slide, 0.76, 6.12, 11.80)
    metric_row(slide, [
        ("目标速度", "推进方向", "沿 ToA/目标速度前进"),
        ("速度边界", "推力约束", "限制过大推力"),
        ("航向对齐", "机头朝向", "视野贴合运动方向"),
    ], 1.02, 6.32, 11.28, size=16.0)


def toa_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.6", "ToA 特权信息")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.72)
    set_textbox(top, [{"level": 0, "text": "ToA map 只在训练期提供全局路径方向，帮助反应式策略从深度图中学习绕开凹形障碍。"}], default_size=18.5)
    add_pic(slide, "fig4_toa_map_paths.png", 0.86, 2.82, 5.54, 2.96)
    equation_rows(slide, [
        {"label": "", "text": "|∇T(x)| F(x) = 1", "size": 23.0},
        {"label": "", "text": "F(x) = md(x) + (v_slow − mr),  d(x) ≤ d_safe", "size": 18.8},
        {"label": "", "text": "F(x) = 1,  d(x) > d_safe", "size": 18.8},
        {"label": "", "text": "m = (d_safe − v_slow) / (d_safe − r)", "size": 18.8},
    ], 6.74, 2.92, 5.60, row_h=0.46)
    add_rule(slide, 0.76, 6.04, 11.80)
    bottom = add_box(slide, 0.78, 6.22, 11.72, 0.86)
    set_textbox(bottom, [
        {"level": 1, "text": "T(x) 是到达目标的最短时间，负梯度给出 v^set_k 的方向；F(x) 在障碍附近降低传播速度，使路径远离表面。"},
    ], default_size=17.0, space_after=0.0)


def total_loss_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.7", "总目标函数与训练过程")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.76)
    set_textbox(top, [{"level": 0, "text": "训练目标是多个物理可解释损失项的加权组合，通过 BPTT 在时间维度累积梯度。"}], default_size=18.7)
    equation_math_rows(slide, [
        {"label": "(14)", "size": 18.4, "pieces": ["L = λ", ("acc", "sub"), "L", ("acc", "sub"), " + λ", ("jerk", "sub"), "L", ("jerk", "sub"), " + λ", ("ω", "sub"), "L", ("ω", "sub"), " + λ", ("v", "sub"), "L", ("v", "sub"), " + λ", ("vmax", "sub"), "L", ("vmax", "sub")]},
        {"label": "", "size": 18.4, "pieces": ["    + λ", ("clearance", "sub"), "L", ("clearance", "sub"), " + λ", ("collision", "sub"), "L", ("collision", "sub"), " + λ", ("yaw", "sub"), "L", ("yaw", "sub")]},
    ], 0.84, 2.62, 11.80, row_h=0.46)
    add_rule(slide, 0.76, 3.74, 11.80)
    metric_row(slide, [
        ("500", "稳定阶段", "空环境 + 无碰撞损失"),
        ("10K", "训练迭代", "BPTT 展开"),
        ("ToA", "离线预计算", "训练环境复用"),
    ], 0.96, 3.94, 11.42, size=16.3)
    add_rule(slide, 0.76, 5.24, 11.80)
    left = add_box(slide, 0.78, 5.48, 5.82, 1.10)
    set_textbox(left, [
        {"level": 1, "text": "训练环境由随机 primitive 障碍构成，ToA map 预先计算并在训练中反复复用。"},
        {"level": 1, "text": "前 500 次迭代使用空环境和无碰撞损失稳定梯度，随后进入完整环境训练。"},
    ], default_size=16.8, space_after=0.015, line_spacing=0.92)
    right = add_box(slide, 6.86, 5.48, 5.62, 1.10)
    set_textbox(right, [
        {"level": 1, "text": "BPTT 展开多个时间步，时间梯度衰减避免远期梯度主导训练。"},
        {"level": 1, "text": "策略学习到的不是显式地图，而是从深度图推断全局绕行线索的反应式行为。"},
    ], default_size=16.8, space_after=0.015, line_spacing=0.92)


def control_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.8", "姿态控制与 sim-to-real")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.72)
    set_textbox(top, [{"level": 0, "text": "实机部署使用带 body rate 反馈的姿态控制器，把策略输出的推力和航向转成可执行的刚体控制。"}], default_size=18.5)
    add_pic(slide, "fig6_attitude_control.png", 0.84, 2.80, 5.54, 2.80)
    equation_rows(slide, [
        {"label": "(15)", "text": "ω_d = J(ϕ, θ) [ ϕ_dot, θ_dot, ψ_dot ]^T", "size": 21.0},
        {"label": "", "text": "R_d: z_B ∥ F,  x_B aligned with predicted yaw ψ_d", "size": 18.3},
    ], 6.78, 2.92, 5.56, row_h=0.48)
    right = add_box(slide, 6.84, 4.24, 5.58, 1.32)
    set_textbox(right, [
        {"level": 1, "text": "只跟踪姿态会产生约 200 ms 控制滞后，快速避障时会推迟转向。"},
        {"level": 1, "text": "加入 ω_d 后，期望 pitch/roll 与实际响应更同步。"},
    ], default_size=17.0, space_after=0.02, line_spacing=0.94)
    add_rule(slide, 0.76, 5.92, 11.80)
    bottom = add_box(slide, 0.78, 6.14, 11.72, 0.92)
    set_textbox(bottom, [
        {"level": 1, "text": "该控制层让训练中的 point-mass 动作能够落到实机姿态跟踪上，是仿真策略飞到真实四旋翼的关键环节。"},
    ], default_size=17.1, space_after=0.0)


def randomization_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 02", "研究方法")
    add_section_title(slide, "2.9", "域随机化设计")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.72)
    set_textbox(top, [{"level": 0, "text": "域随机化让策略在训练中见到推力、初始状态和传感噪声扰动，从而在实机中补偿模型误差。"}], default_size=18.5)
    rows = [
        ["变量", "分布", "作用"],
        ["g", "N(9.81, 1.5²)", "补偿推力-转速映射偏差"],
        ["z", "U(0.5, 2.5)", "覆盖不同起飞高度"],
        ["v_target", "U(1, 5)", "覆盖不同目标速度"],
        ["v_noise", "N(0, 0.05²)", "模拟状态估计误差"],
        ["r_noise", "N(0, 0.02²)", "模拟姿态观测误差"],
    ]
    small_table(slide, 0.82, 2.80, 6.02, 2.92, rows, font_size=12.9)
    right = add_box(slide, 7.18, 2.86, 5.08, 2.58)
    set_textbox(right, [
        {"level": 1, "text": "策略输出质量归一化推力，真实电机效率和电池电压会改变实际加速度。"},
        {"level": 1, "text": "随机化重力迫使策略从速度反馈中学习闭环补偿，而不是死记 1g 悬停推力。"},
        {"level": 1, "text": "状态噪声提升对 VINS、深度图和姿态估计误差的鲁棒性。"},
    ], default_size=17.1, space_after=0.04, line_spacing=0.94)
    add_rule(slide, 0.76, 6.02, 11.80)
    bottom = add_box(slide, 0.78, 6.22, 11.72, 0.82)
    set_textbox(bottom, [{"level": 1, "text": "硬件消融表明，加入重力随机化后策略会主动输出约 1.3g 初始推力并稳定到目标高度。"}], default_size=17.0, space_after=0)


def experiment_setup_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验与结果")
    add_section_title(slide, "3.1", "仿真与实机实验设置")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.72)
    set_textbox(top, [{"level": 0, "text": "实验覆盖 photorealistic 仿真、消融对比和户外实机飞行，用同一策略检验训练信号的泛化能力。"}], default_size=18.6)
    add_rule(slide, 0.76, 2.78, 11.80)
    metric_row(slide, [
        ("1,350", "仿真试验", "11 个环境"),
        ("3 组", "对比方法", "BNL / yaw w/o ToA / Ours"),
        ("20 次", "户外飞行", "无碰撞"),
        ("4 m/s", "最高速度", "实机测试"),
    ], 0.82, 3.10, 11.72, size=17.2)
    add_rule(slide, 0.76, 4.62, 11.80)
    left = add_box(slide, 0.78, 4.88, 5.76, 1.52)
    set_textbox(left, [
        {"level": 1, "text": "仿真环境包括 Cave、Industry、Mine、Sewer 等 OOD 场景，使用 Flightmare 渲染深度图。"},
        {"level": 1, "text": "评估时不访问 ToA map，策略只用深度、目标和状态输出动作。"},
    ], default_size=17.0, space_after=0.02, line_spacing=0.94)
    right = add_box(slide, 6.88, 4.88, 5.58, 1.52)
    set_textbox(right, [
        {"level": 1, "text": "硬件平台配备 RealSense D456、Lightware 测距、VINS 和 Orin NX。"},
        {"level": 1, "text": "测试覆盖白天林地、夜间 LED 辅助场地和起降高度消融。"},
    ], default_size=17.0, space_after=0.02, line_spacing=0.94)


def success_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验与结果")
    add_section_title(slide, "3.2", "仿真成功率与失败模式")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.62)
    set_textbox(top, [{"level": 0, "text": "完整方法在 11 个复杂环境上达到最高总体成功率，并显著降低碰撞失败比例。"}], default_size=18.5)
    add_pic(slide, "fig7_success_modes.png", 0.70, 2.54, 12.05, 2.92)
    add_rule(slide, 0.76, 5.70, 11.80)
    bottom = add_box(slide, 0.78, 5.92, 11.72, 1.02)
    set_textbox(bottom, [
        {"level": 1, "text": "Ours 总体成功率为 86%，高于 BNL 和去掉 ToA 的 yaw w/o ToA。"},
        {"level": 1, "text": "Mine 场景仍然困难，说明单纯反应式记忆对长回退路径的空间推理能力有限。"},
    ], default_size=17.1, space_after=0.02, line_spacing=0.94)


def trajectory_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验与结果")
    add_section_title(slide, "3.3", "轨迹对比：Yaw 与 ToA 的互补作用")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.62)
    set_textbox(top, [{"level": 0, "text": "轨迹图显示，Yaw 对齐解决朝向问题，ToA 梯度进一步提供绕过凹形障碍的方向线索。"}], default_size=18.5)
    add_pic(slide, "fig8_traj_comparison.png", 0.88, 2.56, 5.78, 3.40)
    right = add_box(slide, 7.08, 2.62, 5.20, 3.10)
    set_textbox(right, [
        {"level": 1, "text": "BNL 固定朝向目标，遇到大墙体时轨迹贴近障碍并失败。"},
        {"level": 1, "text": "Yaw w/o ToA 能转向，但在凹形区域仍可能沿局部方向超时。"},
        {"level": 1, "text": "Yaw w ToA 同时具备重定向和全局路径诱导，更容易绕开墙体并到达目标。"},
        {"level": 1, "text": "该结果说明 ToA 的价值不是部署期建图，而是训练期形成隐式绕行先验。"},
    ], default_size=17.0, space_after=0.05, line_spacing=0.94)
    add_rule(slide, 0.76, 6.20, 11.80)
    metric_row(slide, [
        ("BNL", "固定航向", "大障碍易失败"),
        ("Yaw", "改善转向", "仍缺少全局方向"),
        ("Yaw + ToA", "完整方法", "成功绕障"),
    ], 1.04, 6.26, 11.20, size=16.0)


def gravity_result_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验与结果")
    add_section_title(slide, "3.4", "域随机化硬件消融")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.62)
    set_textbox(top, [{"level": 0, "text": "重力随机化让策略学会从速度反馈中补偿推力误差，实机悬停时高度误差明显减小。"}], default_size=18.5)
    add_pic(slide, "fig9_gravity_randomization.png", 0.76, 2.50, 6.24, 3.38)
    right = add_box(slide, 7.34, 2.70, 4.96, 2.72)
    set_textbox(right, [
        {"level": 1, "text": "未随机化策略初始输出接近 1g，进入自主模式后快速掉高。"},
        {"level": 1, "text": "随机化策略起始推力可升至约 1.3g，抵消实机推力模型偏差。"},
        {"level": 1, "text": "论文指出实际悬停可能需要约 1.15g，而不是理想模型中的 1g。"},
    ], default_size=17.1, space_after=0.04, line_spacing=0.94)
    add_rule(slide, 0.76, 6.16, 11.80)
    bottom = add_box(slide, 0.78, 6.36, 11.72, 0.74)
    set_textbox(bottom, [{"level": 1, "text": "该消融把训练随机化和实机稳定性直接对应起来，是本文 sim-to-real 论证中最清晰的证据。"}], default_size=17.0, space_after=0.0)


def outdoor_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 03", "实验与结果")
    add_section_title(slide, "3.5", "户外与夜间实机飞行")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.60)
    set_textbox(top, [{"level": 0, "text": "策略部署到自研四旋翼后，在白天林地和夜间障碍场中完成高速无碰撞飞行。"}], default_size=18.5)
    add_pic(slide, "fig10_forest_navigation.png", 0.78, 2.50, 5.72, 2.84)
    add_pic(slide, "fig11_night_flight.png", 6.86, 2.50, 5.62, 2.84)
    add_rule(slide, 0.76, 5.66, 11.80)
    bottom = add_box(slide, 0.78, 5.88, 11.72, 1.12)
    set_textbox(bottom, [
        {"level": 1, "text": "林地实验中，策略利用深度输入和速度反馈穿过低矮枝叶，速度最高约 3.8 m/s。"},
        {"level": 1, "text": "夜间实验依赖 LED 辅助与长曝光轨迹记录，说明策略在低光照障碍场中仍能闭环避障。"},
    ], default_size=17.0, space_after=0.02, line_spacing=0.94)


def limitation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 04", "结论与启发")
    add_section_title(slide, "4.1", "结果总结与局限")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.78)
    set_textbox(top, [{"level": 0, "text": "本文证明：训练期的全局特权信息可以转化为部署期的反应式视觉策略，而不必在实机上维护显式地图。"}], default_size=18.8)
    add_rule(slide, 0.76, 2.82, 11.80)
    metric_row(slide, [
        ("86%", "仿真成功率", "Ours 总体表现"),
        ("20 次", "户外飞行", "全程无碰撞"),
        ("589 m", "实机里程", "最高 4 m/s"),
    ], 0.96, 3.12, 11.42, size=17.2)
    add_rule(slide, 0.76, 4.40, 11.80)
    left = add_box(slide, 0.78, 4.66, 5.76, 1.16)
    set_textbox(left, [
        {"level": 1, "text": "优势来自 ToA 训练监督与 Yaw 对齐项的组合：一个提供绕行方向，一个提供相机朝向。"},
        {"level": 1, "text": "实机结果说明域随机化和 body rate 控制能有效弥合 point-mass 训练与刚体飞行。"},
    ], default_size=17.0, space_after=0.01, line_spacing=0.90)
    right = add_box(slide, 6.86, 4.66, 5.62, 1.16)
    set_textbox(right, [
        {"level": 1, "text": "Mine 这类需要回退和长时记忆的迷宫环境仍然困难。"},
        {"level": 1, "text": "初始 yaw 振荡仍存在，说明航向预测和时序记忆还有优化空间。"},
    ], default_size=17.0, space_after=0.01, line_spacing=0.90)
    add_rule(slide, 0.76, 6.06, 11.80)
    bottom = add_box(slide, 0.78, 6.26, 11.72, 0.72)
    set_textbox(bottom, [
        {"level": 1, "text": "对扩展研究而言，更强的 learned memory 和显式不确定性建模，是把该类策略扩展到长航程任务的关键。"},
    ], default_size=16.8, space_after=0.0)


def inspiration_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Part. 04", "结论与启发")
    add_section_title(slide, "4.2", "对无人机研究的启发")
    top = add_box(slide, 0.72, 1.78, 11.90, 0.78)
    set_textbox(top, [{"level": 0, "text": "这篇论文的价值不在于单一网络结构，而在于把训练期全局信息、部署期低延迟和实机可控性统一起来。"}], default_size=18.7)
    add_rule(slide, 0.76, 2.82, 11.80)
    left = add_box(slide, 0.78, 3.10, 5.68, 2.14)
    set_textbox(left, [
        {"level": 1, "text": "训练监督可以比部署感知更强：ToA、ESDF、可达域或拓扑图都可作为特权信号。"},
        {"level": 1, "text": "部署策略仍需要保持轻量，适合机载算力、低延迟控制和状态估计误差。"},
        {"level": 1, "text": "损失函数设计要覆盖安全、目标推进、航向一致和控制平滑四类约束。"},
    ], default_size=17.3, space_after=0.04, line_spacing=0.95)
    right = add_box(slide, 6.86, 3.10, 5.62, 2.14)
    set_textbox(right, [
        {"level": 1, "text": "可引入更强的 learned memory，用于处理需要回退的 Mine 类长时规划。"},
        {"level": 1, "text": "可以把任务语义、局部可通行性和风险预测纳入目标函数，扩展到搜索、巡检和跟随任务。"},
        {"level": 1, "text": "实机闭环仍需把策略输出、姿态控制和传感器延迟作为一个整体设计。"},
    ], default_size=17.3, space_after=0.04, line_spacing=0.95)
    add_rule(slide, 0.76, 5.72, 11.80)
    metric_row(slide, [
        ("全局监督", "训练期", "ToA / 可达性 / 拓扑"),
        ("反应式策略", "部署期", "低延迟机载推理"),
        ("闭环控制", "实机端", "姿态与推力补偿"),
    ], 1.04, 5.98, 11.20, size=16.8)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    cover(prs)
    problem_slide(prs)
    related_slide(prs)
    technical_points_slide(prs)
    overview_slide(prs)
    dynamics_slide(prs)
    network_slide(prs)
    loss_safety_slide(prs)
    loss_target_slide(prs)
    toa_slide(prs)
    total_loss_slide(prs)
    control_slide(prs)
    randomization_slide(prs)
    experiment_setup_slide(prs)
    success_slide(prs)
    trajectory_slide(prs)
    gravity_result_slide(prs)
    outdoor_slide(prs)
    limitation_slide(prs)
    inspiration_slide(prs)
    thanks(prs)
    return prs


def assert_layout(prs: Presentation) -> None:
    def bounds(shape):
        return (
            shape.left / 914400,
            shape.top / 914400,
            (shape.left + shape.width) / 914400,
            (shape.top + shape.height) / 914400,
        )

    def area(a, b):
        return max(0, min(a[2], b[2]) - max(a[0], b[0])) * max(0, min(a[3], b[3]) - max(a[1], b[1]))

    failures: list[str] = []
    last_slide = len(prs.slides)
    for si, slide in enumerate(prs.slides, 1):
        if si in (1, last_slide):
            continue
        items = []
        for idx, shape in enumerate(slide.shapes):
            name = getattr(shape, "name", "")
            left, top, right, bottom = bounds(shape)
            if left < -0.03 or right > SLIDE_W + 0.03 or bottom > SLIDE_H - 0.02:
                failures.append(f"slide {si} shape {idx} outside {left:.2f},{top:.2f},{right:.2f},{bottom:.2f}")
            if top < 1.02 or name.startswith(("MATH_", "DIAG_", "DECOR")):
                continue
            label_type = None
            if shape.shape_type == 13:
                label_type = "pic"
            elif shape.shape_type == 19:
                label_type = "table"
            elif hasattr(shape, "text") and shape.text.strip():
                label_type = "txt"
            if label_type:
                items.append((label_type, idx, bounds(shape)))
        for a in range(len(items)):
            for b in range(a + 1, len(items)):
                if area(items[a][2], items[b][2]) > 0.015:
                    failures.append(f"slide {si} {items[a][0]}#{items[a][1]} overlaps {items[b][0]}#{items[b][1]}")
    if failures:
        raise RuntimeError("\n".join(failures))


def main() -> None:
    OUT.parent.mkdir(exist_ok=True)
    prs = build()
    assert_layout(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
