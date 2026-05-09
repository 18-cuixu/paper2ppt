from __future__ import annotations

import re
import shutil
import zipfile
from pathlib import Path
from xml.sax.saxutils import escape

from PIL import Image, ImageChops, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "out" / "uav-paper-report-v60.pptx"
LOGO = ROOT / "header-logo.png"
CROPS = ROOT / "paper_crops"
BODY_FONT = "Times New Roman"
MATH_FONT = "Times New Roman"
BODY_SIZE = 20.0
AUX_SIZE = 18.0
FORMULA_SIZE = 21.0
TABLE_SIZE = 14.0
MATH_REPLACEMENTS: dict[str, dict] = {}

BLACK = RGBColor(0, 0, 0)
RED = RGBColor(255, 0, 0)
GRAY = RGBColor(78, 78, 78)
HEADER_GRAY = RGBColor(242, 242, 242)
BLUE = RGBColor(0, 121, 192)

EMPH = [
    "multi-UAV CPP", "multi‑UAV CPP",
    "48%", "48.0%", "47.1%", "35.1%",
    "24.6%", "313.1m", "150.8m", "1811.9m", "1349.2m", "51.8%",
    "48.3%", "δd", "NP-hard",
    "NP-hard", "99%", "98%",
]


def add_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def split_terms(text: str, terms: list[str]) -> list[tuple[str, bool]]:
    terms = sorted({t for t in terms if t}, key=len, reverse=True)
    out: list[tuple[str, bool]] = []
    i = 0
    while i < len(text):
        hit = next((term for term in terms if text.startswith(term, i)), None)
        if hit:
            out.append((hit, True))
            i += len(hit)
            continue
        j = i + 1
        while j < len(text) and not any(text.startswith(term, j) for term in terms):
            j += 1
        out.append((text[i:j], False))
        i = j
    return out


def set_run_font(run, name: str = BODY_FONT) -> None:
    run.font.name = name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        element = r_pr.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            r_pr.append(element)
        element.set("typeface", name)


def add_runs(paragraph, text: str, size: float, *, bold=False, terms=None) -> None:
    for piece, marked in split_terms(text, EMPH + (terms or [])):
        run = paragraph.add_run()
        run.text = piece
        set_run_font(run)
        run.font.size = Pt(size)
        run.font.bold = bold or marked
        run.font.color.rgb = RED if marked else BLACK


def compact_text(text: str) -> str:
    join = "\u2060"
    replacements = {
        "6 个 OpenStreetMap": "6个 OpenStreetMap",
        "T1/T2 用于": "T1/T2用于",
        "T3-T6 用于": "T3-T6用于",
        "覆盖率要求": "覆盖阈值",
        "159×319×109m": "159×319×109 m",
        "159×319×109 m": "159×319×109\u00a0m",
        "2 m": "2\u00a0m",
        "50 m": "50\u00a0m",
        "70 m": "70\u00a0m",
        "multi-UAV CPP": "multi-UAV\u00a0CPP",
        "path primitive": "path\u00a0primitive",
        "planned path length": "planned\u00a0path\u00a0length",
        "Set Covering Problem": "Set\u00a0Covering\u00a0Problem",
        "Vehicle Routing Problem": "Vehicle\u00a0Routing\u00a0Problem",
        "random keys": "random\u00a0keys",
        "uncertainty handling": "uncertainty\u00a0handling",
        "set covering": "set\u00a0covering",
        "subtour elimination": "subtour\u00a0elimination",
        "path planning and optimization": "path\u00a0planning\u00a0and\u00a0optimization",
        "OpenStreetMap 结构": "OpenStreetMap结构",
        "UAV 数量": "UAV数量",
        "C-PRM": "C‑PRM",
        "SC-VRP": "SC‑VRP",
        "multi-UAV CPP": "multi‑UAV CPP",
        "2-opt": "2‑opt",
        "min-max": "min‑max",
        "set-covering": "set‑covering",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    text = re.sub(r"(?<=\d)\s+(?=[个次架页类])", "", text)
    text = re.sub(r"(?<=[A-Za-z0-9+\-/])\s+(?=[\u4e00-\u9fff])", "", text)
    text = re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[A-Za-z0-9])", "", text)
    text = text.replace("SC-VRP / VRP / orienteering", "SC-VRP/VRP/orienteering")
    for phrase in (
        "覆盖阈值99%", "覆盖阈值98%", "小型建筑", "高层模型",
        "OSM结构", "OSM建筑", "体素重建", "表面重建",
        "UAV系统", "UAV数量", "20-40 min", "n×1", "min-max",
    ):
        text = text.replace(phrase, join.join(phrase))
    return text


def paragraph_spacing(level: int, base: float, item: dict) -> float:
    if "space_after" in item:
        return item["space_after"]
    if base > 0:
        return base
    return 0


def reject_manual_newline(text: str, context: str, *, allow_newlines: bool = False) -> None:
    if "\n" in text and not allow_newlines:
        raise ValueError(f"{context}: manual newline is not allowed in body text")
    if "\n" in text and any(not part.strip() for part in text.split("\n")):
        raise ValueError(f"{context}: empty line inside manual line break")


def set_textbox(box, lines: list[dict], *, default_size=18, terms=None, space_after=0.06, line_spacing=1.02) -> None:
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for idx, item in enumerate(lines):
        text = item["text"]
        reject_manual_newline(text, f"bullet {idx + 1} in {getattr(box, 'name', 'textbox')}")
        level = item.get("level", 0)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(paragraph_spacing(level, space_after, item))
        ppr = p._p.get_or_add_pPr()
        if level == 0:
            ppr.set("marL", "210000")
            ppr.set("indent", "-140000")
            prefix = item.get("marker", "● ")
            size = float(item.get("size", default_size))
        elif level == 1:
            ppr.set("marL", "350000")
            ppr.set("indent", "-120000")
            prefix = item.get("marker", "• ")
            size = float(item.get("size", default_size))
        else:
            ppr.set("marL", "560000")
            ppr.set("indent", "-105000")
            prefix = item.get("marker", "– ")
            size = float(item.get("size", default_size))
        add_runs(p, prefix + compact_text(text), size, bold=item.get("bold", level == 0), terms=terms)


BODY_ADDITIONS = {
    "1.1": [
        {"level": 1, "text": "任务输入为3D mesh，输出为多UAV巡检路径。"},
        {"level": 2, "text": "方法链：3D mesh → C-PRM → SC-VRP → BRKGA+。"},
        {"level": 2, "text": "评价指标是在覆盖达标后压缩最长单机路径。"},
    ],
    "1.2": [
        {"level": 1, "text": "规划同时受到覆盖完整性、续航时间和多机均衡约束。"},
        {"level": 2, "text": "可写成约束优化：coverage ≥ δd，同时最小化 maxₖ Lₖ。"},
        {"level": 2, "text": "因此本文不是单纯“找一条短路”，而是在多架 UAV 间平衡任务负载。"},
    ],
    "1.3": [
        {"level": 1, "text": "viewpoint-based planning把“点”当覆盖单元；本文把“边/路径基元”也纳入覆盖。"},
        {"level": 2, "text": "路径段飞行过程中相机持续观测，若忽略这部分，会造成重复视点或不必要绕行。"},
        {"level": 2, "text": "这也是 C-PRM 相比普通 PRM 的关键差异。"},
    ],
    "1.4": [
        {"level": 1, "text": "SCP负责覆盖集合，VRP负责多机路径。"},
        {"level": 2, "text": "SC-VRP能直接表达多机覆盖巡检，但搜索空间随路径基元和UAV数量增长。"},
        {"level": 2, "text": "精确 ILP 难以直接处理大规模实例，随后采用 BRKGA 搜索可行路径。"},
    ],
    "1.5": [
        {"level": 1, "text": "方法分为表示层C-PRM、建模层SC-VRP和求解层BRKGA+。"},
        {"level": 2, "text": "实验层用Drake + Octomap验证覆盖质量。"},
        {"level": 2, "text": "三层结构共同支撑大规模多UAV结构巡检问题。"},
    ],
    "2.3": [
        {"level": 1, "text": "变量含义：x̂ᵢⱼₖ 是 UAV k 是否选择 i→j，dᵢⱼ 是路径长度，sᵢⱼ 是覆盖收益。"},
        {"level": 2, "text": "目标函数 min maxₖ ΣᵢΣⱼ dᵢⱼ·x̂ᵢⱼₖ 对应“任务完成时间由最慢 UAV 决定”。"},
        {"level": 2, "text": "覆盖约束 ΣᵢΣₖ sᵢⱼ·x̂ᵢⱼₖ ≥ 1 保证每个目标面片至少被观测。"},
        {"level": 2, "text": "子回路消除项把同一 UAV 的边限制成有效巡检序列，避免数学上可行但飞行上断裂的闭环。"},
    ],
    "2.4": [
        {"level": 1, "text": "解码思想：连续随机键不直接代表 ILP 变量，而是按整数部选择 UAV、按小数部选择邻接边。"},
        {"level": 2, "text": "覆盖率 δ 每扩展一条边就重新计算；达到 δd 后停止，减少无效扩展。"},
        {"level": 2, "text": "适应度 F(P) 同时包含最长单机路径和覆盖不足惩罚，使不可行解在进化中被压制。"},
    ],
    "2.5": [
        {"level": 1, "text": "算法结构：BRKGA 先找可行解，局部改进再删除或替换冗余路径基元。"},
        {"level": 2, "text": "接受准则同时检查 coverage ≥ δd 与路径长度下降，避免“短但漏检”。"},
        {"level": 2, "text": "实现工具链 DEAP + NetworkX 体现出工程启发式特征，适合大规模实例。"},
        {"level": 2, "text": "局部搜索只改变少量邻接关系，因此计算代价可控，适合作为 BRKGA 每代后的增强算子。"},
    ],
    "2.6": [
        {"level": 1, "text": "实验公平性：BRKGA/BRKGA+ 都取 10 次平均，降低随机初值影响。"},
        {"level": 2, "text": "T1/T2 用更严格 99% 覆盖率，T3-T6 用 98% 覆盖率测试更大结构。"},
        {"level": 2, "text": "FOV、最大视角、安全距离和覆盖阈值固定后，路径长度比较才有意义。"},
    ],
    "4.1": [
        {"level": 1, "text": "覆盖达标与路径缩短必须同时成立。"},
        {"level": 2, "text": "可复用结构：C-PRM生成候选路径，SC-VRP/BRKGA+完成选择与排序。"},
    ],
    "5.1": [
        {"level": 1, "text": "系统闭环：任务理解 → 约束生成 → 路径规划 → 仿真验证 → 缺口反馈。"},
        {"level": 2, "text": "任务语义需要转成覆盖、禁飞区和多机分配约束。"},
    ],
}


def plain(
    box,
    text: str,
    size: int,
    *,
    bold=False,
    color=BLACK,
    align=PP_ALIGN.LEFT,
    allow_newlines: bool = False,
) -> None:
    reject_manual_newline(text, f"plain({getattr(box, 'name', 'textbox')})", allow_newlines=allow_newlines)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        set_run_font(run)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_header(slide, part: str, title: str) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.87))
    band.fill.solid()
    band.fill.fore_color.rgb = HEADER_GRAY
    band.line.fill.background()
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.49), Inches(0.06), width=Inches(2.65), height=Inches(0.74))
    box = add_box(slide, 3.25, 0.11, 7.7, 0.62)
    box.text_frame.word_wrap = False
    plain(box, f"{part}  {title}", 28, bold=True)


def add_section_title(slide, num: str, title: str) -> None:
    box = add_box(slide, 0.70, 1.12, 8.5, 0.54)
    box.text_frame.word_wrap = False
    plain(box, f"{num} {title}", 20, bold=True)


def add_rule(slide, left: float, top: float, width: float, color=RGBColor(190, 190, 190)) -> None:
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.018))
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()


def caption(slide, text: str, left, top, width, height) -> None:
    box = add_box(slide, left, top, width, height)
    plain(box, text, 18, color=GRAY)


def muted_label(slide, text: str, left, top, width, height) -> None:
    box = add_box(slide, left, top, width, height)
    plain(box, text, 20, bold=True, color=GRAY)


def add_pic(slide, path: Path, left, top, width, height):
    with Image.open(path) as im:
        ratio = im.width / im.height
    box_ratio = width / height
    if ratio > box_ratio:
        w = width
        h = width / ratio
        x = left
        y = top + (height - h) / 2
    else:
        h = height
        w = height * ratio
        x = left + (width - w) / 2
        y = top
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_pic_fill(slide, path: Path, left, top, width, height):
    shape = slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    shape.crop_left = 0
    shape.crop_right = 0
    shape.crop_top = 0
    shape.crop_bottom = 0
    return shape


def set_shape_name(shape, name: str) -> None:
    shape._element.xpath(".//p:cNvPr")[0].set("name", name)


def crop_whitespace(src: Path, dst: Path, pad=8) -> Path:
    img = Image.open(src).convert("RGB")
    bg = Image.new("RGB", img.size, "white")
    diff = ImageChops.difference(img, bg)
    bbox = diff.getbbox()
    if not bbox:
        img.save(dst)
        return dst
    left = max(0, bbox[0] - pad)
    top = max(0, bbox[1] - pad)
    right = min(img.width, bbox[2] + pad)
    bottom = min(img.height, bbox[3] + pad)
    img.crop((left, top, right, bottom)).save(dst)
    return dst


def crop_box(src: Path, dst: Path, rel_box: tuple[float, float, float, float]) -> Path:
    img = Image.open(src).convert("RGB")
    l, t, r, b = rel_box
    img.crop((int(img.width * l), int(img.height * t), int(img.width * r), int(img.height * b))).save(dst)
    return dst


def make_structure_grid(src: Path, dst: Path) -> Path:
    img = Image.open(src).convert("RGB")
    boxes = [
        (0.00, 0.02, 0.31, 0.36),
        (0.33, 0.02, 0.62, 0.36),
        (0.65, 0.02, 0.99, 0.36),
        (0.00, 0.50, 0.32, 0.80),
        (0.34, 0.48, 0.65, 0.80),
        (0.66, 0.48, 0.99, 0.80),
    ]
    crops = []
    for l, t, r, b in boxes:
        item = img.crop((int(img.width * l), int(img.height * t), int(img.width * r), int(img.height * b)))
        bg = Image.new("RGB", item.size, "white")
        bbox = ImageChops.difference(item, bg).getbbox()
        if bbox:
            pad = 10
            item = item.crop((
                max(0, bbox[0] - pad),
                max(0, bbox[1] - pad),
                min(item.width, bbox[2] + pad),
                min(item.height, bbox[3] + pad),
            ))
        crops.append(item)

    canvas = Image.new("RGB", (1900, 1040), "white")
    cell_w, cell_h = 600, 470
    margin_x, margin_y = 40, 45
    gap_x, gap_y = 10, 35
    for idx, item in enumerate(crops):
        row, col = divmod(idx, 3)
        scale = min((cell_w - 36) / item.width, (cell_h - 30) / item.height)
        resized = item.resize((int(item.width * scale), int(item.height * scale)), Image.Resampling.LANCZOS)
        x = margin_x + col * (cell_w + gap_x) + (cell_w - resized.width) // 2
        y = margin_y + row * (cell_h + gap_y) + (cell_h - resized.height) // 2
        canvas.paste(resized, (x, y))
    canvas.save(dst)
    return dst


def make_coverage_grid(src: Path, dst: Path) -> Path:
    img = Image.open(src).convert("RGB")
    boxes = [
        (0.05, 0.03, 0.34, 0.44),
        (0.36, 0.03, 0.64, 0.44),
        (0.66, 0.03, 0.96, 0.44),
        (0.05, 0.57, 0.34, 0.98),
        (0.36, 0.48, 0.64, 0.98),
        (0.66, 0.48, 0.96, 0.98),
    ]
    crops = []
    for l, t, r, b in boxes:
        item = img.crop((int(img.width * l), int(img.height * t), int(img.width * r), int(img.height * b)))
        bg = Image.new("RGB", item.size, "white")
        bbox = ImageChops.difference(item, bg).getbbox()
        if bbox:
            pad = 14
            item = item.crop((
                max(0, bbox[0] - pad),
                max(0, bbox[1] - pad),
                min(item.width, bbox[2] + pad),
                min(item.height, bbox[3] + pad),
            ))
        crops.append(item)

    canvas = Image.new("RGB", (2050, 1120), "white")
    cell_w, cell_h = 650, 520
    margin_x, margin_y = 38, 35
    gap_x, gap_y = 10, 35
    for idx, item in enumerate(crops):
        row, col = divmod(idx, 3)
        scale = min((cell_w - 28) / item.width, (cell_h - 20) / item.height)
        resized = item.resize((int(item.width * scale), int(item.height * scale)), Image.Resampling.LANCZOS)
        x = margin_x + col * (cell_w + gap_x) + (cell_w - resized.width) // 2
        y = margin_y + row * (cell_h + gap_y) + (cell_h - resized.height) // 2
        canvas.paste(resized, (x, y))
    canvas.save(dst)
    return dst


def body_slide(prs, part, part_title, sec, sec_title, bullets, *, default_size=20):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    merged = bullets + BODY_ADDITIONS.get(sec, [])
    box = add_box(slide, 0.46, 1.66, 12.70, 5.76)
    if sec.startswith("1."):
        set_textbox(box, merged, default_size=default_size, space_after=0, line_spacing=1.00)
    elif sec in ("4.1", "5.1"):
        set_textbox(box, merged, default_size=default_size, space_after=0, line_spacing=0.99)
    else:
        set_textbox(box, merged, default_size=default_size, space_after=0, line_spacing=1.00)
    return slide


def formula(slide, left, top, width, height, text: str, size=25, bold=False, color=BLACK) -> None:
    box = add_box(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run_font(run, MATH_FONT)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_formula_panel(slide, lines: list[str], left: float, top: float, width: float, height: float, *, size=22.0, align=PP_ALIGN.CENTER, spacing=1.08) -> None:
    box = add_box(slide, left, top, width, height)
    set_shape_name(box, "MATH_TEXT_VISIBLE")
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_before = Pt(0)
        p.space_after = Pt(2 if idx < len(lines) - 1 else 0)
        run = p.add_run()
        run.text = line
        set_run_font(run, MATH_FONT)
        run.font.size = Pt(size)
        run.font.color.rgb = BLACK


def add_formula_image(slide, key: str, fallback_img: Path, left: float, top: float, width: float, height: float) -> None:
    pic = add_pic(slide, fallback_img, left, top, width, height)
    set_shape_name(pic, f"MATH_VISIBLE_{key}")


def math_placeholder(slide, token: str, left: float, top: float, width: float, height: float, *, size=20) -> None:
    box = add_box(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = token
    set_run_font(run, MATH_FONT)
    run.font.size = Pt(size)
    run.font.color.rgb = BLACK


def _math_rpr(*, sz="2000", italic=True, bold=False) -> str:
    attrs = f'lang="en-US" altLang="zh-CN" sz="{sz}"'
    if bold:
        attrs += ' b="1"'
    attrs += f' i="{1 if italic else 0}" smtClean="0"'
    return (
        f'<a:rPr {attrs}>'
        '<a:latin typeface="Cambria Math" panose="02040503050406030204" pitchFamily="18" charset="0"/>'
        '<a:ea typeface="Times New Roman" pitchFamily="18" charset="0"/>'
        '</a:rPr>'
    )


def _math_run(text: str, *, sz="2000", italic=True, bold=False) -> str:
    safe = escape(text)
    return f'<m:r>{_math_rpr(sz=sz, italic=italic, bold=bold)}<m:t>{safe}</m:t></m:r>'


def _math_ctrl(sz="2000", italic=True) -> str:
    return f'<m:ctrlPr>{_math_rpr(sz=sz, italic=italic)}</m:ctrlPr>'


def _m_sub(base: str, sub: str, *, sz="2000") -> str:
    return (
        f'<m:sSub><m:sSubPr>{_math_ctrl(sz)}</m:sSubPr>'
        f'<m:e>{base}</m:e><m:sub>{sub}</m:sub></m:sSub>'
    )


def _m_sup(base: str, sup: str, *, sz="2000") -> str:
    return (
        f'<m:sSup><m:sSupPr>{_math_ctrl(sz)}</m:sSupPr>'
        f'<m:e>{base}</m:e><m:sup>{sup}</m:sup></m:sSup>'
    )


def _m_frac(num: str, den: str, *, sz="2000") -> str:
    return (
        f'<m:f><m:fPr><m:type m:val="bar"/>{_math_ctrl(sz)}</m:fPr>'
        f'<m:num>{num}</m:num><m:den>{den}</m:den></m:f>'
    )


def _m_sum(sub: str, sup: str, expr: str, *, sz="2000") -> str:
    return (
        f'<m:nary><m:naryPr><m:chr m:val="∑"/><m:limLoc m:val="undOvr"/>'
        f'<m:grow m:val="1"/>{_math_ctrl(sz)}</m:naryPr>'
        f'<m:sub>{sub}</m:sub><m:sup>{sup}</m:sup><m:e>{expr}</m:e></m:nary>'
    )


def _m_delim(content: str, beg: str, end: str, *, sz="2000") -> str:
    return (
        f'<m:d><m:dPr><m:begChr m:val="{escape(beg)}"/><m:endChr m:val="{escape(end)}"/>'
        f'<m:grow m:val="1"/>{_math_ctrl(sz)}</m:dPr><m:e>{content}</m:e></m:d>'
    )


def math_line(text: str, *, sz="2000", spacing="115000", raw=False) -> str:
    body = text if raw else _math_run(text, sz=sz)
    return (
        f'<a:p><a:pPr algn="ctr"><a:lnSpc><a:spcPct val="{spacing}"/></a:lnSpc></a:pPr>'
        '<a14:m xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main">'
        '<m:oMathPara xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">'
        '<m:oMathParaPr><m:jc m:val="centerGroup"/></m:oMathParaPr>'
        f'<m:oMath>{body}</m:oMath></m:oMathPara></a14:m></a:p>'
    )


def native_math_body(key: str, sz="2200") -> str | None:
    r = lambda text, italic=True: _math_run(text, sz=sz, italic=italic)
    sub = lambda base, s: _m_sub(base, s, sz=sz)
    sup = lambda base, s: _m_sup(base, s, sz=sz)
    summ = lambda lo, hi, body: _m_sum(r(lo), r(hi), body, sz=sz)
    xhat = sub(r("x̂"), r("i,j,k"))
    if key == "MODEL":
        dij_x = sub(r("d"), r("i,j")) + r(" ") + xhat
        obj = sub(r("min", False), r("x̂")) + r(" ") + sub(r("max", False), r("k∈K")) + r(" ") + summ("i=1", "n", summ("j=1", "n", dij_x))
        cov = r("s.t.  ", False) + summ("i=1", "n", summ("k=1", "K", sub(r("s"), r("i,j")) + r(" ") + xhat)) + r(" ≥ 1,  ∀j")
        inner = summ("i=1", "n", summ("k=1", "K", sub(r("s"), r("i,j")) + r(" ") + xhat)) + r(" ≥ 1")
        ratio = _m_frac(summ("j=1", "m", r("1[") + inner + r("]")), r("m"), sz=sz) + r(" ≥ ") + sub(r("δ"), r("d"))
        binary = xhat + r(" ∈ ") + _m_delim(r("0,1"), "{", "}", sz=sz) + r(",  ") + sub(r("e"), r("i,j")) + r(" ∈ E,  k∈K")
        return "".join(math_line(part, sz=sz, spacing="105000", raw=True) for part in [obj, cov, ratio, binary])
    if key == "DECODE":
        c_i = sub(r("c"), r("i"))
        c_vec = r("C") + r(" = ") + sup(_m_delim(sub(r("c"), r("1")) + r(", …, ") + sub(r("c"), r("n")), "(", ")", sz=sz), r("T")) + r(",  0 ≤ ") + c_i + r(" ≤ K")
        decode = sub(r("u"), r("i")) + r(" = ⌊") + c_i + r("⌋,  ") + sub(r("q"), r("i")) + r(" = ") + c_i + r(" − ") + sub(r("u"), r("i"))
        edge_sum = _m_sum(sub(r("e"), r("a,b")) + r("∈") + sub(r("E"), r("res")) + _m_delim(r("C"), "(", ")", sz=sz), r(""), sub(r("s"), r("a,b,j")), sz=sz)
        delta = sub(r("δ"), _m_delim(r("C"), "(", ")", sz=sz)) + r(" = ") + _m_frac(summ("j=1", "m", r("1[") + edge_sum + r(" ≥ 1]")), r("m"), sz=sz)
        fit = r("F") + _m_delim(r("C"), "(", ")", sz=sz) + r(" = ") + sub(r("max", False), r("k∈K")) + r(" ") + sub(r("L"), r("k")) + _m_delim(r("C"), "(", ")", sz=sz) + r(" + λ max(0, ") + sub(r("δ"), r("d")) + r(" − ") + sub(r("δ"), _m_delim(r("C"), "(", ")", sz=sz)) + r(")")
        return "".join(math_line(part, sz=sz, spacing="105000", raw=True) for part in [c_vec, decode, delta, fit])
    if key == "LOCAL":
        p_k = sub(r("P"), r("k"))
        p_prime = sub(r("P′"), r("k"))
        first = p_prime + r(" = 2optSwap", False) + _m_delim(p_k + r(", ") + sub(r("v"), r("i")) + r(", ") + sub(r("v"), r("j")), "(", ")", sz=sz) + r(",  ") + sub(r("v"), r("j")) + r("∈") + sub(r("N"), r("G")) + _m_delim(sub(r("v"), r("i")), "(", ")", sz=sz)
        accept = r("accept", False) + _m_delim(r("P′"), "(", ")", sz=sz) + r(" ⇔ ") + sub(r("δ"), _m_delim(r("P′"), "(", ")", sz=sz)) + r(" ≥ ") + sub(r("δ"), r("d")) + r("  and  ") + sub(r("max", False), r("k∈K")) + r(" ") + sub(r("L′"), r("k")) + r(" < ") + sub(r("max", False), r("k∈K")) + r(" ") + sub(r("L"), r("k"))
        update = sup(r("P"), r("t+1")) + r(" = ") + _m_delim(r("P′, if accept(P′);  ") + sup(r("P"), r("t")) + r(", otherwise"), "{", "", sz=sz)
        return "".join(math_line(part, sz=sz, spacing="105000", raw=True) for part in [first, accept, update])
    return None


def editable_formula_text(slide, lines: list[str], left: float, top: float, width: float, height: float, *, size=20, align=PP_ALIGN.CENTER):
    box = add_box(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.22
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        set_run_font(run, MATH_FONT)
        run.font.size = Pt(size)
        run.font.color.rgb = BLACK


def native_formula(slide, key: str, fallback_img: Path, left: float, top: float, width: float, height: float, *, size=22) -> None:
    return
    token = f"__NATIVE_MATH_{key}__"
    native = add_box(slide, left, top, width, height)
    set_shape_name(native, f"MATH_NATIVE_{key}")
    tf = native.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = token
    set_run_font(run, MATH_FONT)
    run.font.size = Pt(1)
    run.font.color.rgb = BLACK

    MATH_REPLACEMENTS[token] = {
        "key": key,
        "fallback_name": "",
        "lines": NATIVE_MATH[key],
        "size": size,
    }


def inject_native_math(pptx_path: Path, replacements: dict[str, dict]) -> None:
    tmp = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if info.filename.startswith("ppt/slides/slide") and info.filename.endswith(".xml"):
                text = data.decode("utf-8")
                for token, spec in replacements.items():
                    if token not in text:
                        continue
                    sz = str(int(round(spec.get("size", FORMULA_SIZE) * 100)))
                    math_xml = native_math_body(spec["key"], sz=sz)
                    if math_xml is None:
                        math_xml = "".join(math_line(line, sz=sz) for line in spec["lines"])
                    tx_body = (
                        '<p:txBody><a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" rtlCol="0">'
                        '<a:noAutofit/></a:bodyPr><a:lstStyle/>'
                        f'{math_xml}</p:txBody>'
                    )
                    native_pat = re.compile(
                        r'<p:sp>(?:(?!</p:sp>).)*?<a:t>' + re.escape(token) + r'</a:t>(?:(?!</p:sp>).)*?</p:sp>',
                        re.S,
                    )
                    native_match = native_pat.search(text)
                    if not native_match:
                        continue
                    native_shape = re.sub(r'<p:txBody>[\s\S]*?</p:txBody>', tx_body, native_match.group(0), count=1)
                    text = text[:native_match.start()] + native_shape + text[native_match.end():]
                data = text.encode("utf-8")
            zout.writestr(info, data)
    shutil.move(tmp, pptx_path)


def render_formula_rows(rows: list[tuple[str, float]], dst: Path, *, width=3300, row_h=260, gap=10, fontsize=31) -> Path:
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    plt.rcParams.update({
        "mathtext.fontset": "stix",
        "font.family": "STIXGeneral",
        "mathtext.rm": "STIXGeneral",
        "mathtext.it": "STIXGeneral:italic",
        "mathtext.bf": "STIXGeneral:bold",
    })
    rendered = []
    transparent = (255, 255, 255, 0)
    for text, x in rows:
        fig = plt.figure(figsize=(12.8, 1.16), dpi=260)
        fig.patch.set_alpha(0)
        ax = fig.add_axes((0, 0, 1, 1))
        ax.axis("off")
        fig.text(x, 0.50, text, ha="left", va="center", fontsize=fontsize, color="black")
        fig.canvas.draw()
        w, h = fig.canvas.get_width_height()
        img = Image.frombytes("RGBA", (w, h), fig.canvas.buffer_rgba())
        bbox = ImageChops.difference(img, Image.new("RGBA", img.size, transparent)).getbbox()
        if bbox:
            img = img.crop((
                max(0, bbox[0] - 8),
                max(0, bbox[1] - 8),
                min(w, bbox[2] + 8),
                min(h, bbox[3] + 8),
            ))
        rendered.append(img)
        plt.close(fig)

    canvas_h = len(rendered) * row_h + (len(rendered) - 1) * gap
    canvas = Image.new("RGBA", (width, canvas_h), transparent)
    y = 0
    for img in rendered:
        scale = min(1.0, (width - 70) / img.width, (row_h - 12) / img.height)
        resized = img.resize((int(img.width * scale), int(img.height * scale)), Image.Resampling.LANCZOS)
        canvas.alpha_composite(resized, (34, y + (row_h - resized.height) // 2))
        y += row_h + gap
    bbox = ImageChops.difference(canvas, Image.new("RGBA", canvas.size, transparent)).getbbox()
    if bbox:
        pad = 18
        canvas = canvas.crop((
            max(0, bbox[0] - pad),
            max(0, bbox[1] - pad),
            min(canvas.width, bbox[2] + pad),
            min(canvas.height, bbox[3] + pad),
        ))
    canvas.save(dst)
    return dst


def make_model_formula(dst: Path) -> Path:
    return render_formula_rows([
        (r"$\min_{\hat{x}}\ \max_{k\in K}\ \sum_{i=1}^{n}\sum_{j=1}^{n} d_{ij}\hat{x}_{ijk}$", 0.06),
        (r"$\mathrm{s.t.}\quad \sum_{i=1}^{n}\sum_{k=1}^{K}s_{ij}\hat{x}_{ijk}\geq1,\quad \forall j$", 0.06),
        (r"$\frac{1}{m}\sum_{j=1}^{m} I_j\geq\delta_d,\quad I_j=1\ \mathrm{if}\ \sum_{i=1}^{n}\sum_{k=1}^{K}s_{ij}\hat{x}_{ijk}\geq1$", 0.06),
        (r"$\hat{x}_{i,j,k}\in\{0,1\},\quad e_{ij}\in E,\ k\in K$", 0.06),
    ], dst, width=3400, row_h=250, gap=8, fontsize=31)


def make_decode_formula(dst: Path) -> Path:
    return render_formula_rows([
        (r"$C=(c_1,\ldots,c_n)^{T},\quad 0\leq c_i\leq K$", 0.06),
        (r"$u_i=\lfloor c_i\rfloor,\quad q_i=c_i-u_i$", 0.06),
        (r"$\delta(C)=\frac{1}{m}\sum_{j=1}^{m} I_j(C),\quad I_j(C)=1\ \mathrm{if}\ \sum_{e_{ab}\in E_{\mathrm{res}}(C)}s_{abj}\geq1$", 0.06),
        (r"$F(C)=\max_{k\in K}L_k(C)+\lambda\max\!\left(0,\delta_d-\delta(C)\right)$", 0.06),
    ], dst, width=3400, row_h=230, gap=8, fontsize=30)


def make_local_formula(dst: Path) -> Path:
    return render_formula_rows([
        (r"$P'_k=\mathrm{2optSwap}(P_k,v_i,v_j),\quad v_j\in N_G(v_i)$", 0.06),
        (r"$\mathrm{accept}(P')\Longleftrightarrow \delta(P')\geq\delta_d\ \wedge\ \max_{k\in K}L'_k<\max_{k\in K}L_k$", 0.06),
        (r"$P^{t+1}=P',\quad \mathrm{if}\ \mathrm{accept}(P')$", 0.06),
        (r"$P^{t+1}=P^t,\quad \mathrm{otherwise}$", 0.06),
    ], dst, width=3400, row_h=215, gap=8, fontsize=28)


def make_model_objective_formula(dst: Path) -> Path:
    return render_formula_rows([
        (r"$\min_{\hat{x}}\ \max_{k\in K}\ \sum_{i=1}^{n}\sum_{j=1}^{m} d_{ij}\hat{x}_{ijk}$", 0.06),
        (r"$\hat{x}_{ijk}\in\{0,1\},\qquad e_{ij}\in E,\quad k\in K$", 0.06),
    ], dst, width=3450, row_h=250, gap=8, fontsize=37)


def make_model_constraints_formula(dst: Path) -> Path:
    return render_formula_rows([
        (r"$\sum_{i=1}^{n}\sum_{k=1}^{K}s_{ij}\hat{x}_{ijk}\geq1,\qquad \forall j$", 0.06),
        (r"$\sum_{i,j:\ v_i,v_j\in V'_k}\hat{x}_{ijk}\leq |L|-1,\qquad |L|\subset V'_k,\ |L|\geq2$", 0.06),
        (r"$V'_k\subset V,\ \mathrm{for}\ \hat{x}_{ijk}=1;\qquad e_{ij}\in E,\ \mathrm{for}\ \hat{x}_{ijk}=1$", 0.06),
        (r"$\hat{x}_{ijk}\in\{0,1\},\qquad (i,j)\in E,\quad k\in K$", 0.06),
    ], dst, width=3600, row_h=195, gap=0, fontsize=29)


def make_model_coverage_ratio_formula(dst: Path) -> Path:
    return render_formula_rows([
        (r"$\frac{\sum_{j=1}^{m}\left(\sum_{i=1}^{n}\sum_{k=1}^{K}s_{ij}\hat{x}_{ijk}\geq1\right)}{m}\geq\delta_d$", 0.06),
    ], dst, width=3450, row_h=285, gap=8, fontsize=36)


def make_decode_encoding_formula(dst: Path) -> Path:
    return render_formula_rows([
        (r"$C=(c_1,\ldots,c_n)^T,\qquad x_i\in[0,K]$", 0.08),
        (r"$x_{\mathrm{int}}=\lfloor x_i\rfloor,\qquad x_{\mathrm{frac}}=x_i-x_{\mathrm{int}}$", 0.06),
    ], dst, width=3450, row_h=250, gap=8, fontsize=37)


def make_decode_fitness_formula(dst: Path) -> Path:
    return render_formula_rows([
        (r"$E_{\mathrm{res}}\leftarrow E_{\mathrm{res}}\cup e_i,\qquad s_{x_{\mathrm{int}}}\leftarrow s_{x_{\mathrm{int}}}+d(e_i)$", 0.06),
        (r"$\delta\leftarrow \mathrm{evaluateCoverage}(E_{\mathrm{res}})$", 0.06),
        (r"$\delta\geq\delta_d\quad\Longrightarrow\quad s_{\max}\leftarrow\max_{k\in K}s_k$", 0.06),
        (r"$f(C)=s_{\max}$", 0.06),
    ], dst, width=3500, row_h=185, gap=4, fontsize=32)


def make_local_accept_formula(dst: Path) -> Path:
    return render_formula_rows([
        (r"$C_o\leftarrow C,\qquad p_k\in P,\qquad v_i\in p_k,\qquad v_j\in \mathrm{neighbour}(v_i,\mathcal{G})$", 0.06),
        (r"$C_o\leftarrow \mathrm{2optSwap}(C_o,v_i,v_j)$", 0.06),
        (r"$\Delta L<0,\qquad \delta(C_o)\geq\delta_d$", 0.06),
    ], dst, width=3600, row_h=210, gap=4, fontsize=33)


def sc_vrp_objective_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    top = add_box(slide, 0.62, 1.70, 12.02, 0.74)
    set_textbox(top, [
        {"level": 0, "text": "SC-VRP 将多 UAV 巡检写成 min-max 组合优化，目标是压缩最长单机路径。"},
        {"level": 1, "text": "决策变量 x̂ᵢⱼₖ 表示第 k 架 UAV 是否选择 C-PRM 边 eᵢⱼ。"},
    ], default_size=20, space_after=0, line_spacing=0.96)

    muted_label(slide, "目标函数  Eq. (2)", 0.74, 2.55, 3.20, 0.34)
    add_formula_panel(slide, [
        "min  maxₖ∈K   Σᵢ₌₁ⁿ Σⱼ₌₁ⁿ dᵢⱼ x̂ᵢⱼₖ",
        "x̂ᵢⱼₖ ∈ {0,1},      eᵢⱼ∈E,      k∈K",
    ], 0.82, 2.88, 11.90, 1.00, size=22.5, spacing=1.04)
    if formula_img is not None:
        native_formula(slide, "MODEL", formula_img, 12.93, 7.20, 0.08, 0.05, size=14)

    add_rule(slide, 0.74, 4.22, 11.85)
    left = add_box(slide, 0.70, 4.44, 6.08, 1.52)
    set_textbox(left, [
        {"level": 1, "text": "dᵢⱼ 是路径基元长度，求和得到第 k 架 UAV 路径 Lₖ。"},
        {"level": 1, "text": "maxₖLₖ 对应任务瓶颈，不是所有 UAV 路径之和。"},
        {"level": 1, "text": "min-max 目标抑制单架 UAV 路径过长。"},
    ], default_size=19.0, space_after=0, line_spacing=0.91)

    right = add_box(slide, 6.92, 4.44, 5.70, 1.52)
    set_textbox(right, [
        {"level": 1, "text": "只最小化总路程，可能让单机承担过多巡检任务。"},
        {"level": 1, "text": "负载均衡直接进入目标函数，更符合多机并行巡检。"},
        {"level": 1, "text": "覆盖约束再保证每个目标面片被观测。"},
    ], default_size=19.0, space_after=0, line_spacing=0.91)

    bottom = add_box(slide, 0.70, 6.16, 11.88, 0.76)
    set_textbox(bottom, [
        {"level": 1, "text": "该目标把“路径短”和“多机均衡”合并，BRKGA 适应度也围绕它设计。"},
    ], default_size=19.0, space_after=0, line_spacing=0.94)
    return slide


def sc_vrp_constraints_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None, ratio_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    top = add_box(slide, 0.62, 1.70, 12.02, 0.74)
    set_textbox(top, [
        {"level": 0, "text": "约束项同时保证目标面片覆盖、路径连通和二进制选择。"},
        {"level": 1, "text": "覆盖约束由 sᵢⱼ 表达路径基元对目标面片的观测收益。"},
    ], default_size=20, space_after=0, line_spacing=0.96)

    muted_label(slide, "覆盖与连通约束  Eq. (3)-(7)", 0.74, 2.52, 4.90, 0.34)
    add_formula_panel(slide, [
        "Σᵢ₌₁ⁿ Σₖ₌₁ᴷ sᵢⱼ x̂ᵢⱼₖ ≥ 1,      ∀j",
        "Σᵢ,ⱼ:ᵥᵢ,ᵥⱼ∈V'ₖ x̂ᵢⱼₖ ≤ |L|-1,      |L|⊂V'ₖ, |L|≥2",
        "V'ₖ⊂V for x̂ᵢⱼₖ=1;      eᵢⱼ∈E for x̂ᵢⱼₖ=1",
    ], 0.74, 2.92, 6.78, 1.34, size=17.4, spacing=1.00)
    muted_label(slide, "覆盖率阈值  Eq. (8)", 7.70, 2.52, 3.60, 0.34)
    add_formula_panel(slide, [
        "1/m · Σⱼ₌₁ᵐ I(Σᵢ₌₁ⁿΣₖ₌₁ᴷ sᵢⱼx̂ᵢⱼₖ ≥ 1) ≥ δd",
    ], 7.38, 3.14, 5.22, 0.68, size=17.4, spacing=1.00)
    if formula_img is not None:
        native_formula(slide, "MODEL", formula_img, 12.93, 7.20, 0.08, 0.05, size=14)

    add_rule(slide, 0.74, 4.54, 11.85)
    left = add_box(slide, 0.70, 4.78, 5.94, 1.76)
    set_textbox(left, [
        {"level": 1, "text": "Eq. (3) 要求每个目标面片至少被一条已选路径基元覆盖。"},
        {"level": 1, "text": "Eq. (4) 是子回路消除约束，避免数学上连通但飞行上断裂的闭环。"},
        {"level": 1, "text": "Eq. (5)-(7) 约束节点访问、边选择和二进制变量。"},
    ], default_size=19.2, space_after=0, line_spacing=0.91)

    right = add_box(slide, 6.82, 4.78, 5.80, 1.76)
    set_textbox(right, [
        {"level": 1, "text": "Eq. (8) 用 δd 控制覆盖阈值，实验中小规模结构取99%，大规模结构取98%。"},
        {"level": 1, "text": "阈值越高，可行路径越难找，但覆盖质量更接近全覆盖。"},
        {"level": 1, "text": "BRKGA 解码时也用该阈值判断何时停止扩展路径。"},
    ], default_size=19.2, space_after=0, line_spacing=0.91)
    return slide


def brkga_encoding_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    top = add_box(slide, 0.62, 1.70, 12.02, 0.86)
    set_textbox(top, [
        {"level": 0, "text": "BRKGA 用连续随机键表示解，再通过解码过程映射到多 UAV 路径。"},
        {"level": 1, "text": "随机键把受约束的组合优化空间转成更容易搜索的编码空间。"},
    ], default_size=20, space_after=0, line_spacing=0.95)

    muted_label(slide, "随机键编码", 0.74, 2.70, 3.20, 0.34)
    add_formula_panel(slide, [
        "C = (c₁, …, cₙ)ᵀ,        cᵢ ∈ [0, K]",
        "uᵢ = floor(cᵢ),        qᵢ = cᵢ - uᵢ",
    ], 0.86, 3.04, 11.76, 0.98, size=22.0, spacing=1.05)
    if formula_img is not None:
        native_formula(slide, "DECODE", formula_img, 12.93, 7.20, 0.08, 0.05, size=14)

    add_rule(slide, 0.74, 4.38, 11.85)
    left = add_box(slide, 0.70, 4.62, 5.90, 1.78)
    set_textbox(left, [
        {"level": 1, "text": "整数部 uᵢ 决定该路径基元分配给哪架 UAV。"},
        {"level": 1, "text": "小数部 qᵢ 在当前邻域内选择下一条可扩展边。"},
        {"level": 1, "text": "这种编码避免直接操作大量 0/1 决策变量。"},
    ], default_size=19.5, space_after=0, line_spacing=0.92)

    right = add_box(slide, 6.78, 4.62, 5.84, 1.78)
    set_textbox(right, [
        {"level": 1, "text": "不是所有节点都必须访问，达到覆盖阈值后可以停止解码。"},
        {"level": 1, "text": "覆盖达标前继续扩展，覆盖达标后比较最长单机路径。"},
        {"level": 1, "text": "这与 set-covering 问题“选足够覆盖集合”的逻辑一致。"},
    ], default_size=19.5, space_after=0, line_spacing=0.92)
    return slide


def brkga_fitness_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    top = add_box(slide, 0.62, 1.70, 12.02, 0.76)
    set_textbox(top, [
        {"level": 0, "text": "Algorithm 2 在解码过程中更新覆盖率；达到阈值后，以最长单机路径作为适应度。"},
        {"level": 1, "text": "适应度函数把 Eq. (2) 的瓶颈路径目标和覆盖不足惩罚合并。"},
    ], default_size=20, space_after=0, line_spacing=0.95)

    muted_label(slide, "适应度评价", 0.74, 2.56, 3.80, 0.34)
    add_formula_panel(slide, [
        "Eᵣₑₛ ← Eᵣₑₛ ∪ eᵢ,        Lᵤᵢ ← Lᵤᵢ + d(eᵢ)",
        "δ(C) = evaluateCoverage(Eᵣₑₛ)",
        "F(C) = maxₖ∈K Lₖ + λ·max(0, δd - δ(C))",
    ], 0.82, 2.90, 11.80, 1.20, size=19.5, spacing=1.02)
    if formula_img is not None:
        native_formula(slide, "DECODE", formula_img, 12.93, 7.20, 0.08, 0.05, size=14)

    add_rule(slide, 0.74, 4.46, 11.85)
    left = add_box(slide, 0.70, 4.70, 6.06, 1.76)
    set_textbox(left, [
        {"level": 1, "text": "每加入路径基元，同时更新 Eᵣₑₛ、δ(C) 和 UAV 累计长度。"},
        {"level": 1, "text": "δ(C)<δd 时，惩罚项会压低不可行解在进化过程中的竞争力。"},
        {"level": 1, "text": "δ(C)≥δd 后，适应度主要由最长单机路径决定。"},
    ], default_size=18.8, space_after=0, line_spacing=0.90)

    right = add_box(slide, 6.90, 4.70, 5.72, 1.76)
    set_textbox(right, [
        {"level": 1, "text": "适应度与 SC-VRP 一致：压缩 maxₖLₖ。"},
        {"level": 1, "text": "覆盖不足不直接丢弃，而是通过惩罚保留搜索梯度。"},
        {"level": 1, "text": "BRKGA 可以从不可行解逐步靠近可行覆盖路径。"},
    ], default_size=18.8, space_after=0, line_spacing=0.90)
    return slide


def local_accept_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    top = add_box(slide, 0.62, 1.70, 12.02, 0.86)
    set_textbox(top, [
        {"level": 0, "text": "局部改进沿用 TSP 中的 2-opt 思想，在覆盖图中寻找可替换顶点。"},
        {"level": 1, "text": "替换操作只有在路径变短且覆盖约束仍满足时才有意义。"},
    ], default_size=20, space_after=0, line_spacing=0.95)

    muted_label(slide, "Algorithm 3 中的局部替换", 0.74, 2.68, 4.80, 0.34)
    add_formula_panel(slide, [
        "P′ₖ = 2optSwap(Pₖ, vᵢ, vⱼ),        vⱼ ∈ neighbour(vᵢ, G)",
        "accept(P′) ⇔ δ(P′) ≥ δd  and  maxₖ∈K L′ₖ < maxₖ∈K Lₖ",
    ], 0.84, 3.02, 11.72, 1.02, size=19.2, spacing=1.04)
    if formula_img is not None:
        native_formula(slide, "LOCAL", formula_img, 12.93, 7.20, 0.08, 0.05, size=14)

    add_rule(slide, 0.74, 4.38, 11.85)
    left = add_box(slide, 0.70, 4.62, 6.04, 1.78)
    set_textbox(left, [
        {"level": 1, "text": "候选顶点来自 C-PRM 邻域，搜索范围受图结构限制。"},
        {"level": 1, "text": "2optSwap 只改动局部片段，代价低于整条路径重解码。"},
        {"level": 1, "text": "如果替换导致覆盖率下降到阈值以下，即使路径更短也不接受。"},
    ], default_size=18.8, space_after=0, line_spacing=0.90)

    right = add_box(slide, 6.90, 4.62, 5.72, 1.78)
    set_textbox(right, [
        {"level": 1, "text": "该过程作为 BRKGA 增强算子，利用可行解附近的局部结构。"},
        {"level": 1, "text": "进化搜索做全局探索，局部改进压缩冗余路径段。"},
        {"level": 1, "text": "最终输出仍是覆盖达标且最长单机路径更短的染色体。"},
    ], default_size=18.8, space_after=0, line_spacing=0.90)
    return slide


def local_implementation_slide(prs, part, part_title, sec, sec_title, algo_img: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    left = add_box(slide, 0.62, 1.74, 5.74, 5.02)
    set_textbox(left, [
        {"level": 0, "text": "局部搜索作为额外算子插入 BRKGA。"},
        {"level": 1, "text": "输入包括染色体 C、覆盖图 G(V,E) 与各 UAV 路径 P。"},
        {"level": 1, "text": "遍历每条 UAV 路径节点，并检查图邻域。"},
        {"level": 1, "text": "调用 2optSwap 得到局部更新后的染色体。"},
        {"level": 1, "text": "与变异、交叉、选择算子配合使用。"},
        {"level": 1, "text": "DEAP 负责进化搜索；NetworkX 维护邻接表。"},
    ], default_size=20, space_after=0, line_spacing=0.92)

    muted_label(slide, "Algorithm 3：Local Improvement Heuristics", 6.72, 1.76, 5.70, 0.34)
    add_pic(slide, algo_img, 6.64, 2.18, 5.62, 3.18)

    note = add_box(slide, 6.64, 5.66, 5.78, 1.10)
    set_textbox(note, [
        {"level": 1, "text": "伪代码依次枚举邻域、尝试替换，并返回更新结果。"},
        {"level": 1, "text": "该步骤进一步缩短 BRKGA 已生成的可行路径。"},
    ], default_size=20, space_after=0, line_spacing=0.94)
    return slide


def model_slide(prs, part, part_title, sec, sec_title, model_formula: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    intro = add_box(slide, 0.62, 1.70, 12.02, 0.92)
    set_textbox(intro, [
        {"level": 0, "text": "SC-VRP把C-PRM边选择、UAV分配和访问顺序统一成min-max组合优化。"},
        {"level": 1, "text": "目标压缩最长路径，约束保证覆盖，每条边由二进制变量决定是否选入某架 UAV 路径。"},
    ], default_size=20, space_after=0, line_spacing=0.99)

    muted_label(slide, "ILP 核心公式", 0.70, 2.66, 2.60, 0.30)
    if model_formula is not None:
        add_formula_image(slide, "MODEL", model_formula, 0.70, 2.88, 7.70, 2.84)
        native_formula(slide, "MODEL", model_formula, 12.93, 7.20, 0.08, 0.05, size=14)
    else:
        editable_formula_text(slide, NATIVE_MATH["MODEL"], 0.64, 2.90, 8.02, 3.04, size=FORMULA_SIZE)

    muted_label(slide, "变量与约束解释", 8.58, 2.66, 2.90, 0.30)
    explain = add_box(slide, 8.38, 3.00, 4.28, 2.72)
    set_textbox(explain, [
        {"level": 1, "text": "xᵢⱼₖ=1 表示第 k 架UAV选择C-PRM边 eᵢⱼ。"},
        {"level": 1, "text": "dᵢⱼ是边长，sᵢⱼₘ表示该边是否覆盖第 m 个目标面片。"},
        {"level": 1, "text": "z 约束所有 Lₖ，目标函数对应任务完成时间瓶颈。"},
        {"level": 1, "text": "覆盖率约束对应论文Eq. (8)，允许用δd控制“全覆盖”或“近似全覆盖”。"},
    ], default_size=20, space_after=0, line_spacing=0.97)

    add_rule(slide, 0.72, 6.06, 11.85)
    bottom = add_box(slide, 0.66, 6.22, 12.00, 0.78)
    set_textbox(bottom, [
        {"level": 0, "text": "方法衔接：原始ILP难以直接求大规模实例，再用BRKGA做随机键解码和路径评价。"},
    ], default_size=20, space_after=0, line_spacing=0.98)
    return slide


def image_below_slide(prs, part, part_title, sec, sec_title, bullets, images, cap):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    box = add_box(slide, 0.68, 1.70, 11.95, 1.28)
    set_textbox(box, bullets, default_size=19, space_after=1, line_spacing=1.02)
    for spec in images:
        add_pic(slide, *spec)
    caption(slide, cap, 0.85, 6.28, 11.2, 0.40)
    return slide


def compact_image_below_slide(prs, part, part_title, sec, sec_title, bullets, images, cap):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    box = add_box(slide, 0.68, 1.70, 11.95, 1.00)
    set_textbox(box, bullets, default_size=18.2, space_after=1, line_spacing=1.02)
    for spec in images:
        add_pic(slide, *spec)
    caption(slide, cap, 0.85, 6.28, 11.2, 0.40)
    return slide


def visual_analysis_slide(prs, part, part_title, sec, sec_title, top_bullets, images, analysis_bullets, cap):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    if sec == "2.1":
        top = add_box(slide, 0.62, 1.70, 12.00, 1.06)
        set_textbox(top, top_bullets, default_size=20, space_after=0, line_spacing=0.99)
        add_pic(slide, images[0][0], 1.10, 2.88, 6.72, 2.52)
        analysis = add_box(slide, 7.96, 2.92, 4.82, 2.56)
        set_textbox(analysis, analysis_bullets, default_size=20, space_after=0, line_spacing=0.97)
        add_rule(slide, 0.74, 5.78, 11.80)
        bottom = add_box(slide, 0.66, 5.96, 12.00, 0.84)
        set_textbox(bottom, [
            {"level": 0, "text": "数据流从 3D 结构采样覆盖图，再把候选边交给 SC-VRP/BRKGA 求多机路径。"},
        ], default_size=20, space_after=0, line_spacing=0.98)
    elif sec == "2.2":
        top = add_box(slide, 0.62, 1.70, 12.00, 0.74)
        set_textbox(top, top_bullets, default_size=20.5, space_after=0, line_spacing=0.94)
        add_pic(slide, images[0][0], 0.92, 2.66, 6.10, 2.70)
        add_pic(slide, images[1][0], 7.34, 2.66, 4.72, 2.92)
        add_rule(slide, 0.74, 5.74, 11.80)
        bottom_note = add_box(slide, 0.70, 5.88, 11.90, 1.04)
        set_textbox(bottom_note, analysis_bullets, default_size=18.8, space_after=0, line_spacing=0.88)
    else:
        top = add_box(slide, 0.68, 1.70, 11.95, 0.74)
        set_textbox(top, top_bullets, default_size=20, space_after=0, line_spacing=0.96)
        for spec in images:
            add_pic(slide, *spec)
        analysis = add_box(slide, 0.72, 5.56, 11.80, 1.02)
        set_textbox(analysis, analysis_bullets, default_size=20, space_after=0, line_spacing=0.98)
    return slide


def decode_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    top = add_box(slide, 0.62, 1.70, 12.02, 1.02)
    set_textbox(top, [
        {"level": 0, "text": "BRKGA把随机键向量解码成多UAV路径，再按覆盖率和最长路径评价。"},
        {"level": 1, "text": "整数部分表示UAV编号，小数部分决定在当前节点邻域中选择哪条C-PRM边。"},
    ], default_size=20, space_after=0, line_spacing=0.98)

    muted_label(slide, "编码、覆盖率与适应度公式", 0.72, 2.78, 4.40, 0.30)
    if formula_img is not None:
        add_formula_image(slide, "DECODE", formula_img, 0.78, 3.04, 6.74, 2.58)
        native_formula(slide, "DECODE", formula_img, 12.93, 7.20, 0.08, 0.05, size=14)
    else:
        editable_formula_text(slide, NATIVE_MATH["DECODE"], 0.82, 3.18, 6.36, 2.24, size=FORMULA_SIZE)

    right = add_box(slide, 7.70, 3.00, 4.92, 2.58)
    set_textbox(right, [
        {"level": 0, "text": "解码后的适应度与 min-max 目标一致。"},
        {"level": 1, "text": "uᵢ取随机键整数部分，决定第i个节点或基元分配给哪架UAV。"},
        {"level": 1, "text": "δ(C)按已选边Eres(C)统计被覆盖面片比例，并与δd比较。"},
        {"level": 1, "text": "F(C)用最长单机路径作为主目标，覆盖不足时加入惩罚项。"},
    ], default_size=20, space_after=0, line_spacing=0.96)

    add_rule(slide, 0.70, 5.86, 11.90)
    bottom = add_box(slide, 0.62, 6.04, 12.05, 0.92)
    set_textbox(bottom, [
        {"level": 1, "text": "BRKGA不是直接优化ILP变量，而是先把随机键解码为路径序列，再围绕可覆盖路径搜索。"},
        {"level": 1, "text": "随机键解码降低了直接求解整数规划的压力，同时保留覆盖约束。"},
    ], default_size=20, space_after=0, line_spacing=0.94)
    return slide


def local_improve_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None, algo_img: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    top = add_box(slide, 0.62, 1.70, 12.02, 1.02)
    set_textbox(top, [
        {"level": 0, "text": "局部改进不重新建模，只在已解码路径附近做2-opt替换。"},
        {"level": 1, "text": "对每条UAV路径，枚举路径节点及其在C-PRM图中的邻居，并只接受同时满足覆盖和长度改进的替换。"},
    ], default_size=20, space_after=0, line_spacing=0.98)

    muted_label(slide, "局部替换判据", 0.72, 2.78, 2.80, 0.30)
    if formula_img is not None:
        add_formula_image(slide, "LOCAL", formula_img, 0.72, 3.18, 6.56, 2.18)
        native_formula(slide, "LOCAL", formula_img, 12.93, 7.20, 0.08, 0.05, size=14)
    else:
        editable_formula_text(slide, NATIVE_MATH["LOCAL"], 0.78, 3.10, 6.08, 2.18, size=FORMULA_SIZE)

    explain = add_box(slide, 7.72, 2.98, 4.92, 2.38)
    set_textbox(explain, [
        {"level": 1, "text": "P'ₖ是在当前路径Pₖ附近生成的2-opt候选路径。"},
        {"level": 1, "text": "δ(P') ≥ δd 保证替换后覆盖率仍达到实验阈值。"},
        {"level": 1, "text": "maxₖL'ₖ < maxₖLₖ 表示接受替换后瓶颈路径确实变短。"},
    ], default_size=20, space_after=0, line_spacing=0.96)

    muted_label(slide, "实现链路", 0.72, 5.62, 2.10, 0.30)
    impl = add_box(slide, 0.68, 5.92, 5.72, 0.98)
    set_textbox(impl, [
        {"level": 1, "text": "DEAP负责进化搜索，NetworkX维护C-PRM邻接关系。"},
        {"level": 1, "text": "局部搜索作为额外算子插入BRKGA流程。"},
    ], default_size=20, space_after=0, line_spacing=0.90)

    muted_label(slide, "对应伪代码", 6.82, 5.58, 2.10, 0.30)
    add_pic(slide, algo_img, 6.82, 5.94, 5.04, 0.92)

    return slide


def visual_side_analysis_slide(prs, part, part_title, sec, sec_title, bullets, images, cap):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    if sec == "3.1":
        box = add_box(slide, 0.62, 1.70, 5.70, 4.62)
        set_textbox(box, bullets, default_size=20, space_after=0, line_spacing=0.96)
        add_pic(slide, images[0][0], 6.58, 1.84, 5.54, 4.26)
        note = add_box(slide, 6.42, 6.18, 6.18, 0.68)
        set_textbox(note, [
            {"level": 1, "text": "右图展示6类测试结构，几何尺度和形状差异用于检验C-PRM采样与多机分配的稳定性。"},
        ], default_size=20, space_after=0, line_spacing=0.95)
    elif sec == "3.4":
        box = add_box(slide, 0.62, 1.70, 5.58, 4.70)
        set_textbox(box, bullets, default_size=20, space_after=0, line_spacing=0.96)
        add_pic(slide, images[0][0], 6.44, 1.78, 5.88, 4.34)
        note = add_box(slide, 6.34, 6.18, 6.20, 0.70)
        set_textbox(note, [
            {"level": 1, "text": "右图是Octomap体素重建结果，用来验证优化路径不仅短，而且能够形成连续表面覆盖。"},
        ], default_size=20, space_after=0, line_spacing=0.95)
    else:
        box = add_box(slide, 0.52, 1.70, 5.18, 5.70)
        set_textbox(box, bullets, default_size=20, space_after=0, line_spacing=0.98)
        for spec in images:
            add_pic(slide, *spec)
    return slide


def split_slide(prs, part, part_title, sec, sec_title, bullets, images, cap):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    box = add_box(slide, 0.50, 1.70, 4.82, 5.35)
    set_textbox(box, bullets, default_size=18.8, space_after=0, line_spacing=0.98)
    for spec in images:
        add_pic(slide, *spec)
    caption(slide, cap, 5.22, 6.18, 7.25, 0.42)
    return slide


def add_table(slide, left, top, width, height, rows, font_size=12):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(242, 242, 242) if ri == 0 else RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                set_run_font(run)
                run.font.size = Pt(TABLE_SIZE)
                run.font.bold = ri == 0 or ci == 0
                run.font.color.rgb = RED if str(val) in {"BRKGA+", "150.8", "177.8", "1349.2", "24.6%"} else BLACK
    return shape


def table_slide(prs, part, part_title, sec, sec_title, bullets, rows, cap):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    box = add_box(slide, 0.50, 1.70, 4.66, 5.50)
    set_textbox(box, bullets, default_size=18.8, space_after=0, line_spacing=0.98)
    add_table(slide, 5.18, 1.92, 7.25, 4.18, rows, 13.4)
    return slide


def wide_table_slide(prs, part, part_title, sec, sec_title, bullets, rows, cap):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)
    box = add_box(slide, 0.50, 1.70, 12.28, 1.20)
    set_textbox(box, bullets, default_size=18.8, space_after=0, line_spacing=0.98)
    add_table(slide, 0.78, 3.08, 11.82, 3.05, rows, 13.3)
    return slide


def parameter_slide(prs, part, part_title, sec, sec_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    top = add_box(slide, 0.56, 1.70, 12.10, 1.02)
    set_textbox(top, [
        {"level": 0, "text": "实验参数用于保证路径长度比较可解释：相机、覆盖率、最大观测距离和运行次数都固定。"},
        {"level": 1, "text": "planned path length取BRKGA/BRKGA+的10次运行平均值，降低随机启发式波动。"},
    ], default_size=18.2, space_after=0, line_spacing=0.98)

    rows = [
        ["参数", "T1/T2 小规模", "T3-T6 大规模", "说明"],
        ["diagonal FOV", "94°", "94°", "相机视野固定，路径差异来自规划算法"],
        ["viewing angle", "75°", "75°", "限制斜视角，保证可观测质量"],
        ["safety distance", "2 m", "2 m", "路径基元需要保持安全间隔"],
        ["max range", "50 m", "70 m", "大结构允许更远观测距离"],
        ["coverage ratio δd", "99%", "98%", "覆盖阈值进入 Eq. (8) 与解码停止条件"],
        ["算法对比", "VPP-TSP / GNS / BRKGA / BRKGA+", "GNS / BRKGA / BRKGA+", "验证 SC-VRP + BRKGA+ 的路径压缩效果"],
    ]
    add_table(slide, 0.66, 2.82, 12.02, 3.22, rows, 12.2)

    bottom = add_box(slide, 0.56, 6.16, 12.10, 0.92)
    set_textbox(bottom, [
        {"level": 1, "text": "评价指标是最长单机路径 maxₖ Lₖ，对应任务完成时间瓶颈；不是所有UAV路径长度之和。"},
        {"level": 1, "text": "增加UAV数量通常会降低cost，但下降幅度不会线性增长，T3-T6表格也体现这一点。"},
    ], default_size=16.8, space_after=0, line_spacing=0.98)
    return slide


def method_intro(slide, lines: list[dict], *, top=1.70, height=0.92, size=21.0) -> None:
    box = add_box(slide, 0.62, top, 12.05, height)
    set_textbox(box, lines, default_size=size, space_after=0, line_spacing=0.95)


def method_label(slide, text: str, left: float, top: float, width: float) -> None:
    box = add_box(slide, left, top, width, 0.34)
    plain(box, text, 20, bold=True, color=GRAY)


def equation_block(slide, rows: list[tuple], left: float, top: float, width: float, height: float, *, size=22.0) -> None:
    box = add_box(slide, left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for idx, row in enumerate(rows):
        label, text = row[0], row[1]
        row_size = float(row[2]) if len(row) > 2 else size
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.04
        p.space_before = Pt(0)
        p.space_after = Pt(2 if idx < len(rows) - 1 else 0)
        r_label = p.add_run()
        r_label.text = f"{label}    "
        set_run_font(r_label, BODY_FONT)
        r_label.font.size = Pt(row_size)
        r_label.font.bold = True
        r_label.font.color.rgb = GRAY
        r_formula = p.add_run()
        r_formula.text = text
        set_run_font(r_formula, BODY_FONT)
        r_formula.font.size = Pt(row_size)
        r_formula.font.color.rgb = BLACK


def method_notes(slide, lines: list[dict], *, top=5.05, height=1.70, size=20.0) -> None:
    box = add_box(slide, 0.78, top, 11.82, height)
    set_textbox(box, lines, default_size=size, space_after=0, line_spacing=0.94)


def sc_vrp_objective_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    method_intro(slide, [
        {"level": 0, "text": "SC-VRP 把 C-PRM 路径基元、UAV 分配和覆盖约束写入同一个 min-max 模型。"},
        {"level": 1, "text": "变量 x̂ᵢⱼₖ 表示第 k 架 UAV 是否选择边 eᵢⱼ，z 表示最长单机路径上界。"},
    ])

    method_label(slide, "目标函数与路径长度上界", 0.82, 2.74, 4.70)
    equation_block(slide, [
        ("(2a)", "min z", 23.0),
        ("(2b)", "Lₖ = Σ dᵢⱼ x̂ᵢⱼₖ,  (i,j)∈E,  ∀ k∈K", 22.0),
        ("(2c)", "Lₖ ≤ z,  ∀ k∈K", 22.0),
        ("(2d)", "x̂ᵢⱼₖ ∈ {0,1},  (i,j)∈E,  k∈K", 22.0),
    ], 1.18, 3.15, 11.15, 1.32, size=22.0)

    add_rule(slide, 0.74, 4.93, 11.85)
    method_notes(slide, [
        {"level": 1, "text": "dᵢⱼ 是路径基元长度；对同一 k 求和得到该 UAV 的巡检路径长度 Lₖ。"},
        {"level": 1, "text": "min z 等价于最小化 maxₖ Lₖ，把任务完成时间瓶颈直接放进目标函数。"},
        {"level": 1, "text": "如果只压缩总路程，可能让单机承担过多路径；min-max 目标会主动均衡多机负载。"},
        {"level": 1, "text": "BRKGA 的适应度仍围绕 Lₖ 和 z 设计，使建模目标与启发式搜索保持一致。"},
    ], top=5.13, height=1.62, size=19.4)
    return slide


def sc_vrp_constraints_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None, ratio_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    method_intro(slide, [
        {"level": 0, "text": "覆盖约束要求每个目标面片至少被一条已选路径基元观测，同时排除断裂子回路。"},
        {"level": 1, "text": "sᵢⱼₘ 记录边 eᵢⱼ 对面片 m 的可见性收益，δd 是实验设定的覆盖率阈值。"},
    ])

    method_label(slide, "覆盖率与连通约束", 0.82, 2.70, 4.40)
    equation_block(slide, [
        ("(3)", "Σₖ Σᵢⱼ sᵢⱼₘ x̂ᵢⱼₖ ≥ 1,  ∀ m∈M", 22.0),
        ("(4)", "Σᵢⱼ∈L x̂ᵢⱼₖ ≤ |L| − 1,  L⊂V′ₖ, |L|≥2", 22.0),
        ("(5)", "δ(P)= |covered faces| / |M|", 22.0),
        ("(6)", "δ(P) ≥ δd", 23.0),
    ], 1.00, 3.12, 11.55, 1.46, size=22.0)

    add_rule(slide, 0.74, 4.82, 11.85)
    method_notes(slide, [
        {"level": 1, "text": "式 (3) 保证目标面片被覆盖；式 (4) 限制同一 UAV 的边形成连通巡检序列。"},
        {"level": 1, "text": "δ(P) 是已覆盖面片比例；T1/T2 取 99%，T3-T6 取 98%。"},
        {"level": 1, "text": "覆盖阈值越高，路径更难缩短，但能避免得到“短而漏检”的规划结果。"},
        {"level": 1, "text": "解码过程中每加入一条边都更新 δ(P)，达到 δd 后再比较最长单机路径。"},
    ], top=5.02, height=1.76, size=20.0)
    return slide


def brkga_encoding_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    method_intro(slide, [
        {"level": 0, "text": "BRKGA 不直接搜索大量 0/1 变量，而是先生成连续随机键，再解码为多 UAV 路径。"},
        {"level": 1, "text": "一个随机键同时决定 UAV 分配和邻域边选择，降低整数规划求解压力。"},
    ])

    method_label(slide, "随机键编码", 0.82, 2.78, 3.20)
    equation_block(slide, [
        ("(7a)", "C=(c₁,c₂,…,cₙ)ᵀ,    cᵢ∈[0,K]", 23.0),
        ("(7b)", "uᵢ=⌊cᵢ⌋,    qᵢ=cᵢ−uᵢ", 23.0),
        ("(7c)", "uᵢ → UAV index,    qᵢ → edge rank in N(vᵢ)", 22.0),
    ], 1.20, 3.18, 11.10, 1.22, size=23.0)

    add_rule(slide, 0.74, 4.70, 11.85)
    method_notes(slide, [
        {"level": 1, "text": "整数部 uᵢ 决定路径基元分配给哪架 UAV；小数部 qᵢ 在当前节点邻域中排序候选边。"},
        {"level": 1, "text": "解码不是访问所有节点，而是边扩展、覆盖更新、阈值判断交替进行。"},
        {"level": 1, "text": "当 δ(P)<δd 时继续扩展路径；当 δ(P)≥δd 时，解的优劣主要由 maxₖ Lₖ 决定。"},
        {"level": 1, "text": "这种编码让遗传算子在连续空间中工作，同时保留 SC-VRP 的覆盖约束含义。"},
    ], top=4.92, height=1.92, size=20.0)
    return slide


def brkga_fitness_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    method_intro(slide, [
        {"level": 0, "text": "适应度函数把覆盖不足惩罚和最长单机路径合并，使不可行解逐步靠近可行路径。"},
        {"level": 1, "text": "Algorithm 2 每扩展一条边就更新 Eres、覆盖率 δ(C) 和各 UAV 的累计路径长度。"},
    ])

    method_label(slide, "覆盖惩罚与 min-max 目标", 0.82, 2.72, 4.90)
    equation_block(slide, [
        ("(8a)", "Eres ← Eres ∪ {eᵢ},    Lᵤᵢ ← Lᵤᵢ+d(eᵢ)", 21.6),
        ("(8b)", "δ(C)=evaluateCoverage(Eres)", 22.0),
        ("(8c)", "F(C)=maxₖ Lₖ + λ·max(0, δd−δ(C))", 22.0),
        ("(8d)", "δ(C)≥δd  ⇒  F(C)=maxₖ Lₖ", 22.0),
    ], 1.08, 3.06, 11.35, 1.58, size=22.0)

    add_rule(slide, 0.74, 4.90, 11.85)
    method_notes(slide, [
        {"level": 1, "text": "δ(C)<δd 时，惩罚项 λ·max(0, δd−δ(C)) 会降低该染色体的竞争力。"},
        {"level": 1, "text": "覆盖达标后，惩罚项为 0，适应度回到最长单机路径目标。"},
        {"level": 1, "text": "保留惩罚而不是直接丢弃不可行解，可以在早期维持搜索梯度。"},
        {"level": 1, "text": "该设计把 SC-VRP 中的覆盖约束和 min-max 目标连接到进化搜索过程。"},
    ], top=5.10, height=1.78, size=20.0)
    return slide


def local_accept_slide(prs, part, part_title, sec, sec_title, formula_img: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    method_intro(slide, [
        {"level": 0, "text": "局部改进在 BRKGA 生成的可行路径附近做 2-opt 替换，只接受覆盖达标且更短的候选解。"},
        {"level": 1, "text": "候选顶点来自 C-PRM 图邻域，搜索范围由图结构限制，计算代价低于整条路径重解码。"},
    ], height=1.02)

    method_label(slide, "Algorithm 3 中的替换判据", 0.82, 2.86, 5.30)
    equation_block(slide, [
        ("(9a)", "P′ₖ = 2optSwap(Pₖ, vᵢ, vⱼ),    vⱼ∈N(vᵢ)", 22.0),
        ("(9b)", "accept(P′) ⇔ δ(P′)≥δd ∧ maxₖ L′ₖ < maxₖ Lₖ", 22.0),
        ("(9c)", "Pᵗ⁺¹=P′, if accept(P′);    Pᵗ⁺¹=Pᵗ, otherwise", 21.6),
    ], 1.02, 3.24, 11.45, 1.30, size=22.0)

    add_rule(slide, 0.74, 4.82, 11.85)
    method_notes(slide, [
        {"level": 1, "text": "2optSwap 只改变局部片段，先尝试替换相邻连接，再重新计算覆盖率和路径长度。"},
        {"level": 1, "text": "如果替换导致 δ(P′)<δd，即使路径变短也不会接受。"},
        {"level": 1, "text": "如果覆盖仍达标且 maxₖ L′ₖ 下降，局部改进会压缩任务瓶颈路径。"},
        {"level": 1, "text": "全局探索由 BRKGA 完成，局部改进负责清理冗余绕行，两者共同改善路径质量。"},
    ], top=5.02, height=1.86, size=20.0)
    return slide


def local_implementation_slide(prs, part, part_title, sec, sec_title, algo_img: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, part, part_title)
    add_section_title(slide, sec, sec_title)

    left = add_box(slide, 0.62, 1.72, 5.80, 3.88)
    set_textbox(left, [
        {"level": 0, "text": "局部搜索作为额外算子插入 BRKGA。"},
        {"level": 1, "text": "输入：染色体 C、覆盖图 G、路径 P。"},
        {"level": 1, "text": "遍历 UAV 路径节点，检查 C-PRM 邻域。"},
        {"level": 1, "text": "用 2optSwap 生成候选染色体 Co。"},
        {"level": 1, "text": "只保留覆盖达标且更短的更新。"},
        {"level": 1, "text": "DEAP 做进化搜索，NetworkX 维护邻接表。"},
    ], default_size=20.0, space_after=0, line_spacing=0.92)

    method_label(slide, "Algorithm 3：局部改进伪代码", 6.82, 1.72, 5.10)
    add_pic(slide, algo_img, 6.70, 2.18, 5.42, 3.32)

    add_rule(slide, 0.74, 5.78, 11.85)
    note = add_box(slide, 0.78, 5.98, 11.82, 0.86)
    set_textbox(note, [
        {"level": 1, "text": "伪代码依次枚举邻域、尝试替换，并返回更新后的染色体。"},
        {"level": 1, "text": "该步骤不是重新求解 SC-VRP，而是在可行解附近压缩冗余路径段。"},
    ], default_size=20.0, space_after=0, line_spacing=0.92)
    return slide


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.10), Inches(0.65), width=Inches(5.08), height=Inches(1.41))
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.42), Inches(13.333), Inches(3.12))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    title = add_box(slide, 0.0, 3.05, 13.33, 0.85)
    plain(title, "Multi-UAV Coverage Path Planning\nfor Structural Inspection", 30, bold=True, align=PP_ALIGN.CENTER, allow_newlines=True)
    subtitle = add_box(slide, 0.0, 4.02, 13.33, 1.05)
    plain(subtitle, "Wei Jing, Di Deng, Yan Wu, Kenji Shimada\nIROS 2020 / arXiv:2007.13065", 18, align=PP_ALIGN.CENTER, allow_newlines=True)
    date = add_box(slide, 5.50, 6.62, 2.35, 0.45)
    plain(date, "2026年5月6日", 18, align=PP_ALIGN.CENTER)
    return slide


def thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.42), Inches(13.333), Inches(3.12))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    box = add_box(slide, 0, 3.44, 13.33, 0.9)
    plain(box, "谢谢！", 44, bold=True, align=PP_ALIGN.CENTER)
    date = add_box(slide, 5.50, 6.62, 2.35, 0.45)
    plain(date, "2026年5月6日", 18, align=PP_ALIGN.CENTER)
    return slide


NATIVE_MATH = {
    "MODEL": [
        "minimize      z",
        "s.t.          L_k = Σ_(i,j)∈E d_ij x_ijk,        ∀ k∈K",
        "              L_k ≤ z,                           ∀ k∈K",
        "              Σ_(i,j)∈E Σ_k∈K s_ijm x_ijk ≥ 1,    ∀ m∈M",
        "              x_ijk ∈ {0,1},                      ∀ (i,j)∈E, k∈K",
    ],
    "DECODE": [
        "C_k(P_k) = ∨_((i,j)∈P_k) s_ijm,                 m∈M",
        "ρ(P) = 1/|M| Σ_(m∈M) 1[Σ_(k∈K) C_k(P_k)_m ≥ 1]",
        "F(P) = max_(k∈K) L_k + λ max(0, δ_d - ρ(P))",
    ],
    "LOCAL": [
        "P'_k = 2optSwap(P_k, v_i, v_j)",
        "ρ(P') ≥ δ_d",
        "max_k L'_k < max_k L_k",
        "P^(t+1)=P' if accepted; otherwise P^(t+1)=P^t",
    ],
}


def assert_layout(prs: Presentation) -> None:
    def bounds(shape):
        return (
            shape.left / 914400,
            shape.top / 914400,
            (shape.left + shape.width) / 914400,
            (shape.top + shape.height) / 914400,
        )

    def area(a, b):
        return max(0, min(a[2], b[2]) - max(a[0], b[0])) * max(0, min(a[3], b[3]) - max(a[1], b[1]))

    failures = []
    last_slide = len(prs.slides)
    for si, slide in enumerate(prs.slides, 1):
        if si in (1, last_slide):
            continue
        items = []
        for idx, shape in enumerate(slide.shapes):
            if getattr(shape, "name", "").startswith("MATH_"):
                continue
            left, top, right, bottom = bounds(shape)
            if left < -0.02 or right > 13.35 or bottom > 7.45:
                failures.append(f"slide {si} shape {idx} outside {left:.2f},{top:.2f},{right:.2f},{bottom:.2f}")
            if top < 1.05 or (left < 0.75 and top < 0.25):
                continue
            label = None
            if shape.shape_type == 13:
                label = "pic"
            elif shape.shape_type == 19:
                label = "table"
            elif hasattr(shape, "text") and shape.text.strip():
                label = "txt"
            if label:
                items.append((label, idx, bounds(shape)))
        for a in range(len(items)):
            for b in range(a + 1, len(items)):
                if area(items[a][2], items[b][2]) > 0.01:
                    failures.append(f"slide {si} {items[a][0]}#{items[a][1]} overlaps {items[b][0]}#{items[b][1]}")
    if failures:
        raise RuntimeError("\n".join(failures))


def main() -> None:
    assets = ROOT / "paper_crops" / "clean"
    assets.mkdir(exist_ok=True)
    framework = crop_box(CROPS / "framework_page2.png", assets / "framework_core.png", (0.00, 0.00, 1.00, 0.70))
    cprm = crop_box(CROPS / "cprm_page3.png", assets / "cprm_core.png", (0.00, 0.00, 1.00, 0.77))
    algo1 = crop_whitespace(CROPS / "algorithm1_page3.png", assets / "algorithm1_clean.png")
    algo3 = crop_box(CROPS / "algorithms_2_3_page4.png", assets / "algorithm3_clean.png", (0.50, 0.00, 1.00, 0.92))
    structures = make_structure_grid(CROPS / "target_structures_page4.png", assets / "structures_grid_clean.png")
    coverage = make_coverage_grid(CROPS / "coverage_fig8_clean.png", assets / "coverage_grid_clean.png")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    body_slide(prs, "Part. 01", "研究背景及动机", "1.1", "论文与作者信息", [
        {"text": "论文研究大型复杂3D结构的多UAV视觉巡检。"},
        {"level": 1, "text": "IROS 2020 / arXiv:2007.13065，作者为 Wei Jing, Di Deng, Yan Wu, Kenji Shimada。"},
        {"level": 1, "text": "目标是在覆盖完整前提下缩短最长单机路径。"},
        {"level": 1, "text": "问题同时涉及覆盖选择、路径连接和多机分配。"},
        {"level": 2, "text": "与只规划离散拍照点不同，本文把路径边也视为可产生覆盖收益的基本单元。"},
        {"level": 2, "text": "内容围绕 SC-VRP 建模、C-PRM 采样、BRKGA 求解和路径长度下降展开。"},
    ])
    body_slide(prs, "Part. 01", "研究背景及动机", "1.2", "研究背景", [
        {"text": "大型复杂结构巡检需要同时保证安全、效率和覆盖完整性。"},
        {"level": 1, "text": "单机续航通常只有20-40 min，多机协同是大尺度巡检的现实方案。"},
        {"level": 1, "text": "路径必须让相机覆盖目标表面，同时减少重复观测和无效飞行。"},
        {"level": 1, "text": "多机规划需要同时处理路径基元选择、访问顺序和UAV分配。"},
        {"level": 2, "text": "优化目标不是总路程最短，而是压缩最长单机路径，避免任务瓶颈。"},
        {"level": 2, "text": "覆盖约束要求目标面片被可见路径基元覆盖，且整体覆盖率达到阈值。"},
    ])
    body_slide(prs, "Part. 01", "研究背景及动机", "1.3", "先前方法：视点规划", [
        {"text": "传统流程通常先生成候选视点，再用 TSP 或局部规划器连接路径。"},
        {"level": 1, "text": "优点：流程清晰，适合单 UAV 小规模任务。"},
        {"level": 1, "text": "不足：只把离散视点作为覆盖单元，忽略飞行路径本身的覆盖收益。"},
        {"level": 1, "text": "本文改为采样 path primitive，并为每条边记录可见性向量。"},
        {"level": 2, "text": "若两个视点之间的飞行段本身可观察大量表面，传统 VPP-TSP 往往无法利用这部分信息。"},
        {"level": 2, "text": "因此本文的覆盖单元从“点”扩展到“边/路径基元”，表达能力更接近真实飞行过程。"},
    ])
    body_slide(prs, "Part. 01", "研究背景及动机", "1.4", "先前方法：多机器人规划", [
        {"text": "多机器人 CPP 能表达任务分配，但很多方法主要针对 2D 区域覆盖。"},
        {"level": 1, "text": "VRP / mTSP 可以优化多车路线，但不能自然表达“某条路径覆盖一组面片”。"},
        {"level": 1, "text": "本文把 SCP 与 VRP 合并：SCP 负责覆盖集合，VRP 负责多机路径。"},
        {"level": 1, "text": "得到的问题称为 Set-Covering Vehicle Routing Problem，即 SC-VRP。"},
        {"level": 2, "text": "SCP 解决“选哪些路径基元才能覆盖全部目标”，VRP 解决“这些基元如何分配给多架 UAV 并排序”。"},
        {"level": 2, "text": "这个组合建模使覆盖约束和路径均衡同时出现，同时带来更高的问题规模和求解难度。"},
    ])
    body_slide(prs, "Part. 01", "研究背景及动机", "1.5", "技术要点", [
        {"text": "本文从采样表示、组合建模和启发式求解三个层面展开。"},
        {"level": 1, "text": "提出 multi-UAV CPP 框架，用于大型复杂 3D 结构巡检。"},
        {"level": 1, "text": "提出 minmax SC-VRP，将 Set Covering Problem 与 Vehicle Routing Problem 结合。"},
        {"level": 1, "text": "提出 BRKGA 编码/解码策略，并加入局部改进启发式。"},
        {"level": 2, "text": "在 6 个 OpenStreetMap 模型上评估，最高降低约 48% 路径长度。"},
        {"level": 2, "text": "实验同时比较路径长度，并用Drake + Octomap验证表面重建效果。"},
        {"level": 2, "text": "关键价值在于把覆盖判定、路径排序和任务分配统一到一个组合优化问题中。"},
    ])

    visual_analysis_slide(prs, "Part. 02", "研究方法", "2.1", "整体框架", [
        {"level": 1, "text": "Coverage Sampling 将 3D mesh 转换为 C-PRM；BRKGA 在覆盖约束下生成多 UAV 路径。"},
        {"level": 2, "text": "覆盖评估和多机路径分配在同一个优化闭环中完成。"},
    ], [(framework, 1.70, 2.50, 9.95, 2.82)], [
        {"level": 1, "text": "左侧输出 C-PRM graph：节点/边来自路径基元采样，每条边关联可见性评估结果。"},
        {"level": 1, "text": "右侧求解 path planning and optimization：先可用 ILP 表达，再用 BRKGA 处理大规模 NP-hard 搜索。"},
    ], "Fig. 2：论文整体框架。")

    visual_analysis_slide(prs, "Part. 02", "研究方法", "2.2", "C-PRM 与覆盖采样", [
        {"level": 1, "text": "路径基元计算可见性向量，DualSampleVP优先补充未覆盖表面。"},
        {"level": 1, "text": "C-PRM的边同时承担连接和覆盖收益。"},
    ], [(cprm, 0.82, 2.58, 5.55, 2.62), (algo1, 7.02, 2.48, 4.28, 2.92)], [
        {"level": 1, "text": "左图：灰色为结构表面，绿色连线为候选路径基元。"},
        {"level": 1, "text": "右图：Algorithm 1根据未覆盖集合更新采样方向。"},
        {"level": 1, "text": "未覆盖集合反馈给DualSampleVP，推动图向覆盖薄弱处扩展。"},
        {"level": 1, "text": "每条边保存可见性向量，SC-VRP据此检查覆盖率。"},
    ], "Fig. 3 + Algorithm 1：覆盖图与双路径基元采样。")

    model_objective_formula = make_model_objective_formula(assets / "sc_vrp_objective_formula.png")
    model_constraints_formula = make_model_constraints_formula(assets / "sc_vrp_constraints_formula.png")
    model_coverage_ratio_formula = make_model_coverage_ratio_formula(assets / "sc_vrp_coverage_ratio_formula.png")
    decode_encoding_formula = make_decode_encoding_formula(assets / "brkga_encoding_formula.png")
    decode_fitness_formula = make_decode_fitness_formula(assets / "brkga_fitness_formula.png")
    local_accept_formula = make_local_accept_formula(assets / "local_accept_formula.png")

    sc_vrp_objective_slide(prs, "Part. 02", "研究方法", "2.3", "SC-VRP 数学建模（一）：目标函数", model_objective_formula)
    sc_vrp_constraints_slide(prs, "Part. 02", "研究方法", "2.4", "SC-VRP 数学建模（二）：覆盖约束", model_constraints_formula, model_coverage_ratio_formula)
    brkga_encoding_slide(prs, "Part. 02", "研究方法", "2.5", "BRKGA 编码与解码（一）：随机键编码", decode_encoding_formula)
    brkga_fitness_slide(prs, "Part. 02", "研究方法", "2.6", "BRKGA 编码与解码（二）：适应度", decode_fitness_formula)
    local_accept_slide(prs, "Part. 02", "研究方法", "2.7", "局部改进（一）：替换判据", local_accept_formula)
    local_implementation_slide(prs, "Part. 02", "研究方法", "2.8", "局部改进（二）：实现链路", algo3)
    parameter_slide(prs, "Part. 02", "研究方法", "2.9", "实验参数设置")

    visual_side_analysis_slide(prs, "Part. 03", "研究结果分析", "3.1", "目标结构与规模", [
        {"text": "实验覆盖6类OSM结构：从小型建筑到高层模型。"},
        {"level": 1, "text": "T1/T2：阈值99%，用于单机/三机对比。"},
        {"level": 1, "text": "T3-T6：阈值98%，用于大规模多机实验。"},
        {"level": 2, "text": "结构越大，路径平滑性与最长单机路径越关键。"},
        {"level": 2, "text": "T6约159×319×109m，检验多机协同与minmax均衡。"},
        {"level": 2, "text": "几何差异会改变可见性图密度、路径基元长度和局部改进空间。"},
    ], [(structures, 6.16, 1.80, 5.12, 4.62)], "")

    table_slide(prs, "Part. 03", "研究结果分析", "3.2", "T1/T2 小规模实验", [
        {"text": "小规模实验显示，BRKGA+ 对单机和多机情形都明显缩短规划路径。"},
        {"level": 1, "text": "T1 三机：GNS 为 313.1m，BRKGA+ 降到 150.8m，降低约 51.8%。"},
        {"level": 1, "text": "T2 三机：GNS 为 344.1m，BRKGA+ 降到 177.8m，降低约 48.3%。"},
        {"level": 2, "text": "多UAV场景下，BRKGA+相比GNS平均降低48.0%。"},
        {"level": 2, "text": "VPP-TSP 只能处理单机 CPP，因此多机列为空；这也凸显 SC-VRP 的建模优势。"},
    ], [
        ["方法", "T1-1", "T1-2", "T1-3", "T2-1", "T2-2", "T2-3"],
        ["VPP-TSP", "507.7", "-", "-", "587.5", "-", "-"],
        ["GNS", "425.6", "325.9", "313.1", "466.2", "441.2", "344.1"],
        ["BRKGA", "277.9", "197.2", "168.9", "335.3", "232.3", "201.8"],
        ["BRKGA+", "271.1", "186.4", "150.8", "308.1", "225.6", "177.8"],
    ], "Table I：T1/T2 planned path length，单位 m，10 次运行平均。")

    wide_table_slide(prs, "Part. 03", "研究结果分析", "3.3", "T3-T6 路径长度", [
        {"level": 1, "text": "大规模结构上 BRKGA+ 相比 GNS 平均降低 24.6%；T6 四机从 1811.9m 降到 1349.2m。"},
        {"level": 2, "text": "增加 UAV 数量会降低最长单机路径，但收益不是线性增长；T5/T6 仍保持 20% 以上降幅。"},
    ], [
        ["对象", "UAV", "GNS", "BRKGA", "BRKGA+", "降幅"],
        ["T3", "3", "527.6", "396.1", "357.0", "32.3%"],
        ["T4", "4", "542.5", "500.3", "430.9", "20.6%"],
        ["T5", "4", "1373.4", "1167.8", "1062.3", "22.6%"],
        ["T6", "4", "1811.9", "1529.9", "1349.2", "25.5%"],
    ], "Table II 摘要：T3-T6 planned path length，单位 m；降幅为 BRKGA+ 相比 GNS。")

    visual_side_analysis_slide(prs, "Part. 03", "研究结果分析", "3.4", "覆盖质量验证", [
        {"text": "论文还用Octomap验证覆盖质量。"},
        {"level": 1, "text": "Drake模拟多UAV飞行与传感器测量。"},
        {"level": 1, "text": "Octomap构建体素地图，检查覆盖质量。"},
        {"level": 2, "text": "该步骤连接优化结果与实际可观测性。"},
        {"level": 2, "text": "只看路径长度可能得到“短但漏检”的解；体素重建可暴露该问题。"},
        {"level": 2, "text": "各结构均有连续表面重建，说明规划没有只覆盖局部。"},
    ], [(coverage, 6.18, 1.88, 5.18, 4.42)], "")

    body_slide(prs, "Part. 04", "结论与思考", "4.1", "结论", [
        {"text": "本文面向大型复杂结构巡检，提出multi-UAV CPP框架。"},
        {"level": 1, "text": "path primitive作为覆盖单元，利用飞行段中的可见性收益。"},
        {"level": 1, "text": "SC-VRP统一表达覆盖选择、路径连接与多机分配。"},
        {"level": 1, "text": "BRKGA+在复杂组合空间中搜索短路径。"},
        {"level": 2, "text": "多机巡检应从覆盖约束出发，再联合优化分配和路径长度。"},
        {"level": 2, "text": "局限：路径平滑性、动力学、通信和在线避障尚未完整纳入模型。"},
        {"level": 2, "text": "可加入smoothness、uncertainty handling和在线重规划。"},
        {"level": 2, "text": "局限之外，本文仍提供了可扩展的大规模结构巡检建模框架。"},
    ])
    body_slide(prs, "Part. 05", "研究方向", "5.1", "研究方向", [
        {"text": "可扩展为语言指令驱动的多UAV巡检规划。"},
        {"level": 1, "text": "LLM解析覆盖区域、优先级、禁飞区和期望覆盖率。"},
        {"level": 1, "text": "优化模块转成SC-VRP/VRP/orienteering等问题。"},
        {"level": 1, "text": "飞控层把路径基元转为可执行轨迹，并在线重规划。"},
        {"level": 2, "text": "任务输入可包含目标类型、重点区域、覆盖阈值和UAV数量。"},
        {"level": 2, "text": "仿真器检查覆盖缺口，并反馈给采样模块。"},
        {"level": 2, "text": "最终目标是从任务描述到可执行巡检路径形成自动闭环。"},
    ])
    thanks(prs)

    assert_layout(prs)
    prs.save(OUT)
    if MATH_REPLACEMENTS:
        inject_native_math(OUT, MATH_REPLACEMENTS)


if __name__ == "__main__":
    main()
