from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
EXAMPLES = ROOT / "assets" / "examples"

BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(72, 72, 72)
RED = RGBColor(220, 0, 0)
PANEL = RGBColor(242, 244, 247)
PANEL_LINE = RGBColor(214, 218, 224)


def inch(value: float):
    return Inches(value)


def set_bounds(shape, left: float, top: float, width: float, height: float) -> None:
    shape.left = inch(left)
    shape.top = inch(top)
    shape.width = inch(width)
    shape.height = inch(height)


def remove_text_body(shape) -> None:
    tx_body = shape._element.find(qn("p:txBody"))
    if tx_body is not None:
        shape._element.remove(tx_body)


def send_panel_back(slide, shape) -> None:
    tree = slide.shapes._spTree
    element = shape._element
    tree.remove(element)
    # Keep the master-covering background behind everything, then place panels
    # before slide content so they cannot cover text after LibreOffice export.
    tree.insert(4, element)


def add_panel(slide, left: float, top: float, width: float, height: float, name: str) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        inch(left),
        inch(top),
        inch(width),
        inch(height),
    )
    panel.name = name
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = PANEL_LINE
    panel.line.width = Pt(0.7)
    remove_text_body(panel)
    send_panel_back(slide, panel)


def add_density_text(slide, left: float, top: float, width: float, height: float, name: str, lines: list[str]) -> None:
    box = slide.shapes.add_textbox(inch(left), inch(top), inch(width), inch(height))
    box.name = name
    set_bullet_text(box, lines)


def set_run_font(run, size: float, *, bold: bool = False, color: RGBColor = BLACK) -> None:
    run.font.name = "Times New Roman"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    rpr.set("lang", "zh-CN")
    for tag in ("latin", "ea", "cs"):
        node = rpr.find(qn(f"a:{tag}"))
        if node is None:
            node = OxmlElement(f"a:{tag}")
            rpr.append(node)
        node.set("typeface", "Times New Roman")


def remove_density_shapes(slide) -> None:
    for item in list(slide.shapes):
        if item.name.startswith("DENSITY_") or item.name.startswith("BODY_DENSITY_"):
            item._element.getparent().remove(item._element)


def add_emphasis_run(paragraph, text: str, size: float, *, bold: bool, emphasis: tuple[str, ...]) -> None:
    remaining = text
    while remaining:
        positions = [(remaining.find(term), term) for term in emphasis if remaining.find(term) >= 0]
        if not positions:
            run = paragraph.add_run()
            run.text = remaining
            set_run_font(run, size, bold=bold)
            return
        index, term = min(positions, key=lambda item: item[0])
        if index:
            run = paragraph.add_run()
            run.text = remaining[:index]
            set_run_font(run, size, bold=bold)
        run = paragraph.add_run()
        run.text = term
        set_run_font(run, size, bold=True, color=RED)
        remaining = remaining[index + len(term) :]


def set_bullet_text(shape, lines: list[str], *, emphasis: tuple[str, ...] = ()) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = inch(0.04)
    tf.margin_right = inch(0.04)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(2)
        paragraph.line_spacing = 1.02
        is_main = line.startswith("●")
        size = 19.2 if is_main else 18.4
        add_emphasis_run(paragraph, line, size, bold=is_main, emphasis=emphasis)


def shape(slide, name: str):
    for item in slide.shapes:
        if item.name == name:
            return item
    raise KeyError(name)


def delete_shape(slide, name: str) -> None:
    for item in slide.shapes:
        if item.name == name:
            item._element.getparent().remove(item._element)
            return


def fix_fov_cbf(path: Path) -> None:
    prs = Presentation(str(path))
    slide = prs.slides[5]
    remove_density_shapes(slide)
    add_panel(slide, 0.78, 2.52, 11.78, 1.26, "DENSITY_FORMULA_PANEL")
    set_bounds(shape(slide, "TextBox 6"), 0.92, 2.66, 11.60, 0.48)
    set_bounds(shape(slide, "TextBox 7"), 0.92, 3.22, 11.60, 0.48)
    set_bounds(shape(slide, "Rectangle 8"), 0.76, 4.06, 11.80, 0.02)
    set_bounds(shape(slide, "TextBox 9"), 0.74, 4.30, 5.76, 1.82)
    set_bounds(shape(slide, "TextBox 10"), 6.86, 4.30, 5.62, 1.82)
    prs.save(path)


def fix_primer(path: Path) -> None:
    prs = Presentation(str(path))
    slide = prs.slides[15]
    remove_density_shapes(slide)
    set_bounds(shape(slide, "Picture 6"), 1.18, 2.42, 5.75, 4.06)
    set_bounds(shape(slide, "TextBox 7"), 7.70, 2.52, 4.72, 3.64)
    set_bounds(shape(slide, "Rectangle 8"), 0.76, 6.50, 11.80, 0.02)
    add_panel(slide, 0.86, 6.58, 11.64, 0.50, "DENSITY_TAKEAWAY_PANEL")
    set_bounds(shape(slide, "TextBox 9"), 0.98, 6.68, 11.28, 0.32)
    prs.save(path)


def fix_spot(path: Path) -> None:
    prs = Presentation(str(path))

    slide = prs.slides[3]
    remove_density_shapes(slide)
    set_bounds(shape(slide, "TextBox 5"), 0.76, 1.74, 5.50, 0.32)
    set_bounds(shape(slide, "TextBox 6"), 0.72, 2.08, 5.70, 3.72)
    set_bounds(shape(slide, "Rectangle 7"), 6.56, 1.86, 0.02, 4.50)
    set_bounds(shape(slide, "TextBox 8"), 6.82, 1.74, 5.50, 0.32)
    set_bounds(shape(slide, "TextBox 9"), 6.78, 2.08, 5.72, 3.72)
    add_panel(slide, 0.78, 6.00, 11.78, 0.62, "DENSITY_COMPARE_PANEL")
    add_density_text(
        slide,
        0.98,
        6.12,
        11.36,
        0.48,
        "BODY_DENSITY_COMPARE",
        ["• SPOT的改进点在于把动态障碍的未来占用写入时空走廊，同时保留目标路径不可达时的backup恢复机制。"],
    )

    slide = prs.slides[10]
    remove_density_shapes(slide)
    set_bullet_text(
        shape(slide, "TextBox 5"),
        [
            "● SPOT沿ST‑RRT*路径段构造时空安全走廊。",
            "• 红到灰表示时间演化，绿色线为最终优化轨迹。",
        ],
        emphasis=("SPOT", "ST‑RRT*"),
    )
    set_bounds(shape(slide, "Picture 6"), 1.02, 2.72, 3.90, 4.10)
    set_bounds(shape(slide, "Rectangle 7"), 5.36, 2.74, 0.02, 3.86)
    set_bounds(shape(slide, "TextBox 8"), 5.70, 2.78, 6.56, 2.56)
    set_bullet_text(
        shape(slide, "TextBox 8"),
        [
            "• 动态障碍不是只在当前帧避开，而是在路径段对应的时间窗口内参与走廊构造。",
            "• 每一段凸多面体同时对应空间位置和到达时间，避免只满足几何可通行。",
            "• 绿色轨迹穿过走廊中心，说明优化结果仍保持连续可执行。",
            "• 该设计把空间安全进一步收紧为时空一致的飞行约束。",
        ],
    )

    slide = prs.slides[15]
    remove_density_shapes(slide)
    set_bounds(shape(slide, "TextBox 5"), 0.62, 1.76, 12.00, 1.86)
    set_bullet_text(
        shape(slide, "TextBox 5"),
        [
            "● backup机制对高动态障碍密度场景贡献最明显。",
            "• 10个障碍时，SPOT和无backup版本成功率接近，差距只有5.8个百分点。",
            "• 30个障碍时，无backup版本成功率下降到52.2%，SPOT仍达到80.2%。",
            "• 无backup时UAV可能在目标路径不可达处悬停，更容易与动态障碍碰撞。",
            "• 消融表说明backup主要承担恢复可行路径和触发重规划的作用。",
        ],
        emphasis=("backup", "SPOT", "52.2%", "80.2%"),
    )
    set_bounds(shape(slide, "Rectangle 6"), 0.78, 3.76, 11.78, 0.02)
    metric_cards = [(0.76, 3.96), (4.62, 3.96), (8.48, 3.96)]
    for idx, (left, top) in enumerate(metric_cards, start=1):
        add_panel(slide, left, top, 3.84, 1.08, f"DENSITY_METRIC_PANEL_{idx}")
    for offset, name in enumerate(["TextBox 7", "TextBox 8", "TextBox 9"]):
        set_bounds(shape(slide, name), 0.86, 4.04 + offset * 0.32, 3.74, 0.30)
    for offset, name in enumerate(["TextBox 10", "TextBox 11", "TextBox 12"]):
        set_bounds(shape(slide, name), 4.72, 4.04 + offset * 0.32, 3.74, 0.30)
    for offset, name in enumerate(["TextBox 13", "TextBox 14", "TextBox 15"]):
        set_bounds(shape(slide, name), 8.58, 4.04 + offset * 0.32, 3.74, 0.30)
    note = slide.shapes.add_textbox(inch(0.82), inch(5.74), inch(11.72), inch(0.78))
    note.name = "BODY_DENSITY_NOTE"
    set_bullet_text(
        note,
        [
            "• backup 主要提升目标路径暂时不可达时的恢复能力，密集动态障碍下收益最明显。",
            "• 消融结果说明该模块不是装饰性后处理，而是保证在线规划不断路的关键机制。",
        ],
        emphasis=("backup",),
    )

    slide = prs.slides[20]
    remove_density_shapes(slide)
    set_bounds(shape(slide, "TextBox 5"), 0.58, 1.76, 12.12, 2.08)
    set_bullet_text(
        shape(slide, "TextBox 5"),
        [
            "● 该方法适合作为未知动态环境中的局部重规划与应急避障模块。",
            "• 任务层可给出目标点、优先级和禁飞约束；SPOT 负责局部时空避障。",
            "• 与语义任务结合时，动态障碍类别、速度和预测可信度可以进入规划代价。",
            "• 多机系统中可将其他 UAV 视作动态障碍，但需要补充通信延迟和互相预测机制。",
            "• 迁移到多机系统时，规划器需要同时处理局部避障、队形目标和通信不确定性。",
        ],
        emphasis=("SPOT",),
    )
    delete_shape(slide, "TextBox 6")
    set_bounds(shape(slide, "Rectangle 7"), 0.76, 4.02, 11.80, 0.02)
    set_bounds(shape(slide, "TextBox 8"), 0.82, 4.20, 4.20, 0.32)
    set_bounds(shape(slide, "TextBox 9"), 0.74, 4.56, 11.86, 2.10)
    set_bullet_text(
        shape(slide, "TextBox 9"),
        [
            "• 无地图时空规划减少地图维护成本，但更依赖实时视觉检测和短期预测稳定性。",
            "• 基于地图融合的动态避障信息更完整，但需要承担建图、更新和查询开销。",
            "• 实际系统可按速度、障碍密度和算力预算选择二者的组合边界。",
            "• 多机任务中还需要把通信延迟、互相预测误差和backup优先级纳入统一调度。",
            "• 工程实现还包含短时预测失效时的降级策略，保证局部重规划连续运行。",
        ],
        emphasis=("backup",),
    )

    prs.save(path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Densify public regression PPTX examples after strict rendered blank scans.")
    parser.add_argument("--examples-dir", type=Path, default=EXAMPLES)
    args = parser.parse_args()
    examples = args.examples_dir.resolve()
    targets = {
        "uav-fov-cbf-regression-report.pptx": fix_fov_cbf,
        "uav-primer-regression-report.pptx": fix_primer,
        "uav-spot-regression-report.pptx": fix_spot,
    }
    for filename, fixer in targets.items():
        path = examples / filename
        if not path.exists():
            raise FileNotFoundError(path)
        fixer(path)
        print(f"densified {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
