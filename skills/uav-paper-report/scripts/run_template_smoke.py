from __future__ import annotations

import argparse
import json
import os
import re
import subprocess
import sys
from collections import Counter
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

import run_template_matrix as matrix_runner


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PAPERS = ROOT / "assets" / "template-profiles" / "stress-papers.json"
DEFAULT_TEMPLATES = ROOT / "assets" / "template-profiles" / "template-smoke.local.example.json"
AUDIT = ROOT / "scripts" / "audit_pptx_text.py"
RENDER = ROOT / "scripts" / "render_pptx_previews.py"
SCAN = ROOT / "scripts" / "scan_rendered_slides.py"

EMU_PER_INCH = 914400
FONT = "Times New Roman"
BLACK = RGBColor(25, 25, 25)
GRAY = RGBColor(96, 105, 116)
LIGHT_GRAY = RGBColor(243, 245, 248)
MID_GRAY = RGBColor(210, 215, 222)
RED = RGBColor(205, 32, 44)
WHITE = RGBColor(255, 255, 255)
NODE_FILL = RGBColor(238, 246, 255)
DIAGRAM_TITLE_SIZE = 14.2
DIAGRAM_LABEL_SIZE = 14.0

PROFILE_SIZE = {
    "compact": {
        "main": 18.2,
        "secondary": 16.8,
        "tertiary": 16.2,
        "formula": 19.0,
        "table": 11.8,
        "metric_value": 15.6,
        "metric_label": 11.8,
    },
    "dense-visual": {
        "main": 18.2,
        "secondary": 16.8,
        "tertiary": 16.2,
        "formula": 19.0,
        "table": 11.8,
        "metric_value": 15.6,
        "metric_label": 11.8,
    },
    "classic-large": {
        "main": 19.2,
        "secondary": 18.4,
        "tertiary": 17.0,
        "formula": 20.0,
        "table": 12.4,
        "metric_value": 16.4,
        "metric_label": 12.2,
    },
}

THEMES = {
    "blue": RGBColor(0, 121, 192),
    "cyan": RGBColor(0, 145, 166),
    "green": RGBColor(35, 145, 95),
    "red": RGBColor(172, 38, 47),
    "purple": RGBColor(92, 78, 152),
    "dark": RGBColor(25, 41, 58),
}


def expand_path(raw: str, *, template_root: Path | None) -> Path:
    text = os.path.expandvars(raw)
    if template_root is not None:
        text = text.replace("${PPTAGENT_TEMPLATE_ROOT}", str(template_root))
    return Path(text).expanduser().resolve()


def rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return color[0], color[1], color[2]


def brightness(color: RGBColor) -> float:
    r, g, b = rgb_tuple(color)
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255.0


def theme_color(template: Path, theme: str | None) -> RGBColor:
    configured = THEMES.get((theme or "").lower())
    if configured is not None:
        return configured

    colors: Counter[tuple[int, int, int]] = Counter()
    try:
        prs = Presentation(str(template))
        for slide in list(prs.slides)[:4]:
            for shape in slide.shapes:
                try:
                    fill = shape.fill
                    if fill.type != 1:
                        continue
                    rgb = fill.fore_color.rgb
                    if rgb is None:
                        continue
                    r, g, b = rgb_tuple(rgb)
                    if max(r, g, b) > 238 or max(r, g, b) < 35:
                        continue
                    if abs(r - g) + abs(g - b) + abs(r - b) < 42:
                        continue
                    colors[(r, g, b)] += 1
                except Exception:
                    continue
    except Exception:
        pass
    if not colors:
        return THEMES["blue"]
    r, g, b = colors.most_common(1)[0][0]
    return RGBColor(r, g, b)


def clear_slides(prs: Presentation) -> None:
    sld_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public slide-delete API.
    for sld_id in list(sld_id_list):
        prs.part.drop_rel(sld_id.rId)
        sld_id_list.remove(sld_id)


def set_typeface(rpr, font: str = FONT) -> None:
    for tag in ("latin", "ea", "cs"):
        node = rpr.find(qn(f"a:{tag}"))
        if node is None:
            node = rpr.makeelement(qn(f"a:{tag}"))
            rpr.append(node)
        node.set("typeface", font)


def set_run(run, size: float, *, bold: bool = False, color: RGBColor = BLACK) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    set_typeface(run._r.get_or_add_rPr(), FONT)  # noqa: SLF001


def add_textbox(slide, left: float, top: float, width: float, height: float, *, name: str):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    shape.name = name
    shape.text_frame.margin_left = Inches(0.03)
    shape.text_frame.margin_right = Inches(0.03)
    shape.text_frame.margin_top = Inches(0.02)
    shape.text_frame.margin_bottom = Inches(0.02)
    shape.text_frame.word_wrap = True
    return shape


def compact_text_terms(text: str) -> str:
    replacements = {
        "time-optimal primitive selection": "时间最优基元选择",
        "Time-optimal primitive selection": "时间最优基元选择",
        "trajectory candidates": "候选轨迹",
        "safe primitive selection": "安全基元筛选",
        "runtime safety assurance flow": "运行时安全保障流程",
        "method pipeline and closed-loop execution": "方法流程与闭环执行",
        "RRT* path": "RRT*路径",
        "LQR tracking": "LQR跟踪",
        "safety filter": "安全滤波",
        "safe input": "安全输入",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    return text


def clean_inline_text(text: str, *, context: str) -> str:
    if "\n" in text or "\r" in text:
        raise ValueError(f"manual newline is not allowed in {context}")
    cleaned = " ".join(compact_text_terms(text).split())
    if not cleaned:
        raise ValueError(f"empty text is not allowed in {context}")
    return cleaned


def clean_formula_text(text: str) -> str:
    if "\n" in text or "\r" in text:
        raise ValueError("manual newline is not allowed in formula runs")
    cleaned = re.sub(r"[ \t]+", " ", compact_text_terms(text))
    if not cleaned.strip():
        raise ValueError("empty formula run is not allowed")
    return cleaned


def add_plain(slide, text: str, left: float, top: float, width: float, height: float, size: float, *,
              name: str, bold: bool = False, color: RGBColor = BLACK, align=PP_ALIGN.LEFT) -> None:
    shape = add_textbox(slide, left, top, width, height, name=name)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = 0.95
    run = p.add_run()
    run.text = clean_inline_text(text, context=name)
    set_run(run, size, bold=bold, color=color)


def set_baseline(run, value: int) -> None:
    run._r.get_or_add_rPr().set("baseline", str(value))  # noqa: SLF001


def add_math_run(paragraph, piece: str | tuple[str, str], size: float) -> None:
    if isinstance(piece, tuple):
        text, mode = piece
    else:
        text, mode = piece, ""
    text = clean_formula_text(text)
    run = paragraph.add_run()
    run.text = text
    run_size = size * 0.72 if mode in {"sub", "sup"} else size
    set_run(run, run_size, color=BLACK)
    if mode == "sub":
        set_baseline(run, -25000)
    elif mode == "sup":
        set_baseline(run, 30000)


def add_bullets(slide, lines: list[str], left: float, top: float, width: float, height: float, sizes: dict[str, float], *,
                name: str, space_after: float = 1.2) -> None:
    emphasis_terms = (
        "安全",
        "时间最优",
        "运行时",
        "CBF",
        "LQR",
        "执行器",
        "多模态",
        "最小间距",
        "覆盖率",
        "平滑",
        "视觉惯性",
        "事件相机",
        "体素",
        "通信",
        "多智能体",
        "优先级目标",
        "deterministic",
    )

    def add_emphasis_runs(paragraph, text: str, size: float, *, base_bold: bool) -> None:
        cursor = 0
        while cursor < len(text):
            hit = None
            hit_pos = len(text)
            for term in emphasis_terms:
                pos = text.find(term, cursor)
                if pos != -1 and pos < hit_pos:
                    hit = term
                    hit_pos = pos
            if hit is None:
                run = paragraph.add_run()
                run.text = text[cursor:]
                set_run(run, size, bold=base_bold, color=BLACK)
                break
            if hit_pos > cursor:
                run = paragraph.add_run()
                run.text = text[cursor:hit_pos]
                set_run(run, size, bold=base_bold, color=BLACK)
            run = paragraph.add_run()
            run.text = hit
            set_run(run, size, bold=True, color=RED)
            cursor = hit_pos + len(hit)

    shape = add_textbox(slide, left, top, width, height, name=name)
    tf = shape.text_frame
    tf.clear()
    if not lines:
        raise ValueError(f"empty bullet block is not allowed in {name}")
    for idx, raw in enumerate(lines):
        line = clean_inline_text(raw, context=name)
        marker = line.lstrip()[0]
        if marker == "●":
            size = sizes["main"]
            indent = Inches(0.28)
            first = Inches(-0.18)
        elif marker == "•":
            size = sizes["secondary"]
            indent = Inches(0.43)
            first = Inches(-0.16)
        else:
            size = sizes["tertiary"]
            indent = Inches(0.60)
            first = Inches(-0.14)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 0.93
        p.left_margin = indent
        p.first_line_indent = first
        bold = marker == "●" and len(line) < 34
        add_emphasis_runs(p, line, size, base_bold=bold)


def strip_empty_text_bodies(prs: Presentation) -> int:
    removed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            element = shape._element  # noqa: SLF001
            tx_body = element.find(qn("p:txBody"))
            if tx_body is None:
                continue
            texts = [node.text or "" for node in tx_body.findall(".//" + qn("a:t"))]
            if not any(text.strip() for text in texts):
                element.remove(tx_body)
                removed += 1
    return removed


def add_rect(slide, left: float, top: float, width: float, height: float, color: RGBColor, *, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_rule(slide, left: float, top: float, width: float, accent: RGBColor) -> None:
    shape = add_rect(slide, left, top, width, 0.018, accent)
    shape.name = "RULE"


def add_header(slide, w: float, h: float, accent: RGBColor, part: str, title: str) -> None:
    # Some public templates keep sample lab names, footers, or date fields on the
    # master. Cover the master canvas before adding our own report layout.
    add_rect(slide, 0.0, 0.0, w, h, WHITE)
    add_rect(slide, 0.0, 0.0, w, 0.48, LIGHT_GRAY)
    add_rect(slide, 0.0, 0.48, w, 0.04, accent)
    add_plain(slide, part, 0.46, 0.12, 1.30, 0.24, 11.0, name="HEADER_PART", bold=True, color=accent)
    add_plain(slide, title, 1.55, 0.12, max(2.0, w - 2.1), 0.24, 11.0, name="HEADER_TITLE", color=GRAY)


def add_section(slide, w: float, y: float, number: str, title: str, accent: RGBColor) -> None:
    add_plain(slide, number, 0.46, y, 0.78, 0.36, 18.5, name="SECTION_NO", bold=True, color=accent)
    add_plain(slide, title, 1.24, y + 0.03, max(2.5, w - 1.74), 0.34, 18.5, name="SECTION_TITLE", bold=True, color=BLACK)


def fit_picture(slide, path: Path, left: float, top: float, width: float, height: float, *, name: str):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = width / height
    if ratio > box_ratio:
        draw_w = width
        draw_h = width / ratio
        x = left
        y = top + (height - draw_h) / 2
    else:
        draw_h = height
        draw_w = height * ratio
        x = left + (width - draw_w) / 2
        y = top
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(draw_w), height=Inches(draw_h))
    pic.name = name
    return pic


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor, *, width_pt: float = 2.2, name: str = "DIAG_LINE") -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.name = name
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)


def add_diagram_node(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str,
    accent: RGBColor,
    *,
    name: str,
    size: float = 12.8,
) -> None:
    node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    node.name = f"{name}_BOX"
    node.fill.solid()
    node.fill.fore_color.rgb = NODE_FILL
    node.line.color.rgb = accent
    node.line.width = Pt(1.5)
    tf = node.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = clean_inline_text(label, context=name)
    set_run(run, size, color=BLACK)


def pipeline_diagram(slide, left: float, top: float, width: float, height: float, accent: RGBColor, *, name: str) -> None:
    add_rect(slide, left, top, width, height, RGBColor(250, 252, 255))
    title = add_textbox(slide, left + 0.18, top + 0.12, width - 0.36, 0.30, name=f"{name}_TITLE")
    p = title.text_frame.paragraphs[0]
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = "方法流程与闭环执行"
    set_run(run, DIAGRAM_TITLE_SIZE, bold=True, color=GRAY)

    node_w = width * 0.24
    node_h = height * 0.17
    row1 = top + height * 0.28
    row2 = top + height * 0.66
    xs = [left + width * 0.08, left + width * 0.38, left + width * 0.68]
    bx = [left + width * 0.24, left + width * 0.54]
    add_line(slide, xs[0] + node_w, row1 + node_h / 2, xs[1], row1 + node_h / 2, accent, name=f"{name}_LINE")
    add_line(slide, xs[1] + node_w, row1 + node_h / 2, xs[2], row1 + node_h / 2, accent, name=f"{name}_LINE")
    add_line(slide, xs[1] + node_w / 2, row1 + node_h, bx[0] + node_w / 2, row2, accent, name=f"{name}_LINE")
    add_line(slide, xs[2] + node_w / 2, row1 + node_h, bx[1] + node_w / 2, row2, accent, name=f"{name}_LINE")
    add_line(slide, bx[0] + node_w, row2 + node_h / 2, bx[1], row2 + node_h / 2, accent, name=f"{name}_LINE")
    for idx, (x, y, label) in enumerate([
        (xs[0], row1, "基元库"),
        (xs[1], row1, "碰撞检测"),
        (xs[2], row1, "代价选择"),
        (bx[0], row2, "轨迹"),
        (bx[1], row2, "控制器"),
    ]):
        add_diagram_node(slide, x, y, node_w, node_h, label, accent, name=f"{name}_LABEL_{idx}", size=DIAGRAM_LABEL_SIZE)


def runtime_flow_diagram(slide, left: float, top: float, width: float, height: float, accent: RGBColor, *, name: str) -> None:
    add_rect(slide, left, top, width, height, RGBColor(250, 252, 255))
    title = add_textbox(slide, left + 0.16, top + 0.12, width - 0.32, 0.30, name=f"{name}_TITLE")
    p = title.text_frame.paragraphs[0]
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = "运行时安全保障流程"
    set_run(run, DIAGRAM_TITLE_SIZE, bold=True, color=GRAY)

    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + width * 0.13), Inches(top + height * 0.10), Inches(width * 0.74), Inches(height * 0.82))
    frame.name = f"{name}_FRAME"
    frame.fill.background()
    frame.line.color.rgb = MID_GRAY
    frame.line.width = Pt(1.4)

    node_w = width * 0.51
    node_h = height * 0.09
    node_x = left + width * 0.30
    dot_x = left + width * 0.20
    y_positions = [top + height * 0.20, top + height * 0.39, top + height * 0.58, top + height * 0.77]
    labels = ["RRT*路径", "LQR跟踪", "安全滤波", "安全输入"]
    for idx, (y, label) in enumerate(zip(y_positions, labels)):
        color = accent if idx < 3 else RED
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(dot_x - 0.055), Inches(y + node_h * 0.18), Inches(0.11), Inches(0.11))
        dot.name = f"{name}_DOT_{idx}"
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        if idx < len(y_positions) - 1:
            add_line(slide, dot_x, y + node_h * 0.30, dot_x, y_positions[idx + 1] + node_h * 0.20, MID_GRAY, width_pt=1.5, name=f"{name}_LINE")
        add_diagram_node(slide, node_x, y, node_w, node_h, label, color, name=f"{name}_LABEL_{idx}", size=DIAGRAM_LABEL_SIZE)


def draw_assets(out_dir: Path, accent: RGBColor) -> dict[str, Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    red_ch, green_ch, blue_ch = rgb_tuple(accent)
    accent_rgb = (red_ch, green_ch, blue_ch)
    muted = (min(255, red_ch + 62), min(255, green_ch + 62), min(255, blue_ch + 62))
    paths: dict[str, Path] = {}

    wide = Image.new("RGB", (1100, 430), (250, 252, 255))
    d = ImageDraw.Draw(wide)
    d.rectangle((55, 52, 1045, 364), outline=(190, 198, 210), width=3)
    for x in range(120, 1000, 125):
        d.line((x, 65, x, 352), fill=(226, 230, 236), width=1)
    for y in range(100, 340, 58):
        d.line((70, y, 1030, y), fill=(226, 230, 236), width=1)
    pts = [(90, 330), (210, 270), (342, 285), (480, 195), (620, 170), (760, 116), (918, 126), (1010, 84)]
    d.line(pts, fill=accent_rgb, width=7, joint="curve")
    for x, y in pts[1:-1]:
        d.ellipse((x - 10, y - 10, x + 10, y + 10), fill=muted, outline=accent_rgb, width=2)
    for rect in [(270, 92, 370, 170), (694, 232, 840, 314)]:
        d.rounded_rectangle(rect, radius=12, fill=(255, 228, 230), outline=(205, 32, 44), width=2)
    path = out_dir / "wide-trajectory.png"
    wide.save(path)
    paths["wide"] = path

    return paths


def add_table(
    slide,
    rows: list[list[str]],
    left: float,
    top: float,
    width: float,
    height: float,
    size: float,
    accent: RGBColor,
    *,
    col_widths: list[float] | None = None,
) -> None:
    graphic = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height))
    graphic.name = "TABLE_RESULTS"
    tbl = graphic.table
    if col_widths:
        if len(col_widths) != len(rows[0]):
            raise ValueError("col_widths must match table column count")
        total = sum(col_widths)
        if total <= 0:
            raise ValueError("col_widths must contain positive widths")
        for ci, raw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(width * raw / total)
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            cell = tbl.cell(ri, ci)
            value = clean_inline_text(value, context="table cell")
            cell.text = value
            cell.margin_left = Inches(0.02)
            cell.margin_right = Inches(0.02)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if ri == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(0)
            for run in p.runs:
                value_is_key = value in {"低", "运行时输入滤波", "时间最优候选", "实机实验"}
                set_run(run, size, bold=ri == 0 or ci == 0 or value_is_key, color=RED if value_is_key else BLACK)


FORMULA_SPECS: dict[str, list[list[str | tuple[str, str]]]] = {
    "primitive-planner-2025": [
        ["p(t) = p", ("0", "sub"), " + v", ("0", "sub"), "t + 1/2 a t", ("2", "sup")],
        ["J(π) = T(π) + λ", ("c", "sub"), " C(π) + λ", ("s", "sub"), " S(π)"],
        ["π", ("*", "sup"), " = arg min", ("π∈Pₛ", "sub"), " J(π)"],
    ],
    "safety-assurance-2025": [
        ["u", ("s", "sub"), " = arg min", ("u", "sub"), " ||u - u", ("n", "sub"), "||", ("2", "sup")],
        ["h(x) ≥ 0,    ḣ(x) + α h(x) ≥ 0"],
        ["ẋ = f(x) + g(x)u", ("s", "sub")],
    ],
    "quad-lcd-2025": [
        ["u = clip(π", ("θ", "sub"), "(x,r), u", ("min", "sub"), ", u", ("max", "sub"), ")"],
        ["L = L", ("track", "sub"), " + λ", ("u", "sub"), "L", ("sat", "sub"), " + λ", ("s", "sub"), "L", ("smooth", "sub")],
        ["e", ("t", "sub"), " = ||p", ("t", "sub"), " - p", ("ref", "sub"), "||", ("2", "sub"), ",    Δu", ("t", "sub"), " = ||u", ("t", "sub"), " - u", ("t-1", "sub"), "||", ("2", "sub")],
    ],
    "mc-swarm-2025": [
        ["m", ("i", "sub"), " = arg max", ("m", "sub"), " q", ("m", "sub"), "(x", ("i", "sub"), ", X", ("N", "sub"), ", g)"],
        ["d", ("ij", "sub"), " = ||p", ("i", "sub"), " - p", ("j", "sub"), "||", ("2", "sub"), " ≥ d", ("min", "sub")],
        ["u", ("i", "sub"), " = u", ("i", "sub"), ("m", "sup"), " + k", ("s", "sub"), " ∇h", ("i", "sub"), "(x)"],
    ],
    "smooth-coverage-2025": [
        ["min ", "Σ", ("i", "sub"), " (L", ("i", "sub"), " + λ", ("k", "sub"), "K", ("i", "sub"), " + λ", ("b", "sub"), "B", ("i", "sub"), ")"],
        ["s.t.  ∪", ("i", "sub"), " C", ("i", "sub"), " = Ω,    C", ("i", "sub"), " ∩ C", ("j", "sub"), " = ∅"],
        ["||p", ("i", "sub"), "(t) - p", ("j", "sub"), "(t)||", ("2", "sub"), " ≥ d", ("min", "sub")],
    ],
    "downfacing-vio-2025": [
        ["T", ("k+1", "sub"), " = T", ("k", "sub"), " exp(ξ", ("k", "sub"), " Δt)"],
        ["r", ("i", "sub"), " = z", ("i", "sub"), " - π(T", ("cw", "sub"), " P", ("i", "sub"), ")"],
        ["E = Σ", ("i", "sub"), " ||r", ("i", "sub"), "||", ("2", "sup"), " + λ", ("b", "sub"), " ||b||", ("2", "sup")],
    ],
    "voxel-esvio-2025": [
        ["p", ("v", "sub"), " = TopK(p | v, s(p))"],
        ["E", ("map", "sub"), " = Σ", ("v", "sub"), " w", ("v", "sub"), " ||e", ("v", "sub"), "||", ("2", "sup")],
        ["x", ("*", "sup"), " = arg min", ("x", "sub"), " (E", ("imu", "sub"), " + E", ("event", "sub"), " + E", ("stereo", "sub"), ")"],
    ],
    "persistent-monitoring-2025": [
        ["a", ("i", "sub"), ("*", "sup"), " = arg max", ("a", "sub"), " Q", ("i", "sub"), "(o", ("i", "sub"), ", m", ("i", "sub"), ", a)"],
        ["R", ("ij", "sub"), "(t) ≥ R", ("min", "sub"), ",    ||p", ("i", "sub"), " - p", ("j", "sub"), "|| ≤ d", ("comm", "sub")],
        ["J = Σ", ("l", "sub"), " γ", ("l", "sub"), " q", ("l", "sub"), "(t) - λ", ("c", "sub"), " C", ("i", "sub"), "(t)"],
    ],
    "terrain-aware-uav-2025": [
        ["z", ("k", "sub"), " = h(x", ("k", "sub"), ", y", ("k", "sub"), ") + d", ("safe", "sub")],
        ["v", ("k", "sub"), " = Π(K, R", ("k", "sub"), ", t", ("k", "sub"), ", X)"],
        ["J = Σ", ("k", "sub"), " (1 - cov(v", ("k", "sub"), ")) + λ", ("h", "sub"), "|Δz", ("k", "sub"), "| + λ", ("ψ", "sub"), "|Δψ", ("k", "sub"), "|"],
    ],
    "zoom-ptz-coverage-2025": [
        ["GSD", ("k", "sub"), " = H", ("k", "sub"), " p / f", ("k", "sub")],
        ["f", ("k", "sub"), ("*", "sup"), " = clip(H", ("k", "sub"), " p / GSD", ("0", "sub"), ", f", ("min", "sub"), ", f", ("max", "sub"), ")"],
        ["J = T + λ", ("L", "sub"), " L + λ", ("z", "sub"), " Σ", ("k", "sub"), " |Δz", ("k", "sub"), "|"],
    ],
    "diff-physics-flight-2025": [
        ["a", ("t", "sub"), " = π", ("θ", "sub"), "(D", ("t", "sub"), ", v", ("t", "sub"), ")"],
        ["v", ("t+1", "sub"), " = v", ("t", "sub"), " + a", ("t", "sub"), " Δt"],
        ["L = Σ", ("t", "sub"), " ℓ", ("obs", "sub"), "(p", ("t", "sub"), ") + λ", ("u", "sub"), " ||a", ("t", "sub"), "||", ("2", "sup")],
    ],
    "shape-adaptive-quad-2025": [
        ["x", ("k+1", "sub"), " = f(x", ("k", "sub"), ", u", ("k", "sub"), ", η", ("k", "sub"), ")"],
        ["J = Σ", ("k", "sub"), " ||p", ("k", "sub"), " - p", ("ref", "sub"), "||", ("2", "sup"), " + λ", ("η", "sub"), " ||Δη", ("k", "sub"), "||", ("2", "sup")],
        ["η", ("*", "sup"), " = arg min", ("η", "sub"), " J", ("path", "sub"), "(x, η)"],
    ],
}


def formula_rows_for(paper: dict[str, Any]) -> list[list[str | tuple[str, str]]]:
    return FORMULA_SPECS.get(paper["id"], [[row] for row in paper["formula_rows"]])


def add_formula_rows(slide, rows: list[list[str | tuple[str, str]]], left: float, top: float, width: float, row_h: float, size: float, accent: RGBColor) -> None:
    for idx, row in enumerate(rows):
        y = top + idx * row_h
        add_plain(slide, f"({idx + 1})", left, y + 0.04, 0.50, 0.32, size - 2.2, name=f"MATH_LABEL_{idx}", color=GRAY, align=PP_ALIGN.RIGHT)
        shape = add_textbox(slide, left + 0.68, y, width - 0.68, 0.46, name=f"MATH_BODY_{idx}")
        shape.text_frame.word_wrap = False
        shape.text_frame.margin_left = Inches(0)
        shape.text_frame.margin_right = Inches(0)
        shape.text_frame.margin_top = Inches(0)
        shape.text_frame.margin_bottom = Inches(0)
        p = shape.text_frame.paragraphs[0]
        p.space_after = Pt(0)
        p.line_spacing = 1.0
        for piece in row:
            add_math_run(p, piece, size)
    add_rule(slide, left + 0.62, top + len(rows) * row_h + 0.08, width - 0.62, accent)


def metric_row(
    slide,
    items: list[list[str]],
    left: float,
    top: float,
    width: float,
    accent: RGBColor,
    size: float,
    *,
    label_size: float | None = None,
) -> None:
    gap = 0.12
    box_w = (width - gap * (len(items) - 1)) / len(items)
    label_pt = label_size if label_size is not None else max(11.6, size - 3.2)
    for idx, item in enumerate(items):
        x = left + idx * (box_w + gap)
        box = add_rect(slide, x, top, box_w, 0.64, RGBColor(248, 250, 253), line=MID_GRAY)
        box.name = f"METRIC_BOX_{idx}"
        add_plain(slide, item[0], x + 0.08, top + 0.08, box_w - 0.16, 0.22, size, name=f"METRIC_VALUE_{idx}", bold=True, color=accent, align=PP_ALIGN.CENTER)
        add_plain(slide, item[1], x + 0.08, top + 0.36, box_w - 0.16, 0.18, label_pt, name=f"METRIC_LABEL_{idx}", color=GRAY, align=PP_ALIGN.CENTER)


def blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower() or "空白" in layout.name:
            return layout
    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


def scrub_placeholders(slide) -> None:
    banned = ("单击此处", "Click to", "‹#›")
    for shape in list(slide.shapes):
        if not getattr(shape, "is_placeholder", False):
            continue
        text = getattr(shape, "text", "")
        if not text or any(token.lower() in text.lower() for token in banned):
            slide.shapes._spTree.remove(shape._element)  # noqa: SLF001


def add_blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(blank_layout(prs))
    scrub_placeholders(slide)
    return slide


def cover_slide(prs: Presentation, paper: dict[str, Any], accent: RGBColor, template_id: str, sizes: dict[str, float]) -> None:
    w = prs.slide_width / EMU_PER_INCH
    h = prs.slide_height / EMU_PER_INCH
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, w, h, RGBColor(246, 249, 252))
    band_h = h * 0.40
    band_y = h * 0.27
    add_rect(slide, 0, band_y, w, band_h, accent)
    title_color = WHITE if brightness(accent) < 0.72 else BLACK
    add_plain(slide, paper["title"], w * 0.08, band_y + band_h * 0.24, w * 0.84, band_h * 0.34, min(28, sizes["main"] + 8.8), name="COVER_TITLE", bold=True, color=title_color, align=PP_ALIGN.CENTER)
    add_plain(slide, f"{paper['topic']} / {paper['year']}", w * 0.10, band_y + band_h * 0.62, w * 0.80, 0.34, sizes["secondary"], name="COVER_SUBTITLE", bold=True, color=title_color, align=PP_ALIGN.CENTER)
    add_plain(slide, "中文学术论文汇报", w * 0.32, h - 0.72, w * 0.36, 0.28, 13.0, name="COVER_FOOTER", color=accent, align=PP_ALIGN.CENTER)


def content_position(prs: Presentation, paper: dict[str, Any], accent: RGBColor, sizes: dict[str, float]) -> None:
    w = prs.slide_width / EMU_PER_INCH
    h = prs.slide_height / EMU_PER_INCH
    slide = add_blank_slide(prs)
    add_header(slide, w, h, accent, "Part. 01", "论文定位")
    add_section(slide, w, 0.78, "1.1", "研究问题与核心思路", accent)
    body_top = 1.35
    if w < 11.0:
        add_bullets(slide, ["● " + paper["problem"], "● " + paper["main_claim"]] + paper["bullets"][:2], 0.55, body_top, w - 1.1, 2.58, sizes, name="BODY_POSITION", space_after=0.78)
        add_table(slide, paper["comparison_rows_narrow"], 0.68, 3.92, w - 1.36, 1.08, sizes["table"], accent)
        add_bullets(slide, paper["position_notes"], 0.72, 5.18, w - 1.44, 0.64, sizes, name="BODY_POSITION_NOTE", space_after=0.0)
        metric_row(slide, paper["metrics"], 0.65, h - 1.16, w - 1.30, accent, sizes["metric_value"], label_size=sizes["metric_label"])
    else:
        compact_profile = sizes["main"] < 19.0
        body_h = 1.94 if compact_profile else 2.05
        rule_y = 3.34 if compact_profile else 3.56
        table_y = rule_y + 0.24
        metric_y = table_y + 1.22
        note_y = metric_y + 0.82
        # Keep this stress slide strictly top-to-bottom. The older side-rail
        # version created narrow paragraphs and made wrapped lines look like
        # manual blank lines after LibreOffice export.
        add_bullets(
            slide,
            ["● " + paper["problem"], "● " + paper["main_claim"]] + paper["bullets"][:2],
            0.78,
            body_top,
            w - 1.56,
            body_h,
            sizes,
            name="BODY_POSITION",
            space_after=0.24,
        )
        add_rule(slide, 0.78, rule_y, w - 1.56, accent)
        add_table(slide, paper["comparison_rows_wide"], 0.92, table_y, w - 1.84, 1.10, sizes["table"], accent)
        metric_row(slide, paper["metrics"], 0.92, metric_y, w - 1.84, accent, sizes["metric_value"], label_size=sizes["metric_label"])
        add_bullets(
            slide,
            paper["position_notes"],
            0.94,
            note_y,
            w - 1.88,
            1.02,
            sizes,
            name="BODY_POSITION_NOTE",
            space_after=0.10,
        )


def method_overview(prs: Presentation, paper: dict[str, Any], assets: dict[str, Path], accent: RGBColor, sizes: dict[str, float]) -> None:
    w = prs.slide_width / EMU_PER_INCH
    h = prs.slide_height / EMU_PER_INCH
    slide = add_blank_slide(prs)
    add_header(slide, w, h, accent, "Part. 02", "研究方法")
    add_section(slide, w, 0.78, "2.1", "方法流程", accent)
    if w < 11.0:
        pipeline_diagram(slide, 0.72, 1.34, 4.22, 3.12, accent, name="DIAG_PIPELINE")
        add_bullets(slide, paper["bullets"][:2], 5.12, 1.34, w - 5.66, 1.82, sizes, name="BODY_METHOD_RIGHT", space_after=0.75)
        add_table(
            slide,
            [["阶段", "输出"], ["候选库", "基元"], ["碰撞检测", "安全集"], ["代价选择", "轨迹"]],
            5.18,
            3.44,
            w - 5.80,
            0.94,
            sizes["table"],
            accent,
        )
        add_rule(slide, 0.68, 4.72, w - 1.36, accent)
        add_bullets(slide, paper["method_notes_narrow"], 0.72, 4.92, w - 1.44, 0.78, sizes, name="BODY_METHOD_BOTTOM", space_after=0.35)
        metric_row(slide, [["流程", "候选生成"], ["约束", "安全筛选"], ["输出", "连续轨迹"]], 0.82, 6.04, w - 1.64, accent, sizes["metric_value"], label_size=sizes["metric_label"])
    else:
        add_bullets(
            slide,
            paper["bullets"][:2] + paper["method_notes_wide"],
            0.82,
            1.36,
            w - 1.64,
            1.18,
            sizes,
            name="BODY_METHOD_TOP",
            space_after=0.30,
        )
        pipeline_diagram(slide, 0.88, 2.48, w * 0.38, 3.58, accent, name="DIAG_PIPELINE")
        add_table(
            slide,
            paper["method_rows"],
            w * 0.52,
            2.54,
            w * 0.40,
            1.06,
            sizes["table"],
            accent,
        )
        add_rule(slide, w * 0.52, 3.90, w * 0.40, accent)
        add_bullets(slide, paper["method_bottom_notes"], w * 0.52, 4.14, w * 0.40, 0.86, sizes, name="BODY_METHOD_NOTE", space_after=0.0)
        metric_row(slide, [["流程", "候选生成"], ["约束", "安全筛选"], ["输出", "连续轨迹"]], w * 0.52, 5.42, w * 0.40, accent, sizes["metric_value"], label_size=sizes["metric_label"])


def formula_slide(prs: Presentation, paper: dict[str, Any], accent: RGBColor, sizes: dict[str, float]) -> None:
    w = prs.slide_width / EMU_PER_INCH
    h = prs.slide_height / EMU_PER_INCH
    slide = add_blank_slide(prs)
    add_header(slide, w, h, accent, "Part. 02", "研究方法")
    add_section(slide, w, 0.78, "2.2", "约束与目标函数", accent)
    formula_rows = formula_rows_for(paper)
    if w < 11.0:
        add_formula_rows(slide, formula_rows, 0.74, 1.38, w - 1.45, 0.58, sizes["formula"], accent)
        rows = [["符号", "含义"]] + paper["terms"][:4]
        add_table(slide, rows, 0.70, 3.35, w - 1.40, 1.92, sizes["table"], accent, col_widths=[0.28, 0.72])
        add_bullets(slide, paper["formula_notes_narrow"] + paper["formula_notes_wide"][1:2], 0.72, 5.48, w - 1.44, 0.88, sizes, name="BODY_FORMULA_NOTE", space_after=0.15)
    else:
        add_formula_rows(slide, formula_rows, 0.86, 1.44, w * 0.54, 0.62, sizes["formula"], accent)
        rows = [["符号", "含义"]] + paper["terms"][:4]
        add_table(slide, rows, w * 0.61, 1.42, w * 0.33, 2.50, sizes["table"], accent, col_widths=[0.28, 0.72])
        add_rule(slide, 0.86, 4.10, w - 1.72, accent)
        add_bullets(slide, paper["formula_notes_wide"], 0.86, 4.32, w * 0.52, 1.36, sizes, name="BODY_FORMULA_EXPLAIN", space_after=0.16)
        add_table(slide, [["对象", "进入方式"], ["目标", "代价函数"], ["约束", "可行集合"], ["输出", "控制/轨迹"]], w * 0.61, 4.28, w * 0.33, 1.30, sizes["table"], accent, col_widths=[0.38, 0.62])
        metric_row(slide, [["目标", "优化量"], ["约束", "安全集"], ["动力学", "状态传播"]], 0.98, 5.86, w - 1.96, accent, sizes["metric_value"], label_size=sizes["metric_label"])


def evidence_slide(prs: Presentation, paper: dict[str, Any], assets: dict[str, Path], accent: RGBColor, sizes: dict[str, float]) -> None:
    w = prs.slide_width / EMU_PER_INCH
    h = prs.slide_height / EMU_PER_INCH
    slide = add_blank_slide(prs)
    add_header(slide, w, h, accent, "Part. 03", "实验与结果")
    add_section(slide, w, 0.78, "3.1", "宽图与表格排版", accent)
    rows = [
        ["维度", "基线", "本文方法"],
        ["在线负担", "中等", paper["metrics"][0][1]],
        ["轨迹性质", "可行候选", paper["metrics"][1][1]],
        ["验证范围", "仿真", paper["metrics"][2][1]],
    ]
    if w < 11.0:
        fit_picture(slide, assets["wide"], 0.72, 1.34, w - 1.44, 2.50, name="FIG_TRAJECTORY")
        add_table(slide, rows, 0.70, 4.10, w - 1.40, 1.40, sizes["table"], accent)
        add_bullets(slide, paper["evidence_notes_narrow"], 0.70, 5.74, w - 1.40, 0.72, sizes, name="BODY_EVIDENCE", space_after=0.0)
    else:
        fit_picture(slide, assets["wide"], 0.74, 1.34, w * 0.53, 2.76, name="FIG_TRAJECTORY")
        add_bullets(slide, paper["evidence_notes_wide"], w * 0.61, 1.44, w * 0.33, 2.36, sizes, name="BODY_EVIDENCE", space_after=0.28)
        add_rule(slide, 0.74, 4.30, w - 1.48, accent)
        add_table(slide, rows, 0.80, 4.54, w - 1.60, 1.22, sizes["table"], accent)
        metric_row(slide, [["轨迹", "候选收敛"], ["对比", "性能差异"], ["结论", "可执行性"]], 0.94, 6.04, w - 1.88, accent, sizes["metric_value"], label_size=sizes["metric_label"])


def narrow_figure_slide(prs: Presentation, paper: dict[str, Any], assets: dict[str, Path], accent: RGBColor, sizes: dict[str, float]) -> None:
    w = prs.slide_width / EMU_PER_INCH
    h = prs.slide_height / EMU_PER_INCH
    slide = add_blank_slide(prs)
    add_header(slide, w, h, accent, "Part. 03", "实验与结果")
    add_section(slide, w, 0.78, "3.2", "纵向流程与术语解释", accent)
    if w < 11.0:
        runtime_flow_diagram(slide, 0.70, 1.42, 3.35, 4.85, accent, name="DIAG_TALL_FLOW")
        add_bullets(slide, paper["tall_figure_notes_narrow"], 4.32, 1.48, w - 4.85, 3.48, sizes, name="BODY_NARROW", space_after=0.8)
        add_table(slide, [["变量", "含义"]] + paper["terms"][:2], 4.34, 5.18, w - 4.90, 0.92, sizes["table"], accent, col_widths=[0.30, 0.70])
    else:
        runtime_flow_diagram(slide, 0.88, 1.40, w * 0.31, 5.42, accent, name="DIAG_TALL_FLOW")
        add_bullets(slide, paper["tall_figure_notes_wide"], w * 0.42, 1.44, w * 0.50, 2.02, sizes, name="BODY_NARROW", space_after=0.16)
        add_table(slide, [["变量", "含义"]] + paper["terms"][:3], w * 0.42, 3.76, w * 0.50, 1.30, sizes["table"], accent, col_widths=[0.30, 0.70])
        add_rule(slide, w * 0.42, 5.32, w * 0.50, accent)
        add_bullets(slide, paper["tall_figure_bottom_notes"], w * 0.42, 5.54, w * 0.50, 0.78, sizes, name="BODY_NARROW_NOTE", space_after=0.0)


def summary_slide(prs: Presentation, paper: dict[str, Any], accent: RGBColor, sizes: dict[str, float]) -> None:
    w = prs.slide_width / EMU_PER_INCH
    h = prs.slide_height / EMU_PER_INCH
    slide = add_blank_slide(prs)
    add_header(slide, w, h, accent, "Part. 04", "总结")
    add_section(slide, w, 0.78, "4.1", "报告结论", accent)
    summary_lines = [
        "● " + paper["main_claim"],
        "• " + paper["summary_method"],
        "• " + paper["summary_result"],
        "• 方法边界主要来自约束建模、候选覆盖范围和真实执行误差。",
    ]
    summary_lines.extend(paper["summary_extra"])
    add_bullets(slide, summary_lines, 0.76, 1.48, w - 1.52, 3.18 if w < 11.0 else 2.42, sizes, name="BODY_SUMMARY", space_after=0.48 if w < 11.0 else 1.8)
    if w < 11.0:
        add_table(slide, [["对象", "约束"], ["轨迹", "连续性"], ["控制", "输入边界"], ["实验", "可执行性"]], 0.88, 3.64, w - 1.76, 1.05, sizes["table"], accent, col_widths=[0.35, 0.65])
        metric_row(slide, [["问题", "规划约束"], ["方法", "在线选择"], ["结果", "飞行验证"]], 0.88, h - 1.84, w - 1.76, accent, sizes["metric_value"], label_size=sizes["metric_label"])
    else:
        add_table(
            slide,
            paper["summary_rows"],
            0.88,
            3.50,
            w - 1.76,
            1.34,
            sizes["table"],
            accent,
        )
        add_bullets(slide, paper["summary_bridge"], 0.96, 5.10, w - 1.92, 0.54, sizes, name="BODY_SUMMARY_BRIDGE", space_after=0.0)
        metric_row(slide, [["问题", "任务约束"], ["方法", "闭环求解"], ["结果", "实验证据"]], 0.88, 5.92, w - 1.76, accent, sizes["metric_value"], label_size=sizes["metric_label"])


def thanks_slide(prs: Presentation, accent: RGBColor, template_id: str) -> None:
    w = prs.slide_width / EMU_PER_INCH
    h = prs.slide_height / EMU_PER_INCH
    slide = add_blank_slide(prs)
    add_rect(slide, 0, 0, w, h, RGBColor(246, 249, 252))
    band_h = h * 0.40
    band_y = h * 0.27
    add_rect(slide, 0, band_y, w, band_h, accent)
    title_color = WHITE if brightness(accent) < 0.72 else BLACK
    add_plain(slide, "谢谢！", w * 0.10, band_y + band_h * 0.33, w * 0.80, 0.64, 42.0, name="THANKS_TITLE", bold=True, color=title_color, align=PP_ALIGN.CENTER)
    add_plain(slide, "学术论文汇报", w * 0.36, h - 0.72, w * 0.28, 0.24, 12.0, name="THANKS_FOOTER", color=accent, align=PP_ALIGN.CENTER)


def build_deck(template: dict[str, Any], paper: dict[str, Any], out_dir: Path, *, template_root: Path | None) -> tuple[Path, str]:
    template_path = expand_path(template["path"], template_root=template_root)
    if not template_path.exists():
        raise FileNotFoundError(template_path)
    template_id = template["id"]
    profile = template.get("audit_profile", "compact")
    sizes = PROFILE_SIZE[profile]
    accent = theme_color(template_path, template.get("theme"))

    prs = Presentation(str(template_path))
    clear_slides(prs)
    assets = draw_assets(out_dir / "assets" / f"{template_id}-{paper['id']}", accent)
    cover_slide(prs, paper, accent, template_id, sizes)
    content_position(prs, paper, accent, sizes)
    method_overview(prs, paper, assets, accent, sizes)
    formula_slide(prs, paper, accent, sizes)
    evidence_slide(prs, paper, assets, accent, sizes)
    narrow_figure_slide(prs, paper, assets, accent, sizes)
    summary_slide(prs, paper, accent, sizes)
    thanks_slide(prs, accent, template_id)
    strip_empty_text_bodies(prs)

    deck_dir = out_dir / f"{template_id}-{paper['id']}"
    deck_dir.mkdir(parents=True, exist_ok=True)
    deck_path = deck_dir / "input.pptx"
    prs.save(deck_path)
    return deck_path, profile


def run(cmd: list[str], *, timeout: int) -> subprocess.CompletedProcess[str]:
    print("+", " ".join(cmd), flush=True)
    return subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, timeout=timeout)


def render_case(pptx: Path, case_dir: Path, *, timeout: int) -> int:
    office = matrix_runner.soffice_path()
    if not office:
        print("LibreOffice not found; smoke render skipped.")
        return 2
    try:
        pdf = matrix_runner.export_pdf(pptx, case_dir, office, timeout=timeout)
    except Exception as exc:
        print(f"LibreOffice export failed: {exc}")
        return 1
    py = matrix_runner.script_python()
    png_dir = case_dir / "png"
    preview = case_dir / "preview-grid.png"
    proc = run([py, str(RENDER), str(pdf), "--out-dir", str(png_dir), "--preview-grid", str(preview)], timeout=120)
    print(proc.stdout)
    if proc.returncode != 0:
        return proc.returncode
    proc = run([
        py,
        str(SCAN),
        str(png_dir),
        "--fail-on-warning",
        "--ignore-edge-slides",
        "--blank-warn",
        "0.74",
        "--min-band-fraction",
        "0.10",
        "--body-blank-warn",
        "0.89",
    ], timeout=60)
    print(proc.stdout)
    return proc.returncode


def load_templates(path: Path, *, template_root: Path | None) -> list[dict[str, Any]]:
    data = json.loads(path.read_text(encoding="utf-8"))
    templates = data.get("templates", [])
    resolved = []
    for template in templates:
        try:
            template_path = expand_path(template["path"], template_root=template_root)
        except Exception:
            continue
        if template_path.exists():
            resolved.append(template)
        else:
            print(f"skip missing template: {template.get('id', template.get('path'))} -> {template_path}")
    return resolved


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate small paper-report decks across local PPTX templates.")
    parser.add_argument("--papers", type=Path, default=DEFAULT_PAPERS)
    parser.add_argument("--templates", type=Path, default=DEFAULT_TEMPLATES)
    parser.add_argument("--template-root", type=Path, help="Root containing PPTAgent-style template subfolders.")
    parser.add_argument("--out-dir", type=Path, default=ROOT / "out" / "template-smoke")
    parser.add_argument("--case-limit", type=int, default=0, help="Limit generated cases for quick debugging.")
    parser.add_argument("--skip-render", action="store_true")
    parser.add_argument("--keep-going", action="store_true")
    parser.add_argument("--export-timeout", type=int, default=90)
    args = parser.parse_args()

    papers = json.loads(args.papers.read_text(encoding="utf-8")).get("papers", [])
    templates = load_templates(args.templates, template_root=args.template_root)
    if not papers:
        print("no stress papers configured", file=sys.stderr)
        return 2
    if not templates:
        print("no readable templates configured", file=sys.stderr)
        return 2

    py = matrix_runner.script_python()
    failures: list[str] = []
    count = 0
    args.out_dir.mkdir(parents=True, exist_ok=True)
    for template in templates:
        for paper in papers:
            if args.case_limit and count >= args.case_limit:
                break
            count += 1
            case_id = f"{template['id']}-{paper['id']}"
            print(f"\n== smoke {case_id} ==", flush=True)
            try:
                pptx, profile = build_deck(template, paper, args.out_dir, template_root=args.template_root)
            except Exception as exc:
                print(f"build failed: {exc}")
                failures.append(f"{case_id}: build failed")
                if not args.keep_going:
                    break
                continue
            proc = run([
                py,
                str(AUDIT),
                str(pptx),
                "--strict-body-hierarchy",
                "--profile",
                profile,
                "--fail-on-warning",
            ], timeout=60)
            print(proc.stdout)
            if proc.returncode != 0:
                failures.append(f"{case_id}: PPTX audit failed")
                if not args.keep_going:
                    break
                continue
            if not args.skip_render:
                rc = render_case(pptx, pptx.parent, timeout=args.export_timeout)
                if rc != 0:
                    failures.append(f"{case_id}: render/scan failed")
                    if not args.keep_going:
                        break
        if args.case_limit and count >= args.case_limit:
            break
        if failures and not args.keep_going:
            break

    if failures:
        print("\nTemplate smoke failed:")
        for failure in failures:
            print(" -", failure)
        return 1
    print(f"\nTemplate smoke passed ({count} case(s)).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
